#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""deck.json → HTML / PPTX ビルダー v2（パラメーター駆動）。

使い方:
    python3 tools/build_deck.py decks/<デッキ名>          # HTML + PPTX
    python3 tools/build_deck.py decks/<デッキ名> --html   # HTML のみ
    python3 tools/build_deck.py decks/<デッキ名> --pptx   # PPTX のみ

設計:
- 仮想キャンバス 1280x720px。全オブジェクトの位置・サイズ・フォント・色を px で持つ
- 解決順: templates/layouts/<layout>.json の既定値
          ← deck.meta.layout_overrides（デッキ全体の調整）
          ← slides[i].style（スライド個別の微修正）
- HTML は px をそのまま使用、PPTX は 96dpi 換算（px/96 インチ、フォント px*0.75 pt）
  なので両形式で座標・サイズが一致する
- スキーマは docs/deck-schema.md を参照

注意: PPTX にはスピーカーノートを出力しない（python-pptx の notes_slide は
macOS Keynote の互換性を壊す既知問題があるため）。notes は HTML でのみ表示される。
"""

import argparse
import base64
import copy
import html as html_mod
import json
import mimetypes
import re
import sys
from pathlib import Path
from string import Template

ROOT = Path(__file__).resolve().parent.parent
CANVAS_W, CANVAS_H = 1280, 720


# ---------------------------------------------------------------------------
# 読み込みとスタイル解決
# ---------------------------------------------------------------------------

def load_json(path: Path):
    if not path.exists():
        sys.exit(f"error: {path} が見つかりません")
    return json.loads(path.read_text(encoding="utf-8"))


def load_deck(deck_dir: Path):
    deck = load_json(deck_dir / "deck.json")
    meta = deck.get("meta", {})
    theme = load_json(ROOT / "templates" / "themes" / f"{meta.get('theme', 'default')}.json")
    layout = load_json(ROOT / "templates" / "layouts" / f"{meta.get('layout', 'default')}.json")
    return deck, theme, layout


def deep_merge(base, over):
    if not isinstance(base, dict) or not isinstance(over, dict):
        return copy.deepcopy(over)
    out = copy.deepcopy(base)
    for k, v in over.items():
        out[k] = deep_merge(out[k], v) if k in out else copy.deepcopy(v)
    return out


def resolve_style(layout, deck, slide):
    t = slide["type"]
    types = layout.get("types", {})
    if t not in types:
        sys.exit(f"error: 未知のスライドタイプ: {t}")
    st = types[t]
    st = deep_merge(st, deck.get("meta", {}).get("layout_overrides", {}).get(t, {}))
    st = deep_merge(st, slide.get("style", {}))
    return st


def col(theme, c):
    """テーマトークン名 → hex。hex 直書きはそのまま返す。"""
    return theme["colors"].get(c, c)


def norm_items(items):
    """bullets 配列を正規化。要素は文字列 or {text, children, size, color, bold}。"""
    out = []
    for it in items or []:
        if isinstance(it, str):
            out.append({"text": it, "children": []})
        else:
            d = dict(it)
            d.setdefault("text", "")
            kids = []
            for c in d.get("children") or []:
                kids.append({"text": c} if isinstance(c, str) else dict(c))
            d["children"] = kids
            out.append(d)
    return out


def fit_code_size(code: str, r: dict):
    """コード領域に収まるフォントサイズを計算（HTML/PPTX 共通ロジック）。"""
    lines = code.splitlines() or [""]
    size = r["size"]
    lh = r["line_height"]
    avail = r["h"] - 2 * r["pad"]
    while size > r.get("min_size", 12) and len(lines) * size * lh > avail:
        size -= 1
    return size, lines


def svg_size(svg_text: str):
    """SVG の描画サイズを viewBox / width / height 属性から取得する。"""
    m = re.search(
        r'viewBox\s*=\s*"\s*[-\d.]+[\s,]+[-\d.]+[\s,]+([\d.]+)[\s,]+([\d.]+)\s*"',
        svg_text,
    )
    if m:
        return float(m.group(1)), float(m.group(2))
    wm = re.search(r'<svg[^>]*?\swidth\s*=\s*"([\d.]+)(?:px)?"', svg_text)
    hm = re.search(r'<svg[^>]*?\sheight\s*=\s*"([\d.]+)(?:px)?"', svg_text)
    if wm and hm:
        return float(wm.group(1)), float(hm.group(1))
    return 1136.0, 440.0


def svg_to_png(svg_path: Path, out_png: Path):
    """SVG を透過 PNG にラスタライズ（PPTX 埋め込み用、2x 解像度）。"""
    from playwright.sync_api import sync_playwright

    svg_text = svg_path.read_text(encoding="utf-8")
    w, h = svg_size(svg_text)
    with sync_playwright() as pw:
        browser = pw.chromium.launch()
        page = browser.new_page(
            viewport={"width": max(1, int(w)), "height": max(1, int(h))},
            device_scale_factor=2,
        )
        page.set_content(
            "<!doctype html><style>html,body{margin:0;background:transparent}"
            "svg{display:block}</style>" + svg_text
        )
        page.wait_for_timeout(120)
        page.screenshot(path=str(out_png), omit_background=True,
                        clip={"x": 0, "y": 0, "width": w, "height": h})
        browser.close()


def flip_regions(st: dict, keys) -> dict:
    """image_side: "left" 用に、指定領域の x をキャンバス中心で左右反転する。"""
    out = copy.deepcopy(st)
    for k in keys:
        r = out.get(k)
        if isinstance(r, dict) and "x" in r and "w" in r:
            r["x"] = CANVAS_W - r["x"] - r["w"]
    return out


def image_source(deck_dir: Path, rel_path: str, for_pptx: bool) -> Path:
    """image スライドの実ファイルを解決。PPTX 用に SVG は PNG 化してキャッシュする。"""
    img_path = deck_dir / rel_path
    if for_pptx and img_path.suffix.lower() == ".svg":
        cache = deck_dir / "build" / ".svg-cache"
        cache.mkdir(parents=True, exist_ok=True)
        png = cache / (img_path.stem + ".png")
        if not png.exists() or png.stat().st_mtime < img_path.stat().st_mtime:
            svg_to_png(img_path, png)
        return png
    return img_path


# ---------------------------------------------------------------------------
# HTML ビルド
# ---------------------------------------------------------------------------

def esc(s):
    return html_mod.escape(str(s), quote=True)


def pos(r, extra=""):
    s = (f"position:absolute;left:{r['x']}px;top:{r['y']}px;width:{r['w']}px;")
    if "h" in r:
        s += f"height:{r['h']}px;"
    return s + extra


def text_style(r, theme):
    s = f"font-size:{r['size']}px;color:{col(theme, r.get('color', 'text'))};"
    if r.get("bold"):
        s += "font-weight:700;"
    if r.get("align"):
        s += f"text-align:{r['align']};"
    if r.get("line_height"):
        s += f"line-height:{r['line_height']};"
    return s


def div(r, theme, inner, extra=""):
    return f'<div style="{pos(r)}{text_style(r, theme)}{extra}">{inner}</div>'


def rect(r, theme, extra=""):
    fill = col(theme, r.get("fill", "accent"))
    radius = f"border-radius:{r['radius']}px;" if r.get("radius") else ""
    return f'<div style="{pos(r)}background:{fill};{radius}{extra}"></div>'


def bullets_html(items, r, theme):
    mk = col(theme, r.get("marker_color", "accent"))
    parts = [
        f'<ul style="margin:0;padding-left:{r["indent"]}px;list-style:disc;'
        f'font-size:{r["size"]}px;line-height:{r["line_height"]};'
        f'color:{col(theme, r["color"])};--mk:{mk};">'
    ]
    for it in items:
        li = f'font-size:{it["size"]}px;' if it.get("size") else ""
        li += f'color:{col(theme, it["color"])};' if it.get("color") else ""
        li += "font-weight:700;" if it.get("bold") else ""
        parts.append(f'<li style="margin-bottom:{r["gap"]}px;{li}">{esc(it["text"])}')
        if it["children"]:
            parts.append(
                f'<ul style="margin:{r["child_gap"]}px 0 0;padding-left:{r["indent"] - 6}px;'
                f'list-style:circle;font-size:{r["child_size"]}px;'
                f'color:{col(theme, r["child_color"])};">'
            )
            for c in it["children"]:
                cs = f'font-size:{c["size"]}px;' if c.get("size") else ""
                cs += f'color:{col(theme, c["color"])};' if c.get("color") else ""
                parts.append(f'<li style="margin-bottom:{r["child_gap"]}px;{cs}">{esc(c["text"])}</li>')
            parts.append("</ul>")
        parts.append("</li>")
    parts.append("</ul>")
    return "".join(parts)


def footer_html(st, deck, theme, page, total):
    out = ""
    if "footer_l" in st:
        out += div(st["footer_l"], theme, esc(deck["meta"].get("title", "")))
    if "footer_r" in st:
        out += div(st["footer_r"], theme, f"{page} / {total}")
    return out


def html_slide_body(slide, st, deck, theme, deck_dir, page, total):
    t = slide["type"]
    fonts = theme["fonts"]
    parts = []

    def chrome():
        parts.append(div(st["title"], theme, esc(slide.get("title", ""))))
        parts.append(rect(st["rule"], theme))
        if slide.get("lead") and "lead" in st:
            parts.append(div(st["lead"], theme, esc(slide["lead"])))
        parts.append(footer_html(st, deck, theme, page, total))

    if t == "title":
        parts.append(rect(st["bar"], theme))
        parts.append(div(st["title"], theme, esc(slide.get("title", ""))))
        if slide.get("subtitle"):
            parts.append(div(st["subtitle"], theme, esc(slide["subtitle"])))
        if slide.get("meta"):
            parts.append(div(st["meta"], theme, esc(slide["meta"])))

    elif t == "section":
        if slide.get("number"):
            parts.append(div(st["number"], theme, esc(slide["number"])))
        parts.append(div(st["title"], theme, esc(slide.get("title", ""))))
        parts.append(rect(st["rule"], theme))
        if slide.get("subtitle"):
            parts.append(div(st["subtitle"], theme, esc(slide["subtitle"])))

    elif t == "bullets":
        chrome()
        parts.append(
            f'<div style="{pos(st["body"])}">'
            + bullets_html(norm_items(slide.get("bullets")), st["body"], theme)
            + "</div>"
        )

    elif t == "two_column":
        chrome()
        for side in ("left", "right"):
            box = st[side]
            colc = slide.get(side, {})
            radius = f"border-radius:{box.get('radius', 0)}px;"
            inner = ""
            hd = st["col_heading"]
            if colc.get("heading"):
                inner += (
                    f'<div style="font-size:{hd["size"]}px;color:{col(theme, hd["color"])};'
                    f'font-weight:700;margin-bottom:{hd["gap_below"]}px;">{esc(colc["heading"])}</div>'
                )
            inner += bullets_html(norm_items(colc.get("bullets")), st["col_body"], theme)
            parts.append(
                f'<div style="{pos(box)}background:{col(theme, box["fill"])};{radius}'
                f'padding:{box["pad"]}px;box-sizing:border-box;overflow:hidden;">{inner}</div>'
            )

    elif t == "table":
        chrome()
        tr = st["table"]
        cols = slide.get("columns", [])
        widths = tr.get("col_widths")
        colgroup = ""
        if widths:
            total_w = sum(widths)
            colgroup = "<colgroup>" + "".join(
                f'<col style="width:{w / total_w * 100:.2f}%;">' for w in widths
            ) + "</colgroup>"
        cells_pad = f"padding:{tr['pad_y']}px {tr['pad_x']}px;"
        html = [
            f'<div style="{pos(tr)}"><table style="width:100%;border-collapse:collapse;'
            f'line-height:{tr["line_height"]};">{colgroup}<thead><tr>'
        ]
        for c in cols:
            html.append(
                f'<th style="background:{col(theme, tr["header_fill"])};'
                f'color:{col(theme, tr["header_color"])};font-size:{tr["header_size"]}px;'
                f'text-align:left;font-weight:600;height:{tr["header_h"]}px;'
                f'{cells_pad}box-sizing:border-box;">{esc(c)}</th>'
            )
        html.append("</tr></thead><tbody>")
        for i, row in enumerate(slide.get("rows", [])):
            bg = col(theme, tr["row_alt_fill"]) if i % 2 else col(theme, "background")
            html.append("<tr>")
            for cell in row:
                html.append(
                    f'<td style="background:{bg};color:{col(theme, tr["cell_color"])};'
                    f'font-size:{tr["cell_size"]}px;height:{tr["row_h"]}px;{cells_pad}'
                    f'box-sizing:border-box;border-bottom:1px solid rgba(0,0,0,0.08);">{esc(cell)}</td>'
                )
            html.append("</tr>")
        html.append("</tbody></table></div>")
        parts.append("".join(html))

    elif t == "code":
        chrome()
        r = st["code"]
        size, _lines = fit_code_size(slide.get("code", ""), r)
        radius = f"border-radius:{r.get('radius', 0)}px;"
        parts.append(
            f'<pre style="{pos(r)}margin:0;background:{col(theme, r["fill"])};{radius}'
            f'padding:{r["pad"]}px;box-sizing:border-box;overflow:hidden;">'
            f'<code style="font-family:\'{fonts["code"]}\',monospace;font-size:{size}px;'
            f'line-height:{r["line_height"]};color:{col(theme, r["color"])};">'
            f"{esc(slide.get('code', ''))}</code></pre>"
        )

    elif t == "quote":
        parts.append(div(st["mark"], theme, "“"))
        parts.append(div(st["text"], theme, esc(slide.get("text", ""))))
        if slide.get("attribution"):
            parts.append(div(st["attribution"], theme, "— " + esc(slide["attribution"])))

    elif t == "image":
        if slide.get("title"):
            parts.append(div(st["title"], theme, esc(slide["title"])))
            parts.append(rect(st["rule"], theme))
        r = st["img"]
        img_path = deck_dir / slide["path"]
        mime = mimetypes.guess_type(img_path.name)[0] or "image/png"
        data = base64.b64encode(img_path.read_bytes()).decode()
        parts.append(
            f'<div style="{pos(r)}"><img src="data:{mime};base64,{data}" alt="" '
            f'style="width:100%;height:100%;object-fit:contain;"></div>'
        )
        if slide.get("caption"):
            parts.append(div(st["caption"], theme, esc(slide["caption"])))
        parts.append(footer_html(st, deck, theme, page, total))

    elif t == "image_text":
        if slide.get("image_side") == "left":
            st = flip_regions(st, ("body", "img", "caption"))
        parts.append(div(st["title"], theme, esc(slide.get("title", ""))))
        parts.append(rect(st["rule"], theme))
        if slide.get("punch"):
            parts.append(div(st["punch"], theme, esc(slide["punch"])))
        parts.append(
            f'<div style="{pos(st["body"])}">'
            + bullets_html(norm_items(slide.get("bullets")), st["body"], theme)
            + "</div>"
        )
        r = st["img"]
        img_path = deck_dir / slide["path"]
        mime = mimetypes.guess_type(img_path.name)[0] or "image/png"
        data = base64.b64encode(img_path.read_bytes()).decode()
        parts.append(
            f'<div style="{pos(r)}"><img src="data:{mime};base64,{data}" alt="" '
            f'style="width:100%;height:100%;object-fit:contain;"></div>'
        )
        if slide.get("caption"):
            parts.append(div(st["caption"], theme, esc(slide["caption"])))
        parts.append(footer_html(st, deck, theme, page, total))

    elif t == "closing":
        parts.append(div(st["title"], theme, esc(slide.get("title", ""))))
        parts.append(rect(st["rule"], theme))
        if slide.get("bullets"):
            parts.append(
                f'<div style="{pos(st["body"])}">'
                + bullets_html(norm_items(slide.get("bullets")), st["body"], theme)
                + "</div>"
            )
        if slide.get("message"):
            parts.append(div(st["message"], theme, esc(slide["message"])))

    return "".join(parts)


HTML_TEMPLATE = Template("""<!doctype html>
<html lang="ja">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>$TITLE</title>
<style>
* { margin: 0; padding: 0; box-sizing: border-box; }
html, body { height: 100%; background: #101014; overflow: hidden; }
#stage { position: absolute; top: 50%; left: 50%; width: ${W}px; height: ${H}px; }
.slide {
  position: absolute; inset: 0; width: ${W}px; height: ${H}px;
  font-family: "$F_BODY", "Hiragino Sans", sans-serif;
  display: none; overflow: hidden;
}
.slide.active { display: block; }
.slide li::marker { color: var(--mk, inherit); }
#hud {
  position: fixed; right: 16px; bottom: 12px; color: #9a94ab;
  font: 13px/1 "$F_BODY", sans-serif; user-select: none; z-index: 10;
}
#notes {
  position: fixed; left: 0; right: 0; bottom: 0; max-height: 32vh; overflow: auto;
  background: rgba(16, 10, 26, 0.94); color: #e6def3; padding: 16px 28px 20px;
  font: 15px/1.7 "$F_BODY", sans-serif; z-index: 20; border-top: 2px solid $C_ACCENT;
}
#notes .label { color: $C_ACCENT; font-size: 12px; font-weight: 700; margin-bottom: 6px; }
</style>
</head>
<body>
<div id="stage">
$SLIDES
</div>
<div id="hud"><span id="pageno"></span>&nbsp;&nbsp;[&larr;/&rarr;] 移動 [N] ノート</div>
<div id="notes" hidden><div class="label">SPEAKER NOTES</div><div id="notes-body"></div></div>
<script>
var slides = Array.prototype.slice.call(document.querySelectorAll('.slide'));
var stage = document.getElementById('stage');
var pageno = document.getElementById('pageno');
var notes = document.getElementById('notes');
var notesBody = document.getElementById('notes-body');
var idx = 0;

function fit() {
  var s = Math.min(window.innerWidth / $W, window.innerHeight / $H);
  stage.style.transform = 'translate(-50%, -50%) scale(' + s + ')';
}
function show(i) {
  idx = Math.max(0, Math.min(slides.length - 1, i));
  slides.forEach(function (s, j) { s.classList.toggle('active', j === idx); });
  pageno.textContent = (idx + 1) + ' / ' + slides.length;
  location.replace('#' + (idx + 1));
  var n = slides[idx].getAttribute('data-notes') || '';
  notesBody.textContent = n || '（このスライドにノートはありません）';
}
function toggleNotes() { notes.hidden = !notes.hidden; }

window.addEventListener('resize', fit);
window.addEventListener('keydown', function (e) {
  if (e.key === 'ArrowRight' || e.key === 'PageDown' || e.key === ' ') { show(idx + 1); e.preventDefault(); }
  else if (e.key === 'ArrowLeft' || e.key === 'PageUp') { show(idx - 1); e.preventDefault(); }
  else if (e.key === 'Home') { show(0); }
  else if (e.key === 'End') { show(slides.length - 1); }
  else if (e.key === 'n' || e.key === 'N') { toggleNotes(); }
});
window.addEventListener('click', function (e) {
  if (notes.contains(e.target)) { return; }
  if (e.clientX > window.innerWidth / 2) { show(idx + 1); } else { show(idx - 1); }
});

var start = parseInt((location.hash || '#1').slice(1), 10);
fit();
show(isNaN(start) ? 0 : start - 1);
</script>
</body>
</html>
""")


def slide_bg(slide, st, theme):
    if "bg" in st:
        return col(theme, st["bg"].get("fill", "primary"))
    return col(theme, "background")


def build_html(deck, theme, layout, deck_dir: Path, out_path: Path):
    slides = deck.get("slides", [])
    total = len(slides)
    rendered = []
    for i, slide in enumerate(slides):
        st = resolve_style(layout, deck, slide)
        body = html_slide_body(slide, st, deck, theme, deck_dir, i + 1, total)
        notes_attr = f' data-notes="{esc(slide["notes"])}"' if slide.get("notes") else ""
        bg = slide_bg(slide, st, theme)
        rendered.append(
            f'<section class="slide" style="background:{bg};color:{col(theme, "text")};"'
            f"{notes_attr}>{body}</section>"
        )
    doc = HTML_TEMPLATE.safe_substitute(
        TITLE=esc(deck["meta"].get("title", deck["meta"].get("id", ""))),
        SLIDES="\n".join(rendered),
        W=str(CANVAS_W), H=str(CANVAS_H),
        C_ACCENT=col(theme, "accent"),
        F_BODY=theme["fonts"]["body"],
    )
    out_path.write_text(doc, encoding="utf-8")


# ---------------------------------------------------------------------------
# PPTX ビルド
# ---------------------------------------------------------------------------

def build_pptx(deck, theme, layout, deck_dir: Path, out_path: Path):
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from pptx.util import Emu, Inches, Pt

    fonts = theme["fonts"]

    def IN(px):
        return Inches(px / 96.0)

    def PTS(px):
        return Pt(px * 0.75)

    def C(name):
        return RGBColor.from_string(col(theme, name).lstrip("#"))

    ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}

    prs = Presentation()
    prs.slide_width = IN(CANVAS_W)
    prs.slide_height = IN(CANVAS_H)
    blank = prs.slide_layouts[6]

    def set_run(run, text, size_px, color, bold=False, name=None):
        run.text = text
        f = run.font
        f.name = name or fonts["body"]
        f.size = PTS(size_px)
        f.bold = bold
        f.color.rgb = color
        rPr = run._r.get_or_add_rPr()
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = rPr.makeelement(qn("a:ea"), {})
            rPr.append(ea)
        ea.set("typeface", name or fonts["body"])

    def add_box(slide, r):
        box = slide.shapes.add_textbox(IN(r["x"]), IN(r["y"]), IN(r["w"]), IN(r.get("h", 40)))
        tf = box.text_frame
        tf.word_wrap = True
        for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
            setattr(tf, m, Emu(0))
        return tf

    def para(tf, used):
        if not used[0]:
            used[0] = True
            return tf.paragraphs[0]
        return tf.add_paragraph()

    def text_region(slide, r, text, font=None):
        tf = add_box(slide, r)
        p = tf.paragraphs[0]
        if r.get("align"):
            p.alignment = ALIGN[r["align"]]
        if r.get("line_height"):
            p.line_spacing = r["line_height"]
        set_run(p.add_run(), text, r["size"], C(r.get("color", "text")),
                bold=r.get("bold", False), name=font)
        return tf

    def add_rect(slide, x, y, w, h, fill, radius=0):
        shape = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        sp = slide.shapes.add_shape(shape, IN(x), IN(y), IN(w), IN(h))
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
        sp.line.fill.background()
        sp.shadow.inherit = False
        return sp

    def rect_region(slide, r):
        add_rect(slide, r["x"], r["y"], r["w"], r["h"], C(r.get("fill", "accent")),
                 radius=r.get("radius", 0))

    def para_indent(p, px):
        p._p.get_or_add_pPr().set("marL", str(Emu(IN(px))))

    def add_bullets(slide, items, r):
        tf = add_box(slide, r)
        used = [False]
        for it in items:
            p = para(tf, used)
            p.space_after = PTS(r["gap"])
            p.line_spacing = r["line_height"]
            size = it.get("size", r["size"])
            color = C(it["color"]) if it.get("color") else C(r["color"])
            set_run(p.add_run(), "•  ", size, C(r["marker_color"]), bold=True)
            set_run(p.add_run(), it["text"], size, color, bold=it.get("bold", False))
            for c in it["children"]:
                p = para(tf, used)
                p.space_after = PTS(r["child_gap"])
                p.line_spacing = max(1.0, r["line_height"] - 0.15)
                para_indent(p, r["indent"] + 6)
                csize = c.get("size", r["child_size"])
                ccolor = C(c["color"]) if c.get("color") else C(r["child_color"])
                set_run(p.add_run(), "–  ", csize, C(r["marker_color"]))
                set_run(p.add_run(), c["text"], csize, ccolor)
        return tf

    def footer(slide, st, page, total):
        if "footer_l" in st:
            text_region(slide, st["footer_l"], deck["meta"].get("title", ""))
        if "footer_r" in st:
            text_region(slide, st["footer_r"], f"{page} / {total}")

    def chrome(slide, data, st, page, total):
        text_region(slide, st["title"], data.get("title", ""))
        rect_region(slide, st["rule"])
        if data.get("lead") and "lead" in st:
            text_region(slide, st["lead"], data["lead"])
        footer(slide, st, page, total)

    def s_title(slide, data, st, page, total):
        rect_region(slide, st["bar"])
        text_region(slide, st["title"], data.get("title", ""))
        if data.get("subtitle"):
            text_region(slide, st["subtitle"], data["subtitle"])
        if data.get("meta"):
            text_region(slide, st["meta"], data["meta"])

    def s_section(slide, data, st, page, total):
        add_rect(slide, 0, 0, CANVAS_W, CANVAS_H, C(st["bg"]["fill"]))
        if data.get("number"):
            text_region(slide, st["number"], data["number"])
        text_region(slide, st["title"], data.get("title", ""))
        rect_region(slide, st["rule"])
        if data.get("subtitle"):
            text_region(slide, st["subtitle"], data["subtitle"])

    def s_bullets(slide, data, st, page, total):
        chrome(slide, data, st, page, total)
        add_bullets(slide, norm_items(data.get("bullets")), st["body"])

    def s_two_column(slide, data, st, page, total):
        chrome(slide, data, st, page, total)
        for side in ("left", "right"):
            box = st[side]
            colc = data.get(side, {})
            add_rect(slide, box["x"], box["y"], box["w"], box["h"],
                     C(box["fill"]), radius=box.get("radius", 0))
            pad = box["pad"]
            iy = box["y"] + pad
            if colc.get("heading"):
                hd = st["col_heading"]
                hr = {"x": box["x"] + pad, "y": iy, "w": box["w"] - 2 * pad, "h": hd["size"] * 1.4,
                      "size": hd["size"], "color": hd["color"], "bold": hd.get("bold", False)}
                text_region(slide, hr, colc["heading"])
                iy += hd["size"] * 1.4 + hd["gap_below"]
            body = dict(st["col_body"])
            body.update({"x": box["x"] + pad, "y": iy, "w": box["w"] - 2 * pad,
                         "h": box["y"] + box["h"] - pad - iy})
            add_bullets(slide, norm_items(colc.get("bullets")), body)

    def s_table(slide, data, st, page, total):
        from pptx.enum.text import MSO_ANCHOR
        chrome(slide, data, st, page, total)
        tr = st["table"]
        cols = data.get("columns", [])
        rows = data.get("rows", [])
        n_rows, n_cols = len(rows) + 1, len(cols)
        total_h = tr["header_h"] + tr["row_h"] * len(rows)
        gfx = slide.shapes.add_table(n_rows, n_cols, IN(tr["x"]), IN(tr["y"]),
                                     IN(tr["w"]), IN(total_h))
        table = gfx.table
        table.first_row = False
        table.horz_banding = False
        widths = tr.get("col_widths")
        if widths:
            unit = tr["w"] / sum(widths)
            for j, w in enumerate(widths):
                table.columns[j].width = IN(w * unit)
        table.rows[0].height = IN(tr["header_h"])
        for i in range(len(rows)):
            table.rows[i + 1].height = IN(tr["row_h"])
        for j, name in enumerate(cols):
            cell = table.cell(0, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = C(tr["header_fill"])
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = IN(tr["pad_x"])
            cell.margin_right = IN(tr["pad_x"])
            p = cell.text_frame.paragraphs[0]
            set_run(p.add_run(), str(name), tr["header_size"], C(tr["header_color"]), bold=True)
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                cell = table.cell(i + 1, j)
                cell.fill.solid()
                cell.fill.fore_color.rgb = C(tr["row_alt_fill"]) if i % 2 else C("background")
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.margin_left = IN(tr["pad_x"])
                cell.margin_right = IN(tr["pad_x"])
                p = cell.text_frame.paragraphs[0]
                p.line_spacing = tr["line_height"]
                set_run(p.add_run(), str(val), tr["cell_size"], C(tr["cell_color"]))

    def s_code(slide, data, st, page, total):
        chrome(slide, data, st, page, total)
        r = st["code"]
        size, lines = fit_code_size(data.get("code", ""), r)
        add_rect(slide, r["x"], r["y"], r["w"], r["h"], C(r["fill"]), radius=r.get("radius", 0))
        tf = add_box(slide, {"x": r["x"] + r["pad"], "y": r["y"] + r["pad"],
                             "w": r["w"] - 2 * r["pad"], "h": r["h"] - 2 * r["pad"]})
        used = [False]
        for line in lines:
            p = para(tf, used)
            p.line_spacing = r["line_height"]
            set_run(p.add_run(), line if line else " ", size, C(r["color"]),
                    name=fonts["code"])

    def s_quote(slide, data, st, page, total):
        text_region(slide, st["mark"], "“")
        text_region(slide, st["text"], data.get("text", ""))
        if data.get("attribution"):
            text_region(slide, st["attribution"], "— " + data["attribution"])

    def s_image(slide, data, st, page, total):
        if data.get("title"):
            text_region(slide, st["title"], data["title"])
            rect_region(slide, st["rule"])
        r = st["img"]
        pic = slide.shapes.add_picture(str(image_source(deck_dir, data["path"], True)),
                                       IN(r["x"]), IN(r["y"]))
        ratio = min(IN(r["w"]) / pic.width, IN(r["h"]) / pic.height, 1.0)
        pic.width = int(pic.width * ratio)
        pic.height = int(pic.height * ratio)
        pic.left = int(IN(r["x"]) + (IN(r["w"]) - pic.width) / 2)
        pic.top = int(IN(r["y"]) + (IN(r["h"]) - pic.height) / 2)
        if data.get("caption"):
            text_region(slide, st["caption"], data["caption"])
        footer(slide, st, page, total)

    def s_image_text(slide, data, st, page, total):
        if data.get("image_side") == "left":
            st = flip_regions(st, ("body", "img", "caption"))
        text_region(slide, st["title"], data.get("title", ""))
        rect_region(slide, st["rule"])
        if data.get("punch"):
            text_region(slide, st["punch"], data["punch"])
        add_bullets(slide, norm_items(data.get("bullets")), st["body"])
        r = st["img"]
        pic = slide.shapes.add_picture(str(image_source(deck_dir, data["path"], True)),
                                       IN(r["x"]), IN(r["y"]))
        ratio = min(IN(r["w"]) / pic.width, IN(r["h"]) / pic.height, 1.0)
        pic.width = int(pic.width * ratio)
        pic.height = int(pic.height * ratio)
        pic.left = int(IN(r["x"]) + (IN(r["w"]) - pic.width) / 2)
        pic.top = int(IN(r["y"]) + (IN(r["h"]) - pic.height) / 2)
        if data.get("caption"):
            text_region(slide, st["caption"], data["caption"])
        footer(slide, st, page, total)

    def s_closing(slide, data, st, page, total):
        add_rect(slide, 0, 0, CANVAS_W, CANVAS_H, C(st["bg"]["fill"]))
        text_region(slide, st["title"], data.get("title", ""))
        rect_region(slide, st["rule"])
        if data.get("bullets"):
            add_bullets(slide, norm_items(data.get("bullets")), st["body"])
        if data.get("message"):
            text_region(slide, st["message"], data["message"])

    builders = {
        "title": s_title, "section": s_section, "bullets": s_bullets,
        "two_column": s_two_column, "table": s_table, "code": s_code,
        "quote": s_quote, "image": s_image, "image_text": s_image_text,
        "closing": s_closing,
    }

    slides = deck.get("slides", [])
    total = len(slides)
    for i, data in enumerate(slides):
        st = resolve_style(layout, deck, data)
        slide = prs.slides.add_slide(blank)
        if "bg" not in st:
            add_rect(slide, 0, 0, CANVAS_W, CANVAS_H, C("background"))
        builders[data["type"]](slide, data, st, i + 1, total)
        # 注意: notes は意図的に PPTX へ出力しない（Keynote 互換性問題のため）

    prs.save(str(out_path))


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main():
    ap = argparse.ArgumentParser(description="deck.json → HTML / PPTX ビルダー v2")
    ap.add_argument("deck_dir", help="デッキのディレクトリ（deck.json を含む）")
    ap.add_argument("--html", action="store_true", help="HTML のみ生成")
    ap.add_argument("--pptx", action="store_true", help="PPTX のみ生成")
    args = ap.parse_args()

    deck_dir = Path(args.deck_dir).resolve()
    deck, theme, layout = load_deck(deck_dir)
    deck_id = deck.get("meta", {}).get("id") or deck_dir.name

    out_dir = deck_dir / "build"
    out_dir.mkdir(exist_ok=True)

    if args.html or not args.pptx:
        out = out_dir / f"{deck_id}.html"
        build_html(deck, theme, layout, deck_dir, out)
        print(f"HTML: {out}")
    if args.pptx or not args.html:
        out = out_dir / f"{deck_id}.pptx"
        build_pptx(deck, theme, layout, deck_dir, out)
        print(f"PPTX: {out}")


if __name__ == "__main__":
    main()
