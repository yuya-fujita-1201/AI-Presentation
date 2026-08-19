#!/usr/bin/env python3
"""Verify the shared visual contract and rendered artifacts for AI Engineering 03-05."""

from __future__ import annotations

import colorsys
import json
import re
import sys
import zipfile
from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
DECK_ROOT = ROOT / "decks"
REFERENCE = DECK_ROOT / "01-prompt-engineering"
TARGETS = {
    "03-harness-engineering": 66,
    "04-loop-engineering": 55,
    "05-graph-engineering": 48,
}
INTRO_TYPES = ["title", "image_text", "image_text", "section"]
FIVE_LAYER_LABELS = (
    "頼み方（プロンプト）",
    "渡す情報（コンテキスト）",
    "安全な環境（ハーネス）",
    "確認の繰り返し（ループ）",
    "処理のつながり（グラフ）",
)
FIVE_LAYER_BULLETS = (
    "頼み方（プロンプト）：1回の指示を設計する",
    "渡す情報（コンテキスト）：判断に使う情報を設計する",
    "安全な環境（ハーネス）：安全に動き、検証できる環境を設計する",
    "確認の繰り返し（ループ）：実行・観察・改善を時間軸でつなぐ",
    "処理のつながり（グラフ）：複数の役割を構造としてつなぐ",
)
FIVE_LAYER_SUBLABELS = (
    "1回の指示",
    "資料・前提",
    "権限・ツール・実行環境",
    "作る→確かめる→直す",
    "順序・分岐・役割分担",
)
FIVE_LAYER_DISCLAIMER = (
    "この5層モデルは、直す場所を見つけるための本シリーズ独自の整理で、"
    "業界標準の分類ではありません"
)
FOCUS_INDEX = {
    "03-harness-engineering": 2,
    "04-loop-engineering": 3,
    "05-graph-engineering": 4,
}
PURPLE_TOKENS = (
    "accenture-purple",
    "#460073",
    "#a100ff",
    "#d9bbf2",
    "#f7f0fc",
    "#f1e5fb",
)
PLACEHOLDER_RE = re.compile(r"(?:\bTODO\b|\bTBD\b|lorem|ipsum|xxxx)", re.I)


def load_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def referenced_asset_paths(deck_dir: Path, data: dict[str, Any]) -> set[Path]:
    paths: set[Path] = set()
    for slide in data.get("slides", []):
        for key in ("image", "path", "src"):
            value = slide.get(key)
            if isinstance(value, str) and value.startswith("assets/"):
                paths.add(deck_dir / value)
    return paths


def purple_ratio(path: Path) -> float:
    with Image.open(path) as source:
        image = source.convert("RGB")
        image.thumbnail((256, 256))
        purple = 0
        total = image.width * image.height
        pixels = (
            image.get_flattened_data()
            if hasattr(image, "get_flattened_data")
            else image.getdata()
        )
        for red, green, blue in pixels:
            hue, saturation, value = colorsys.rgb_to_hsv(
                red / 255.0, green / 255.0, blue / 255.0
            )
            degrees = hue * 360.0
            if 260.0 <= degrees <= 315.0 and saturation >= 0.18 and value >= 0.18:
                purple += 1
        return purple / max(total, 1)


def verify_deck(
    name: str,
    expected_count: int,
    reference_layout: dict[str, Any],
) -> list[str]:
    errors: list[str] = []
    deck_dir = DECK_ROOT / name
    deck_path = deck_dir / "deck.json"
    data = load_json(deck_path)
    meta = data.get("meta", {})
    slides = data.get("slides", [])

    if len(slides) != expected_count:
        errors.append(f"{name}: deck.json slide count {len(slides)} != {expected_count}")
    if meta.get("theme") != "warm-terracotta":
        errors.append(f"{name}: theme is {meta.get('theme')!r}, expected 'warm-terracotta'")
    if meta.get("layout_overrides") != reference_layout:
        errors.append(f"{name}: meta.layout_overrides differs from AI Engineering 01")

    actual_intro = [slide.get("type") for slide in slides[:4]]
    if actual_intro != INTRO_TYPES:
        errors.append(f"{name}: intro types {actual_intro!r} != {INTRO_TYPES!r}")

    if len(slides) >= 2:
        layer_slide = slides[1]
        layer_bullets = layer_slide.get("bullets", [])
        bullet_texts = tuple(
            bullet.get("text", "") if isinstance(bullet, dict) else str(bullet)
            for bullet in layer_bullets[:5]
        )
        if bullet_texts != FIVE_LAYER_BULLETS:
            errors.append(f"{name}: slide 2 five-layer bullet copy is not canonical")
        if len(layer_bullets) < 6 or not isinstance(layer_bullets[5], dict):
            errors.append(f"{name}: slide 2 is missing the canonical disclaimer")
        elif layer_bullets[5].get("text") != FIVE_LAYER_DISCLAIMER:
            errors.append(f"{name}: slide 2 disclaimer is not canonical")
        focus_index = FOCUS_INDEX[name]
        for index, bullet in enumerate(layer_bullets[:5]):
            is_focus = (
                isinstance(bullet, dict)
                and bullet.get("bold") is True
                and bullet.get("color") == "accent"
            )
            if is_focus != (index == focus_index):
                errors.append(
                    f"{name}: slide 2 focus styling is incorrect at layer {index + 1}"
                )
        layer_ref = next(
            (
                layer_slide.get(key)
                for key in ("image", "path", "src")
                if isinstance(layer_slide.get(key), str)
            ),
            None,
        )
        if not layer_ref:
            errors.append(f"{name}: slide 2 has no five-layer image reference")
        else:
            layer_path = deck_dir / layer_ref
            if not layer_path.is_file():
                errors.append(f"{name}: slide 2 image is missing: {layer_ref}")
            elif layer_path.suffix.lower() != ".svg":
                errors.append(f"{name}: slide 2 five-layer image is not SVG: {layer_ref}")
            else:
                layer_text = layer_path.read_text(encoding="utf-8")
                for label in FIVE_LAYER_LABELS:
                    if label not in layer_text:
                        errors.append(f"{name}: five-layer diagram is missing label {label!r}")
                for sublabel in FIVE_LAYER_SUBLABELS:
                    if sublabel not in layer_text:
                        errors.append(
                            f"{name}: five-layer diagram is missing sublabel {sublabel!r}"
                        )

    raw_deck = deck_path.read_text(encoding="utf-8")
    if PLACEHOLDER_RE.search(raw_deck):
        errors.append(f"{name}: placeholder-like text remains in deck.json")

    assets = referenced_asset_paths(deck_dir, data)
    for asset in sorted(assets):
        if not asset.is_file():
            errors.append(f"{name}: referenced asset is missing: {asset.relative_to(deck_dir)}")

    if name == "05-graph-engineering":
        text_files = [deck_path, *sorted(asset for asset in assets if asset.suffix.lower() == ".svg")]
        for path in text_files:
            lowered = path.read_text(encoding="utf-8").lower()
            for token in PURPLE_TOKENS:
                if token in lowered:
                    errors.append(
                        f"{name}: purple token {token!r} remains in {path.relative_to(deck_dir)}"
                    )
        for asset in sorted(asset for asset in assets if asset.suffix.lower() == ".png"):
            ratio = purple_ratio(asset)
            if ratio >= 0.025:
                errors.append(
                    f"{name}: referenced PNG remains purple-dominant: "
                    f"{asset.relative_to(deck_dir)} ({ratio:.1%} purple pixels)"
                )

    build_dir = deck_dir / "build"
    artifact_id = meta.get("id", name)
    pptx_path = build_dir / f"{artifact_id}.pptx"
    if not pptx_path.is_file():
        errors.append(f"{name}: PPTX is missing: {pptx_path.relative_to(deck_dir)}")
    else:
        try:
            with zipfile.ZipFile(pptx_path) as archive:
                bad_member = archive.testzip()
                if bad_member:
                    errors.append(f"{name}: corrupt PPTX member: {bad_member}")
                notes_parts = [item for item in archive.namelist() if item.startswith("ppt/notesSlides/")]
                if notes_parts:
                    errors.append(f"{name}: PPTX contains {len(notes_parts)} notesSlide parts")
            presentation = Presentation(pptx_path)
            if len(presentation.slides) != len(slides):
                errors.append(
                    f"{name}: PPTX slide count {len(presentation.slides)} != deck.json {len(slides)}"
                )
            notes_flags = [index for index, slide in enumerate(presentation.slides, 1) if slide.has_notes_slide]
            if notes_flags:
                errors.append(f"{name}: has_notes_slide=True on slides {notes_flags}")
        except Exception as exc:  # pragma: no cover - diagnostic boundary
            errors.append(f"{name}: PPTX reparse failed: {exc}")

    previews = sorted((build_dir / "preview").glob("slide-*.png"))
    if len(previews) != len(slides):
        errors.append(f"{name}: preview count {len(previews)} != deck.json {len(slides)}")
    for preview in previews:
        try:
            with Image.open(preview) as image:
                if image.size != (1280, 720):
                    errors.append(
                        f"{name}: {preview.name} size {image.size} != (1280, 720)"
                    )
        except Exception as exc:  # pragma: no cover - diagnostic boundary
            errors.append(f"{name}: cannot read {preview.name}: {exc}")

    return errors


def main() -> int:
    reference = load_json(REFERENCE / "deck.json")
    reference_layout = reference["meta"]["layout_overrides"]
    all_errors: list[str] = []
    for name, expected_count in TARGETS.items():
        errors = verify_deck(name, expected_count, reference_layout)
        status = "OK" if not errors else f"FAIL ({len(errors)})"
        print(f"{name}: {status}")
        all_errors.extend(errors)

    if all_errors:
        print("\nSeries verification errors:")
        for error in all_errors:
            print(f"- {error}")
        return 1

    print("\nAI Engineering 03-05 series verification: OK")
    return 0


if __name__ == "__main__":
    sys.exit(main())
