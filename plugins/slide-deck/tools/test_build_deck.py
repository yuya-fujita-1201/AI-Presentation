#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""build_deck.py の回帰テスト（unittest）。

対象:
  - 見本デッキ（examples/template-sample）のHTML/PPTXビルドが通ること
  - 生成PPTXにスピーカーノートが一切無いこと（has_notes_slide が全スライドで False）
  - expand_slides() が自動挿入込みで想定枚数（swimlaneの凡例1枚を含め20枚）になること
  - テーマの extends 継承と、明示していないトークンの派生規則（spec.md の mix() 式）
  - meta.layout_overrides.common と slides[i].style の重ね合わせ（resolve_style）
  - validate_deck() が検出すべき代表的なエラー
      * type 欠落 → fail
      * 色トークンの typo → fail
      * table の行セル数と columns 数の不一致 → fail
      * image の path がデッキ外を指す／実体が無い → fail
      * BOM付き UTF-8 の deck.json → 成功する（utf-8-sig 読み込み）

実行:
    python -m unittest discover -s plugins/slide-deck/tools -p "test_*.py"
    または
    python plugins/slide-deck/tools/test_build_deck.py

ビルド成果物を検証する系のテストは、実際のユーザー操作（CLI呼び出し）を模すため
sys.executable で build_deck.py をサブプロセス実行して確認する。テーマ導出・
expand_slides・resolve_style は、シグネチャ維持が約束されている公開API
（load_json / theme_dirs / load_theme / load_deck / expand_slides / resolve_style / col）
を直接importして検証する。
"""
from __future__ import annotations

import json
import subprocess
import sys
import tempfile
import unittest
import zipfile
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
TOOLS = ROOT / "tools"
BUILD_DECK = TOOLS / "build_deck.py"
SAMPLE_DECK = ROOT / "examples" / "template-sample"

sys.path.insert(0, str(TOOLS))
import build_deck  # noqa: E402  (テーマ/expand_slides/resolve_style 直接検証用)


def run_build(deck_dir, *extra_args):
    """build_deck.py <deck_dir> [args...] をサブプロセスで実行し CompletedProcess を返す。"""
    cmd = [sys.executable, str(BUILD_DECK), str(deck_dir), *extra_args]
    return subprocess.run(cmd, cwd=ROOT, capture_output=True, text=True)


def write_deck(tmp_path, deck, deck_id="t", bom=False):
    deck_dir = tmp_path / deck_id
    deck_dir.mkdir(parents=True, exist_ok=True)
    text = json.dumps(deck, ensure_ascii=False)
    encoding = "utf-8-sig" if bom else "utf-8"
    (deck_dir / "deck.json").write_text(text, encoding=encoding)
    return deck_dir


def minimal_deck():
    return {
        "meta": {"id": "t", "title": "T", "theme": "default", "layout": "default"},
        "slides": [{"type": "title", "title": "Hello", "subtitle": "World"}],
    }


# ---------------------------------------------------------------------------
# 見本デッキのビルド・ノート不在・展開後枚数
# ---------------------------------------------------------------------------

class SampleDeckBuildTests(unittest.TestCase):
    """examples/template-sample は公式の見本として常にビルドできる状態を保つ。"""

    @classmethod
    def setUpClass(cls):
        cls.result = run_build(SAMPLE_DECK)
        cls.html_path = SAMPLE_DECK / "build" / "template-sample.html"
        cls.pptx_path = SAMPLE_DECK / "build" / "template-sample.pptx"

    def test_build_succeeds(self):
        self.assertEqual(self.result.returncode, 0, msg=self.result.stderr)

    def test_html_and_pptx_are_produced(self):
        self.assertTrue(self.html_path.exists(), msg=self.result.stderr)
        self.assertTrue(self.pptx_path.exists(), msg=self.result.stderr)

    def test_pptx_is_a_valid_zip(self):
        with zipfile.ZipFile(self.pptx_path) as z:
            self.assertIsNone(z.testzip())

    def test_pptx_reparses_with_python_pptx(self):
        from pptx import Presentation
        prs = Presentation(str(self.pptx_path))
        self.assertGreater(len(prs.slides), 0)

    def test_no_speaker_notes_anywhere(self):
        """PPTX にスピーカーノートを絶対に入れない、というプロジェクトの厳守事項。"""
        from pptx import Presentation
        prs = Presentation(str(self.pptx_path))
        flags = [slide.has_notes_slide for slide in prs.slides]
        self.assertTrue(all(f is False for f in flags), msg=f"has_notes_slide: {flags}")

    def test_slide_count_matches_expand_slides_and_is_20(self):
        """swimlane の凡例自動挿入で deck.json の19枚から20枚に増える想定。"""
        from pptx import Presentation
        prs = Presentation(str(self.pptx_path))
        deck, _theme, layout = build_deck.load_deck(SAMPLE_DECK)
        expanded = build_deck.expand_slides(deck, layout)
        self.assertEqual(len(expanded["slides"]), len(prs.slides))
        self.assertEqual(len(expanded["slides"]), 20)


class ManualDeckBuildTests(unittest.TestCase):
    """ユーザーマニュアル（manual/）も同じビルダーで壊れずに通る（USER-GUIDE.html の元）。"""

    @classmethod
    def setUpClass(cls):
        cls.manual_dir = ROOT / "manual"
        cls.result = run_build(cls.manual_dir)

    def test_build_succeeds(self):
        if not (self.manual_dir / "deck.json").exists():
            self.skipTest("manual/deck.json が無い環境")
        self.assertEqual(self.result.returncode, 0, msg=self.result.stderr)


# ---------------------------------------------------------------------------
# expand_slides() の枚数計算そのもの（凡例だけの最小デッキで挙動を切り分け）
# ---------------------------------------------------------------------------

class ExpandSlidesTests(unittest.TestCase):
    def test_swimlane_with_legend_adds_one_slide(self):
        layout = build_deck.load_json(ROOT / "templates" / "layouts" / "default.json")
        deck = {
            "meta": {"id": "t"},
            "slides": [
                {"type": "title", "title": "T"},
                {
                    "type": "swimlane",
                    "title": "Flow",
                    "legend": True,
                    "legend_items": [["rect", "task", "作業", ""]],
                    "lanes": [{"name": "Lane1"}],
                    "nodes": [{"id": "n1", "lane": 0, "col": 0, "text": "Start", "shape": "terminal"}],
                    "edges": [],
                },
            ],
        }
        expanded = build_deck.expand_slides(deck, layout)
        # title(1) + 凡例自動挿入(1) + swimlane本体(1) = 3。凡例は本体の直前に挿入される。
        self.assertEqual(len(expanded["slides"]), 3)
        self.assertEqual(expanded["slides"][-2]["type"], "swimlane_legend")
        self.assertEqual(expanded["slides"][-1]["type"], "swimlane")

    def test_swimlane_without_legend_flag_does_not_add_slide(self):
        layout = build_deck.load_json(ROOT / "templates" / "layouts" / "default.json")
        deck = {
            "meta": {"id": "t"},
            "slides": [
                {
                    "type": "swimlane",
                    "title": "Flow",
                    "legend": False,
                    "lanes": [{"name": "Lane1"}],
                    "nodes": [{"id": "n1", "lane": 0, "col": 0, "text": "Start", "shape": "terminal"}],
                    "edges": [],
                },
            ],
        }
        expanded = build_deck.expand_slides(deck, layout)
        self.assertEqual(len(expanded["slides"]), 1)


# ---------------------------------------------------------------------------
# テーマ: extends 継承と派生トークン（spec.md の mix() 式どおりに導出されるか）
# ---------------------------------------------------------------------------

def _hex_to_rgb(h):
    h = h.lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    return tuple(int(h[i:i + 2], 16) for i in (0, 2, 4))


def _mix_hex(a, b, t):
    ra, ga, ba = _hex_to_rgb(a)
    rb, gb, bb = _hex_to_rgb(b)
    return (
        round(ra * (1 - t) + rb * t),
        round(ga * (1 - t) + gb * t),
        round(ba * (1 - t) + bb * t),
    )


def _assert_close_rgb(test, actual_hex, expected_rgb, msg, tol=2):
    actual_rgb = _hex_to_rgb(actual_hex)
    diff = max(abs(a - b) for a, b in zip(actual_rgb, expected_rgb))
    test.assertLessEqual(diff, tol, msg=f"{msg}: actual={actual_hex} expected~=#{expected_rgb} diff={diff}")


class ThemeExtendsAndDerivationTests(unittest.TestCase):
    """extends 継承と、明示していないトークンの基本色からの自動導出（spec.md）を検証する。"""

    def setUp(self):
        self.tmp = tempfile.TemporaryDirectory()
        self.themes_dir = Path(self.tmp.name) / "themes"
        self.themes_dir.mkdir()
        self.default_theme = build_deck.load_theme("default", build_deck.theme_dirs(None))

    def tearDown(self):
        self.tmp.cleanup()

    def _load_child(self, colors, fonts=None):
        data = {"name": "child", "colors": colors, "fonts": fonts or {}}
        (self.themes_dir / "child.json").write_text(json.dumps(data, ensure_ascii=False), encoding="utf-8")
        dirs = [self.themes_dir] + build_deck.theme_dirs(None)
        return build_deck.load_theme("child", dirs)

    def _load_child_extending(self, parent_name, name, colors, fonts=None):
        data = {"name": name, "extends": parent_name, "colors": colors, "fonts": fonts or {}}
        (self.themes_dir / f"{name}.json").write_text(json.dumps(data, ensure_ascii=False), encoding="utf-8")
        dirs = [self.themes_dir] + build_deck.theme_dirs(None)
        return build_deck.load_theme(name, dirs)

    def test_explicit_token_is_kept(self):
        theme = self._load_child({"primary": "#123456"})
        self.assertEqual(theme["colors"]["primary"], "#123456")

    def test_unspecified_token_is_derived_not_inherited_verbatim(self):
        """primary だけ変えたら、table_header_bg は default 固定値ではなく新しい primary から追随する（M-81）。"""
        theme = self._load_child({"primary": "#123456"})
        self.assertNotEqual(theme["colors"]["table_header_bg"], self.default_theme["colors"]["table_header_bg"])
        self.assertEqual(theme["colors"]["table_header_bg"], "#123456")
        self.assertEqual(theme["colors"]["table_header_text"], theme["colors"]["on_primary"])

    def test_on_primary_soft_and_muted_formulas(self):
        primary = "#123456"
        theme = self._load_child({"primary": primary})
        bg = theme["colors"]["background"]
        _assert_close_rgb(self, theme["colors"]["on_primary_soft"], _mix_hex(primary, "#FFFFFF", 0.90),
                           "on_primary_soft = mix(primary, white, 0.90)")
        _assert_close_rgb(self, theme["colors"]["on_primary_muted"], _mix_hex(primary, "#FFFFFF", 0.75),
                           "on_primary_muted = mix(primary, white, 0.75)")
        _assert_close_rgb(self, theme["colors"]["code_bg"], _mix_hex(primary, "#000000", 0.55),
                           "code_bg = mix(primary, black, 0.55)")
        # 面の視認性対応（surface がディスプレイによって背景に溶けて見えなくなるのを防ぐため
        # 0.93 -> 0.88 に濃くした。check_theme.py の surface/background 面チェック参照）。
        _assert_close_rgb(self, theme["colors"]["surface"], _mix_hex(primary, bg, 0.88),
                           "surface(未指定) = mix(primary, background, 0.88)")

    def test_muted_and_border_and_highlight_formulas(self):
        theme = self._load_child({"primary": "#123456", "text": "#101018", "accent": "#00AACC"})
        bg = theme["colors"]["background"]
        text = theme["colors"]["text"]
        accent = theme["colors"]["accent"]
        _assert_close_rgb(self, theme["colors"]["muted"], _mix_hex(text, bg, 0.35),
                           "muted(未指定) = mix(text, background, 0.35)")
        # 面の視認性対応: border/highlight_fill も同様に濃い側へ調整（0.82->0.72 / 0.85->0.75）。
        _assert_close_rgb(self, theme["colors"]["border"], _mix_hex(text, bg, 0.72),
                           "border = mix(text, background, 0.72)")
        _assert_close_rgb(self, theme["colors"]["highlight_fill"], _mix_hex(accent, bg, 0.75),
                           "highlight_fill = mix(accent, background, 0.75)")
        _assert_close_rgb(self, theme["colors"]["accent_on_primary"], _mix_hex(accent, "#FFFFFF", 0.45),
                           "accent_on_primary = mix(accent, white, 0.45)")

    def test_default_theme_itself_is_unaffected_by_derivation(self):
        """default 自身は全トークンを明示しているので、派生規則が誤って上書きしてはいけない。"""
        theme = build_deck.load_theme("default", build_deck.theme_dirs(None))
        raw = json.loads((ROOT / "templates" / "themes" / "default.json").read_text(encoding="utf-8-sig"))
        for key, val in raw["colors"].items():
            self.assertEqual(theme["colors"][key], val, msg=f"default.colors.{key} が変化しています")

    def test_extends_non_default_theme_still_derives_unspecified_tokens(self):
        """extends 先が default 以外（accenture-purple）でも、そのテーマが全トークンを明示して
        いる場合に子テーマの primary/accent 変更だけで派生トークンが追随すること（P0 回帰）。
        accenture-purple.json は18トークン全てを明示しているため、旧実装では子テーマの
        table_header_bg 等が accenture-purple の固定値（#460073 系）のまま変化しなかった。"""
        purple = build_deck.load_theme("accenture-purple", build_deck.theme_dirs(None))
        new_primary = "#123456"
        new_accent = "#00AACC"
        theme = self._load_child_extending("accenture-purple", "child-of-purple",
                                            {"primary": new_primary, "accent": new_accent})
        self.assertEqual(theme["colors"]["primary"], new_primary)
        # table_header_bg は primary の固定値。親(accenture-purple)の値のまま固まってはいけない。
        self.assertNotEqual(theme["colors"]["table_header_bg"], purple["colors"]["table_header_bg"])
        self.assertEqual(theme["colors"]["table_header_bg"], new_primary)
        # surface / muted / highlight_fill / border / code_bg / accent_on_primary も同様に
        # 親の固定値からではなく、新しい基本色から再導出されること。
        bg = theme["colors"]["background"]
        text = theme["colors"]["text"]
        self.assertNotEqual(theme["colors"]["surface"], purple["colors"]["surface"])
        _assert_close_rgb(self, theme["colors"]["surface"], _mix_hex(new_primary, bg, 0.88),
                           "surface(未指定) = mix(新primary, background, 0.88)")
        self.assertNotEqual(theme["colors"]["highlight_fill"], purple["colors"]["highlight_fill"])
        _assert_close_rgb(self, theme["colors"]["highlight_fill"], _mix_hex(new_accent, bg, 0.75),
                           "highlight_fill(未指定) = mix(新accent, background, 0.75)")
        self.assertNotEqual(theme["colors"]["code_bg"], purple["colors"]["code_bg"])
        _assert_close_rgb(self, theme["colors"]["code_bg"], _mix_hex(new_primary, "#000000", 0.55),
                           "code_bg(未指定) = mix(新primary, black, 0.55)")
        # 親テーマ(accenture-purple)自身の値は変化しない。
        self.assertEqual(purple["colors"]["table_header_bg"], "#460073")

    def test_extends_circular_reference_fails(self):
        (self.themes_dir / "a.json").write_text(
            json.dumps({"name": "a", "extends": "b", "colors": {}}), encoding="utf-8")
        (self.themes_dir / "b.json").write_text(
            json.dumps({"name": "b", "extends": "a", "colors": {}}), encoding="utf-8")
        dirs = [self.themes_dir] + build_deck.theme_dirs(None)
        with self.assertRaises(SystemExit):
            build_deck.load_theme("a", dirs)


# ---------------------------------------------------------------------------
# layout_overrides: common.brand と slides[i].style の重ね合わせ（resolve_style）
# ---------------------------------------------------------------------------

class LayoutOverridesTests(unittest.TestCase):
    def setUp(self):
        self.layout = build_deck.load_json(ROOT / "templates" / "layouts" / "default.json")

    def test_common_brand_override_applies_to_every_type(self):
        deck = minimal_deck()
        deck["meta"]["layout_overrides"] = {"common": {"brand": {"text": "ACME"}}}
        st = build_deck.resolve_style(self.layout, deck, deck["slides"][0])
        self.assertIn("brand", st)
        self.assertEqual(st["brand"]["text"], "ACME")

    def test_slide_style_overrides_layout_overrides(self):
        deck = minimal_deck()
        deck["meta"]["layout_overrides"] = {"title": {"title": {"size": 50}}}
        deck["slides"][0]["style"] = {"title": {"size": 40}}
        st = build_deck.resolve_style(self.layout, deck, deck["slides"][0])
        self.assertEqual(st["title"]["size"], 40)

    def test_layout_overrides_alone_applies_when_slide_has_no_style(self):
        deck = minimal_deck()
        deck["meta"]["layout_overrides"] = {"title": {"title": {"size": 50}}}
        st = build_deck.resolve_style(self.layout, deck, deck["slides"][0])
        self.assertEqual(st["title"]["size"], 50)


# ---------------------------------------------------------------------------
# validate_deck(): 代表的なエラー・成功系（CLI経由でユーザー操作を模す）
# ---------------------------------------------------------------------------

class ValidateDeckErrorTests(unittest.TestCase):
    def setUp(self):
        self.tmp = tempfile.TemporaryDirectory()
        self.tmp_path = Path(self.tmp.name)

    def tearDown(self):
        self.tmp.cleanup()

    def test_missing_type_fails_with_slide_number_in_message(self):
        deck = minimal_deck()
        deck["slides"] = [{"title": "type無し"}]
        deck_dir = write_deck(self.tmp_path, deck)
        result = run_build(deck_dir)
        self.assertNotEqual(result.returncode, 0)
        self.assertIn("type", result.stderr)
        self.assertIn("1枚目", result.stderr)

    def test_unknown_color_token_fails_with_friendly_message(self):
        deck = minimal_deck()
        deck["slides"] = [{"type": "bullets", "title": "T", "bullets": ["a"],
                            "style": {"title": {"color": "primry"}}}]
        deck_dir = write_deck(self.tmp_path, deck)
        result = run_build(deck_dir)
        self.assertNotEqual(result.returncode, 0)
        self.assertIn("不明な色指定", result.stderr)
        self.assertIn("primry", result.stderr)

    def test_table_row_cell_count_mismatch_fails(self):
        deck = minimal_deck()
        deck["slides"] = [{"type": "table", "title": "T", "columns": ["A", "B"], "rows": [["1"]]}]
        deck_dir = write_deck(self.tmp_path, deck)
        result = run_build(deck_dir)
        self.assertNotEqual(result.returncode, 0)
        self.assertIn("columns", result.stderr)

    def test_image_path_outside_deck_dir_fails(self):
        deck = minimal_deck()
        deck["slides"] = [{"type": "image", "title": "T", "path": "../../../etc/passwd.png"}]
        deck_dir = write_deck(self.tmp_path, deck)
        result = run_build(deck_dir)
        self.assertNotEqual(result.returncode, 0)

    def test_image_path_missing_file_fails(self):
        deck = minimal_deck()
        deck["slides"] = [{"type": "image", "title": "T", "path": "assets/does-not-exist.png"}]
        deck_dir = write_deck(self.tmp_path, deck)
        result = run_build(deck_dir)
        self.assertNotEqual(result.returncode, 0)

    def test_bom_prefixed_deck_json_still_builds(self):
        """M-70: BOM付きUTF-8 は utf-8-sig で読み、クラッシュせずビルドできる。"""
        deck_dir = write_deck(self.tmp_path, minimal_deck(), bom=True)
        result = run_build(deck_dir)
        self.assertEqual(result.returncode, 0, msg=result.stderr)
        self.assertTrue((deck_dir / "build" / "t.html").exists())


if __name__ == "__main__":
    unittest.main()
