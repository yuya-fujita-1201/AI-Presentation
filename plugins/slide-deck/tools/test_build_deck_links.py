import json
import subprocess
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
BUILD_DECK = ROOT / "tools" / "build_deck.py"


class TableUrlLinkTest(unittest.TestCase):
    def test_complete_http_url_cells_are_clickable_in_html_and_pptx(self):
        with tempfile.TemporaryDirectory() as tmp:
            deck_dir = Path(tmp) / "link-deck"
            deck_dir.mkdir()
            url = "https://example.com/source?a=1&b=2"
            title_url = "https://example.com/linked-title"
            deck = {
                "meta": {
                    "id": "link-deck",
                    "title": "Link deck",
                    "theme": "accenture-purple",
                    "layout": "default",
                },
                "slides": [
                    {
                        "type": "table",
                        "title": "Sources",
                        "columns": ["Title", "URL"],
                        "rows": [
                            ["Source", url],
                            [{"text": "Linked title", "url": title_url}, "URLなし"],
                        ],
                    }
                ],
            }
            (deck_dir / "deck.json").write_text(
                json.dumps(deck, ensure_ascii=False), encoding="utf-8"
            )

            subprocess.run(
                ["python3", str(BUILD_DECK), str(deck_dir)],
                cwd=ROOT,
                check=True,
                capture_output=True,
                text=True,
            )

            html = (deck_dir / "build" / "link-deck.html").read_text(encoding="utf-8")
            self.assertIn(
                'href="https://example.com/source?a=1&amp;b=2"',
                html,
            )
            self.assertIn('target="_blank" rel="noopener noreferrer"', html)
            self.assertNotIn('href="URLなし"', html)
            self.assertIn(
                '<a href="https://example.com/linked-title"',
                html,
            )
            self.assertIn('>Linked title</a>', html)

            prs = Presentation(deck_dir / "build" / "link-deck.pptx")
            links = []
            for shape in prs.slides[0].shapes:
                if not getattr(shape, "has_table", False):
                    continue
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            links.extend(
                                run.hyperlink.address
                                for run in paragraph.runs
                                if run.hyperlink.address
                            )
            self.assertEqual([url, title_url], links)


if __name__ == "__main__":
    unittest.main()
