#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""diagram_engine.py（図解エンジン）と図解タイプのビルド統合の回帰テスト（unittest）。

対象:
  - 文字幅推定・折り返し（全角 1.0 / 半角 0.55）
  - 格子ルーター: 直線 / 障害物回避 / 戻り線 / via / route:straight の横切り診断 / 辺指定 / ポート分散 / 双方向
  - ラベル配置: label_box の有無、label_at の尊重
  - 診断: duplicate-id / unknown-endpoint / cell-collision / invalid-position / too-dense / group-leak / out-of-grid
  - 凡例の auto 判定、stages / lanes の見出し帯
  - sequence: 順序、自己メッセージ、activations / segments の index 解決、too-dense
  - 決定論性（同じ入力 → 同じ出力）
  - 4 タイプ＋swimlane を含むデッキが HTML / PPTX ともビルドでき、PPTX にノートが無いこと
  - check_diagram.py が exit 0 / --strict で警告時に exit 1 を返すこと

実行:
    python -m unittest plugins/slide-deck/tools/test_diagram_engine.py
"""
from __future__ import annotations

import json
import subprocess
import sys
import tempfile
import unittest
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
TOOLS = ROOT / "tools"
sys.path.insert(0, str(TOOLS))
import diagram_engine as de  # noqa: E402

AREA = {"x": 72, "y": 190, "w": 1136, "h": 468, "node_wr": 0.78, "node_hr": 0.62,
        "label_size": 13, "node_size": 14, "sub_size": 11, "node_pad": 8, "icon_size": 16}


def grid_and_boxes(nodes, rows=None, cols=None):
    """テスト用: nodes[{id,row,col}] から Grid と箱を作る。"""
    rows = rows or (max(n["row"] for n in nodes) + 1)
    cols = cols or (max(n["col"] for n in nodes) + 1)
    g = de.Grid(AREA["x"], AREA["y"], AREA["w"], AREA["h"], rows, cols)
    boxes = {}
    for n in nodes:
        b = g.node_box(n["row"], n["col"])
        b.update({"id": n["id"], "label": n.get("label", n["id"]), "row": n["row"], "col": n["col"],
                  "cx": b["x"] + b["w"] / 2, "cy": b["y"] + b["h"] / 2})
        boxes[n["id"]] = b
    return g, boxes


def is_orthogonal(points):
    for (x1, y1), (x2, y2) in zip(points, points[1:]):
        if abs(x1 - x2) > 1e-6 and abs(y1 - y2) > 1e-6:
            return False
    return True


def codes(diags):
    return [d["code"] for d in diags]


class TextTests(unittest.TestCase):
    def test_width_cjk_vs_ascii(self):
        self.assertAlmostEqual(de.text_width("日本", 10), 20.0)
        self.assertAlmostEqual(de.text_width("ab", 10), 11.0)

    def test_wrap_respects_width_and_newline(self):
        self.assertEqual(de.wrap_text("あいうえお", 10, 30), ["あいう", "えお"])
        self.assertEqual(de.wrap_text("a\nb", 10, 100), ["a", "b"])

    def test_fit_font_size_shrinks_to_min(self):
        size, lines = de.fit_font_size("とても長いラベルの文字列です", 14, 60, 30, min_size=10)
        self.assertGreaterEqual(size, 10)
        self.assertLessEqual(len(lines) * size * 1.15, 30 + 0.5 + 14 * 1.15)  # 最小サイズ到達時は超過を許容


class RouterTests(unittest.TestCase):
    def test_same_row_adjacent_is_straight(self):
        g, boxes = grid_and_boxes([{"id": "a", "row": 0, "col": 0}, {"id": "b", "row": 0, "col": 1}])
        routed = de.route_all(g, boxes, [{"from": "a", "to": "b"}])
        pts = routed[0]["points"]
        self.assertEqual(len(pts), 2)
        self.assertAlmostEqual(pts[0][1], pts[1][1])
        self.assertAlmostEqual(pts[0][0], boxes["a"]["x"] + boxes["a"]["w"])
        self.assertAlmostEqual(pts[1][0], boxes["b"]["x"])

    def test_same_col_is_vertical(self):
        g, boxes = grid_and_boxes([{"id": "a", "row": 0, "col": 0}, {"id": "b", "row": 1, "col": 0}])
        pts = de.route_all(g, boxes, [{"from": "a", "to": "b"}])[0]["points"]
        self.assertEqual(len(pts), 2)
        self.assertAlmostEqual(pts[0][0], pts[1][0])

    def test_obstacle_between_same_row_is_avoided(self):
        nodes = [{"id": "a", "row": 0, "col": 0}, {"id": "x", "row": 0, "col": 1}, {"id": "b", "row": 0, "col": 2}]
        g, boxes = grid_and_boxes(nodes)
        diags = []
        routed = de.route_all(g, boxes, [{"from": "a", "to": "b"}], diags=diags)
        pts = routed[0]["points"]
        self.assertTrue(is_orthogonal(pts))
        self.assertGreater(len(pts), 2, "障害物を迂回して折れ線になるはず")
        for p, q in zip(pts, pts[1:]):
            self.assertFalse(de.segment_hits_rect(p, q, boxes["x"]), "中間ノード x を横切ってはいけない")
        self.assertNotIn("edge-through-node", codes(diags))

    def test_back_edge_routes_without_crossing_nodes(self):
        nodes = [{"id": "a", "row": 1, "col": 0}, {"id": "m", "row": 1, "col": 1}, {"id": "b", "row": 1, "col": 2},
                 {"id": "f", "row": 2, "col": 1}]
        g, boxes = grid_and_boxes(nodes)
        diags = []
        routed = de.route_all(g, boxes, [{"from": "a", "to": "m"}, {"from": "m", "to": "b"},
                                         {"from": "b", "to": "f"}, {"from": "f", "to": "a"}], diags=diags)
        self.assertEqual(len(routed), 4)
        self.assertNotIn("edge-through-node", codes(diags))
        for e in routed:
            self.assertTrue(is_orthogonal(e["points"]))

    def test_route_straight_through_obstacle_is_diagnosed(self):
        nodes = [{"id": "a", "row": 0, "col": 0}, {"id": "x", "row": 0, "col": 1}, {"id": "b", "row": 0, "col": 2}]
        g, boxes = grid_and_boxes(nodes)
        diags = []
        de.route_all(g, boxes, [{"from": "a", "to": "b", "route": "straight"}], diags=diags)
        self.assertIn("edge-through-node", codes(diags))
        d = [d for d in diags if d["code"] == "edge-through-node"][0]
        self.assertEqual(d["level"], "error")
        self.assertTrue(d["fixes"], "修理指示（fixes）が付く")
        self.assertEqual(d["evidence"]["node"]["id"], "x")

    def test_via_is_orthogonalized(self):
        g, boxes = grid_and_boxes([{"id": "a", "row": 0, "col": 0}, {"id": "b", "row": 1, "col": 2}])
        routed = de.route_all(g, boxes, [{"from": "a", "to": "b", "via": [[600, 250]]}])
        pts = routed[0]["points"]
        self.assertTrue(is_orthogonal(pts))
        self.assertTrue(any(abs(x - 600) < 1e-6 for x, _ in pts) or any(abs(y - 250) < 1e-6 for _, y in pts))

    def test_from_side_and_to_side_are_honored(self):
        g, boxes = grid_and_boxes([{"id": "a", "row": 0, "col": 0}, {"id": "b", "row": 0, "col": 2}])
        routed = de.route_all(g, boxes, [{"from": "a", "to": "b", "from_side": "bottom", "to_side": "bottom"}])
        pts = routed[0]["points"]
        self.assertAlmostEqual(pts[0][1], boxes["a"]["y"] + boxes["a"]["h"], delta=0.01)
        self.assertAlmostEqual(pts[-1][1], boxes["b"]["y"] + boxes["b"]["h"], delta=0.01)

    def test_port_spread_separates_two_edges_from_same_side(self):
        nodes = [{"id": "a", "row": 0, "col": 0}, {"id": "b", "row": 0, "col": 1}, {"id": "c", "row": 1, "col": 1}]
        g, boxes = grid_and_boxes(nodes)
        routed = de.route_all(g, boxes, [{"from": "a", "to": "b"}, {"from": "a", "to": "c"}])
        starts = sorted(e["points"][0] for e in routed)
        # 2 本とも a の右辺から出るなら y がずれている（同じ点から重なって出ない）
        if all(abs(p[0] - (boxes["a"]["x"] + boxes["a"]["w"])) < 1e-6 for p in starts):
            self.assertNotAlmostEqual(starts[0][1], starts[1][1])

    def test_bidirectional_pair_becomes_parallel_lines(self):
        g, boxes = grid_and_boxes([{"id": "a", "row": 0, "col": 0}, {"id": "b", "row": 0, "col": 1}])
        routed = de.route_all(g, boxes, [{"from": "a", "to": "b"}, {"from": "b", "to": "a"}])
        y0 = routed[0]["points"][0][1]
        y1 = routed[1]["points"][0][1]
        self.assertGreater(abs(y0 - y1), 3.0, "双方向の線は重ならず並行にずれる")

    def test_self_loop_has_points(self):
        g, boxes = grid_and_boxes([{"id": "a", "row": 0, "col": 0}])
        routed = de.route_all(g, boxes, [{"from": "a", "to": "a", "label": "再試行"}])
        self.assertGreaterEqual(len(routed[0]["points"]), 4)

    def test_label_box_present_and_label_at_honored(self):
        g, boxes = grid_and_boxes([{"id": "a", "row": 0, "col": 0}, {"id": "b", "row": 0, "col": 2}])
        routed = de.route_all(g, boxes, [{"from": "a", "to": "b", "label": "HTTPS"},
                                         {"from": "b", "to": "a", "label": "ack", "label_at": [500, 400]}])
        self.assertIsNotNone(routed[0]["label_box"])
        self.assertAlmostEqual(routed[1]["label_box"]["cx"], 500)
        self.assertAlmostEqual(routed[1]["label_box"]["cy"], 400)

    def test_unknown_endpoint_edges_are_skipped(self):
        g, boxes = grid_and_boxes([{"id": "a", "row": 0, "col": 0}])
        routed = de.route_all(g, boxes, [{"from": "a", "to": "zzz"}])
        self.assertEqual(routed, [])

    def test_deterministic(self):
        nodes = [{"id": f"n{i}", "row": i % 3, "col": i // 3} for i in range(9)]
        edges = [{"from": f"n{i}", "to": f"n{(i * 4 + 1) % 9}", "label": f"e{i}"} for i in range(9) if i != (i * 4 + 1) % 9]
        outs = []
        for _ in range(2):
            g, boxes = grid_and_boxes(nodes)
            routed = de.route_all(g, boxes, edges)
            outs.append(json.dumps([(e["from"], e["to"], e["points"], e["label_box"]) for e in routed], sort_keys=True))
        self.assertEqual(outs[0], outs[1])


class ValidationTests(unittest.TestCase):
    def slide(self, **over):
        s = {"type": "architecture",
             "nodes": [{"id": "a", "label": "A", "row": 0, "col": 0, "type": "backend"},
                       {"id": "b", "label": "B", "row": 0, "col": 1, "type": "database"}],
             "edges": [{"from": "a", "to": "b"}]}
        s.update(over)
        return s

    def test_clean_slide_has_no_errors(self):
        self.assertEqual([d for d in de.validate_grid_diagram(self.slide(), "architecture") if d["level"] == "error"], [])

    def test_duplicate_id(self):
        s = self.slide(nodes=[{"id": "a", "label": "A", "row": 0, "col": 0}, {"id": "a", "label": "A2", "row": 0, "col": 1}])
        self.assertIn("duplicate-id", codes(de.validate_grid_diagram(s, "architecture")))

    def test_unknown_endpoint(self):
        s = self.slide(edges=[{"from": "a", "to": "nope"}])
        self.assertIn("unknown-endpoint", codes(de.validate_grid_diagram(s, "architecture")))

    def test_cell_collision(self):
        s = self.slide(nodes=[{"id": "a", "label": "A", "row": 0, "col": 0}, {"id": "b", "label": "B", "row": 0, "col": 0}])
        d = [d for d in de.validate_grid_diagram(s, "architecture") if d["code"] == "cell-collision"]
        self.assertTrue(d and d[0]["level"] == "error")

    def test_invalid_position_types(self):
        s = self.slide(nodes=[{"id": "a", "label": "A", "row": "0", "col": 0}, {"id": "b", "label": "B", "row": 0, "col": -1}])
        self.assertEqual(codes(de.validate_grid_diagram(s, "architecture")).count("invalid-position"), 2)

    def test_unknown_type_and_variant_are_warnings(self):
        s = self.slide(nodes=[{"id": "a", "label": "A", "row": 0, "col": 0, "type": "server", "variant": "bold"}], edges=[])
        diags = de.validate_grid_diagram(s, "architecture")
        self.assertIn("unknown-type", codes(diags))
        self.assertIn("unknown-variant", codes(diags))
        self.assertTrue(all(d["level"] == "warning" for d in diags))

    def test_too_dense(self):
        nodes = [{"id": f"n{i}", "label": str(i), "row": i % 3, "col": i // 3} for i in range(13)]
        self.assertIn("too-dense", codes(de.validate_grid_diagram(self.slide(nodes=nodes, edges=[]), "architecture")))

    def test_group_unknown_member_is_error(self):
        s = self.slide(groups=[{"label": "G", "nodes": ["a", "zzz"]}])
        d = [d for d in de.validate_grid_diagram(s, "architecture") if d["code"] == "unknown-endpoint"]
        self.assertTrue(d and d[0]["level"] == "error")

    def test_lifecycle_uses_states_and_lane_range(self):
        s = {"type": "lifecycle", "lanes": [{"name": "main"}],
             "states": [{"id": "s1", "label": "S1", "kind": "start", "lane": 0, "col": 0},
                        {"id": "s2", "label": "S2", "kind": "weird", "lane": 3, "col": 1}],
             "transitions": [{"from": "s1", "to": "s2"}]}
        diags = de.validate_grid_diagram(s, "lifecycle")
        self.assertIn("unknown-kind", codes(diags))
        self.assertIn("invalid-position", codes(diags))


class GridLayoutTests(unittest.TestCase):
    def test_legend_auto_needs_two_types(self):
        one = {"type": "architecture", "nodes": [{"id": "a", "label": "A", "row": 0, "col": 0, "type": "backend"},
                                                 {"id": "b", "label": "B", "row": 0, "col": 1, "type": "backend"}]}
        geo = de.layout_grid_diagram(one, "architecture", AREA)
        self.assertEqual(geo["legend"], [])
        two = dict(one, nodes=[dict(one["nodes"][0]), dict(one["nodes"][1], type="database")])
        geo2 = de.layout_grid_diagram(two, "architecture", AREA)
        self.assertEqual([i["type"] for i in geo2["legend"]], ["backend", "database"])
        geo3 = de.layout_grid_diagram(dict(two, legend=False), "architecture", AREA)
        self.assertEqual(geo3["legend"], [])

    def test_stages_make_col_headers_and_lanes_make_row_headers(self):
        df = {"type": "dataflow", "stages": ["A", "B", "C"],
              "nodes": [{"id": "x", "label": "X", "col": 0, "row": 0}, {"id": "y", "label": "Y", "stage": 2, "row": 0}],
              "edges": [{"from": "x", "to": "y", "label": "data", "classification": "PII"}]}
        geo = de.layout_grid_diagram(df, "dataflow", AREA)
        self.assertEqual(len(geo["col_headers"]), 3)
        self.assertEqual(geo["cols"], 3)
        self.assertEqual(geo["nodes"]["y"]["col"], 2)
        lc = {"type": "lifecycle", "lanes": [{"name": "L1"}, {"name": "L2"}],
              "states": [{"id": "s", "label": "S", "kind": "start", "lane": 1, "col": 0}]}
        geo2 = de.layout_grid_diagram(lc, "lifecycle", AREA)
        self.assertEqual(len(geo2["row_headers"]), 2)
        self.assertEqual(geo2["nodes"]["s"]["row"], 1)
        self.assertEqual(geo2["nodes"]["s"]["shape"], "pill")

    def test_out_of_grid_expands_with_warning(self):
        s = {"type": "architecture", "cols": 2,
             "nodes": [{"id": "a", "label": "A", "row": 0, "col": 0}, {"id": "b", "label": "B", "row": 0, "col": 4}]}
        diags = []
        geo = de.layout_grid_diagram(s, "architecture", AREA, diags)
        self.assertEqual(geo["cols"], 5)
        self.assertIn("out-of-grid", codes(diags))

    def test_group_leak_warning(self):
        s = {"type": "architecture",
             "nodes": [{"id": "a", "label": "A", "row": 0, "col": 0}, {"id": "b", "label": "B", "row": 0, "col": 2},
                       {"id": "m", "label": "M", "row": 0, "col": 1}],
             "groups": [{"label": "G", "nodes": ["a", "b"]}]}
        diags = []
        geo = de.layout_grid_diagram(s, "architecture", AREA, diags)
        self.assertEqual(len(geo["groups"]), 1)
        self.assertIn("group-leak", codes(diags))

    def test_node_text_overflow_warning(self):
        s = {"type": "architecture", "cols": 6, "rows": 4,
             "nodes": [{"id": "a", "label": "とても長いコンポーネント名がここに入ります", "sublabel": "さらに長い補足説明の文章", "row": 0, "col": 0}]}
        diags = []
        de.layout_grid_diagram(s, "architecture", AREA, diags)
        self.assertIn("node-text-overflow", codes(diags))


class SequenceTests(unittest.TestCase):
    def base(self):
        return {"type": "sequence",
                "participants": [{"id": "u", "label": "User"}, {"id": "w", "label": "Web"}, {"id": "a", "label": "API"}],
                "messages": [{"id": "m0", "from": "u", "to": "w", "label": "open"},
                             {"id": "m1", "from": "w", "to": "a", "label": "GET"},
                             {"id": "m2", "from": "a", "to": "a", "label": "cache"},
                             {"id": "m3", "from": "a", "to": "w", "label": "200", "variant": "return"}],
                "activations": [{"participant": "a", "from": "m1", "to": 3}],
                "segments": [{"from": 0, "to": "m1", "label": "req"}]}

    def test_messages_in_order_and_self_message(self):
        geo = de.layout_sequence(self.base(), AREA)
        ys = [m["y"] for m in geo["messages"]]
        self.assertEqual(ys, sorted(ys))
        self_msg = geo["messages"][2]
        self.assertEqual(len(self_msg["points"]), 4)

    def test_activation_and_segment_indices(self):
        geo = de.layout_sequence(self.base(), AREA)
        self.assertEqual(len(geo["activations"]), 1)
        a = geo["activations"][0]
        self.assertLess(a["y"], geo["messages"][1]["y"])
        self.assertGreater(a["y"] + a["h"], geo["messages"][3]["y"])
        self.assertEqual(len(geo["segments"]), 1)

    def test_unknown_participant_is_error(self):
        s = self.base()
        s["messages"].append({"from": "a", "to": "nope", "label": "x"})
        diags = []
        de.layout_sequence(s, AREA, diags)
        self.assertIn("unknown-endpoint", codes(diags))

    def test_too_many_messages_warns(self):
        s = self.base()
        s["messages"] = [{"from": "u", "to": "w", "label": f"m{i}"} for i in range(14)]
        diags = []
        de.layout_sequence(s, AREA, diags)
        self.assertIn("too-dense", codes(diags))

    def test_long_label_warns(self):
        s = self.base()
        s["participants"] = [{"id": f"p{i}", "label": f"P{i}"} for i in range(7)]
        s["messages"] = [{"from": "p0", "to": "p1", "label": "とてもとても長いメッセージラベルで矢印より長い"}]
        s["activations"] = []; s["segments"] = []
        diags = []
        de.layout_sequence(s, AREA, diags)
        self.assertIn("label-collision", codes(diags))


class SwimlaneRouterRegressionTests(unittest.TestCase):
    """見本デッキの swimlane を共通ルーターに通したときの回帰（レビュー指摘: 旧実装では t4→conn と
    t4→mail が同じ辺の同じ点から出て最初の区間が完全に重なっていた）。"""

    def test_sample_swimlane_edges_do_not_overlap_and_avoid_nodes(self):
        import build_deck
        sample = ROOT / "examples" / "template-sample"
        deck, theme, layout = build_deck.load_deck(sample)
        sw = [s for s in deck["slides"] if s.get("type") == "swimlane"]
        self.assertTrue(sw, "見本デッキに swimlane が無い")
        slide = sw[0]
        st = build_deck.resolve_style(layout, deck, slide)
        geo = build_deck.swimlane_geometry(slide, st)
        edges = geo["edges"]
        self.assertGreaterEqual(len(edges), 5)
        # 1) 無関係な線同士が同一直線上で 8px 以上重ならない（ポート分散・回廊分散の効果）
        for i in range(len(edges)):
            for j in range(i + 1, len(edges)):
                a, b = edges[i], edges[j]
                pa, pb = a["points"], b["points"]
                for k in range(len(pa) - 1):
                    for m in range(len(pb) - 1):
                        ov = de.collinear_overlap(pa[k], pa[k + 1], pb[m], pb[m + 1])
                        self.assertLess(ov, de.MIN_SEGMENT,
                                        f"{a['from']}→{a['to']} と {b['from']}→{b['to']} が {ov:.0f}px 重なっている")
        # 2) どの線も端点以外のノードを横切らない
        for e in edges:
            for p, q in zip(e["points"], e["points"][1:]):
                for nid, nd in geo["nodes"].items():
                    if nid in (e["from"], e["to"]):
                        continue
                    self.assertFalse(de.segment_hits_rect(p, q, nd), f"{e['from']}→{e['to']} が {nid} を横切る")
        # 3) 診断に error が無い（見本は常にクリーンであること）
        self.assertEqual([d for d in geo["diagnostics"] if d["level"] == "error"], [])

    def test_sample_swimlane_geometry_is_deterministic(self):
        import build_deck
        sample = ROOT / "examples" / "template-sample"
        outs = []
        for _ in range(2):
            deck, theme, layout = build_deck.load_deck(sample)
            slide = [s for s in deck["slides"] if s.get("type") == "swimlane"][0]
            st = build_deck.resolve_style(layout, deck, slide)
            geo = build_deck.swimlane_geometry(slide, st)
            outs.append(json.dumps([(e["from"], e["to"], e["points"], e["label_box"]) for e in geo["edges"]], sort_keys=True))
        self.assertEqual(outs[0], outs[1])


class BuildIntegrationTests(unittest.TestCase):
    """4 タイプ＋swimlane を含む最小デッキが HTML/PPTX ともにビルドでき、check_diagram が動く。"""

    @classmethod
    def setUpClass(cls):
        cls.tmp = tempfile.TemporaryDirectory()
        cls.deck_dir = Path(cls.tmp.name) / "d"
        cls.deck_dir.mkdir()
        deck = {"meta": {"id": "d", "title": "T"}, "slides": [
            {"type": "architecture", "title": "A",
             "nodes": [{"id": "web", "type": "frontend", "label": "Web", "row": 0, "col": 0},
                       {"id": "api", "type": "backend", "label": "API", "row": 0, "col": 1, "variant": "emphasis"},
                       {"id": "db", "type": "database", "label": "DB", "row": 0, "col": 2}],
             "groups": [{"label": "VPC", "kind": "region", "nodes": ["api", "db"]}],
             "edges": [{"from": "web", "to": "api", "label": "HTTPS"}, {"from": "api", "to": "db", "label": "SQL"}]},
            {"type": "dataflow", "title": "D", "stages": ["S1", "S2"],
             "nodes": [{"id": "a", "label": "A", "col": 0, "row": 0}, {"id": "b", "label": "B", "col": 1, "row": 0}],
             "edges": [{"from": "a", "to": "b", "label": "x", "classification": "c"}]},
            {"type": "lifecycle", "title": "L", "lanes": [{"name": "main"}, {"name": "err"}],
             "states": [{"id": "s", "kind": "start", "label": "S", "lane": 0, "col": 0, "step": "01"},
                        {"id": "r", "kind": "decision", "label": "R", "lane": 0, "col": 1},
                        {"id": "ok", "kind": "success", "label": "OK", "lane": 0, "col": 2},
                        {"id": "ng", "kind": "failure", "label": "NG", "lane": 1, "col": 1}],
             "transitions": [{"from": "s", "to": "r"}, {"from": "r", "to": "ok", "label": "Y"}, {"from": "r", "to": "ng", "label": "N"}]},
            {"type": "sequence", "title": "S",
             "participants": [{"id": "u", "label": "U", "type": "external"}, {"id": "a", "label": "A", "type": "backend"}],
             "messages": [{"from": "u", "to": "a", "label": "req"}, {"from": "a", "to": "u", "label": "res", "variant": "return"}]},
            {"type": "swimlane", "title": "W", "legend": False, "lanes": [{"name": "L1"}, {"name": "L2"}],
             "nodes": [{"id": "t1", "lane": 0, "col": 0, "shape": "task", "text": "作業"},
                       {"id": "d1", "lane": 0, "col": 1, "shape": "decision", "text": "判定"},
                       {"id": "t2", "lane": 1, "col": 2, "shape": "task", "text": "次"}],
             "edges": [{"from": "t1", "to": "d1"}, {"from": "d1", "to": "t2", "label": "Y", "variant": "emphasis"}]},
        ]}
        (cls.deck_dir / "deck.json").write_text(json.dumps(deck, ensure_ascii=False), encoding="utf-8")
        cls.result = subprocess.run([sys.executable, str(TOOLS / "build_deck.py"), str(cls.deck_dir)],
                                    cwd=ROOT, capture_output=True, text=True)

    @classmethod
    def tearDownClass(cls):
        cls.tmp.cleanup()

    def test_build_succeeds(self):
        self.assertEqual(self.result.returncode, 0, msg=self.result.stderr)
        self.assertTrue((self.deck_dir / "build" / "d.html").exists())
        self.assertTrue((self.deck_dir / "build" / "d.pptx").exists())

    def test_pptx_has_no_notes_and_reparses(self):
        from pptx import Presentation
        prs = Presentation(str(self.deck_dir / "build" / "d.pptx"))
        self.assertEqual(len(prs.slides), 5)
        self.assertTrue(all(s.has_notes_slide is False for s in prs.slides))

    def test_html_contains_diagram_markup(self):
        html = (self.deck_dir / "build" / "d.html").read_text(encoding="utf-8")
        self.assertIn("<polyline", html)
        self.assertIn("marker-end", html)

    def test_check_diagram_runs(self):
        r = subprocess.run([sys.executable, str(TOOLS / "check_diagram.py"), str(self.deck_dir)],
                           cwd=ROOT, capture_output=True, text=True)
        self.assertEqual(r.returncode, 0, msg=r.stdout + r.stderr)
        self.assertIn("check_diagram:", r.stdout)

    def test_check_diagram_strict_fails_on_warning(self):
        bad_dir = Path(self.tmp.name) / "bad"
        bad_dir.mkdir()
        deck = {"meta": {"id": "bad", "title": "B"}, "slides": [
            {"type": "architecture", "title": "A", "cols": 1,
             "nodes": [{"id": "a", "label": "A", "row": 0, "col": 0}, {"id": "b", "label": "B", "row": 0, "col": 3}]}]}
        (bad_dir / "deck.json").write_text(json.dumps(deck, ensure_ascii=False), encoding="utf-8")
        r = subprocess.run([sys.executable, str(TOOLS / "check_diagram.py"), str(bad_dir), "--strict"],
                           cwd=ROOT, capture_output=True, text=True)
        self.assertEqual(r.returncode, 1)
        self.assertIn("out-of-grid", r.stdout)

    def test_schema_error_stops_build(self):
        err_dir = Path(self.tmp.name) / "err"
        err_dir.mkdir()
        deck = {"meta": {"id": "err", "title": "E"}, "slides": [
            {"type": "architecture", "title": "A",
             "nodes": [{"id": "a", "label": "A", "row": 0, "col": 0}], "edges": [{"from": "a", "to": "ghost"}]}]}
        (err_dir / "deck.json").write_text(json.dumps(deck, ensure_ascii=False), encoding="utf-8")
        r = subprocess.run([sys.executable, str(TOOLS / "build_deck.py"), str(err_dir), "--html"],
                           cwd=ROOT, capture_output=True, text=True)
        self.assertEqual(r.returncode, 1)
        self.assertIn("unknown-endpoint", r.stderr)


if __name__ == "__main__":
    unittest.main()
