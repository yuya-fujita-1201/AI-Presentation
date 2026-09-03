#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""deck.json → HTML / PPTX ビルダー v2（パラメーター駆動）。

使い方:
    python tools/build_deck.py <deck_dir>          # HTML + PPTX
    python tools/build_deck.py <deck_dir> --html   # HTML のみ
    python tools/build_deck.py <deck_dir> --pptx   # PPTX のみ
（環境により python3 / py -3 を使う）

設計:
- 仮想キャンバス 1280x720px。全オブジェクトの位置・サイズ・フォント・色を px で持つ
- 解決順: templates/layouts/<layout>.json の既定値
          ← deck.meta.layout_overrides（デッキ全体の調整）
          ← slides[i].style（スライド個別の微修正）
- HTML は px をそのまま使用、PPTX は 96dpi 換算（px/96 インチ、フォント px*0.75 pt）
  なので両形式で座標・サイズが一致する
- スキーマは references/deck-schema.md を参照

注意: PPTX にはスピーカーノートを出力しない（python-pptx の notes_slide は
macOS Keynote の互換性を壊す既知問題があるため）。notes は HTML でのみ表示される。
"""

from __future__ import annotations

import argparse
import base64
import copy
import difflib
import hashlib
import html as html_mod
import json
import mimetypes
import os
import re
import sys
import unicodedata
from pathlib import Path
from string import Template

ROOT = Path(__file__).resolve().parent.parent
CANVAS_W, CANVAS_H = 1280, 720

# 図解タイプ（architecture / dataflow / lifecycle / sequence / swimlane の配線）共通エンジン。
# tools/ 直下の同名モジュール。スクリプト実行時は tools/ が sys.path[0] に入るが、
# 他ツールから import build_deck されるケースでも解決できるよう明示的に追加する。
sys.path.insert(0, str(Path(__file__).resolve().parent))
import diagram_engine as de  # noqa: E402


# ---------------------------------------------------------------------------
# エラー・警告の共通ヘルパー（他ツールからも `import build_deck` して再利用可）
# ---------------------------------------------------------------------------

_WARNED: set = set()


def fail(msg: str):
    """stderr に error: <msg> を出して exit code 1 で終了する。"""
    sys.exit(f"error: {msg}")


def warn(msg: str):
    """stderr に warning: <msg> を出す。ビルドは続行する。同一メッセージは1回だけ表示。"""
    if msg in _WARNED:
        return
    _WARNED.add(msg)
    print(f"warning: {msg}", file=sys.stderr)


def note(msg: str):
    """stderr に note: <msg> を出す。"""
    print(f"note: {msg}", file=sys.stderr)


def setup_console():
    """stdout/stderr を UTF-8 に強制する（Windows の既定コードページ対策）。失敗は無視。"""
    for s in (sys.stdout, sys.stderr):
        try:
            s.reconfigure(encoding="utf-8", errors="replace")
        except Exception:
            pass


# ---------------------------------------------------------------------------
# 読み込みとスタイル解決
# ---------------------------------------------------------------------------

def load_json(path: Path):
    if not path.exists():
        fail(f"{path} が見つかりません")
    try:
        return json.loads(path.read_text(encoding="utf-8-sig"))
    except json.JSONDecodeError as e:
        fail(f"{path} の JSON が不正です: {e}")
    except PermissionError:
        fail(f"{path} を読み込めません（権限がありません）")


def theme_dirs(deck_dir: Path):
    """テーマの探索ディレクトリを優先順に返す。
    1. 環境変数 SLIDE_DECK_THEMES（複数可・OS区切り）… ユーザー追加テーマ置き場
    2. <deck_dir>/themes/                           … デッキ／プロジェクト同梱テーマ
    3. 同梱テーマ（templates/themes/）              … default・accenture-purple 等
    先に見つかったものが勝つ。プラグイン更新で消したくないテーマは 1 か 2 に置く。"""
    dirs = []
    env = os.environ.get("SLIDE_DECK_THEMES")
    if env:
        dirs += [Path(p) for p in env.split(os.pathsep) if p.strip()]
    if deck_dir is not None:
        dirs.append(Path(deck_dir) / "themes")
    dirs.append(ROOT / "templates" / "themes")
    return dirs


def _hex_to_rgb(h: str):
    h = h.lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    return tuple(int(h[i:i + 2], 16) for i in (0, 2, 4))


def _rgb_to_hex(rgb) -> str:
    return "#" + "".join(f"{max(0, min(255, round(c))):02X}" for c in rgb)


def mix(a: str, b: str, t: float) -> str:
    """hex カラー a・b を (1-t):t の割合で線形混合した hex を返す。"""
    ra, ga, ba = _hex_to_rgb(a)
    rb, gb, bb = _hex_to_rgb(b)
    return _rgb_to_hex((ra + (rb - ra) * t, ga + (gb - ga) * t, ba + (bb - ba) * t))


# default 以外のテーマで明示指定が無い場合、default の固定 hex をそのまま継承せず、
# そのテーマ自身の基本色（background/surface/text/muted/primary/accent/on_primary）から
# 導出するトークン群（themes.md 参照）。
DERIVABLE_COLOR_TOKENS = (
    "surface", "muted", "table_header_bg", "table_header_text", "table_row_alt",
    "on_primary_soft", "on_primary_muted", "code_bg", "code_text",
    "highlight_fill", "border", "accent_on_primary", "heading_text",
)


def _luminance(hex_color: str) -> float:
    """WCAG 相対輝度。"""
    def ch(c):
        c = c / 255
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    r, g, bl = _hex_to_rgb(hex_color)
    return 0.2126 * ch(r) + 0.7152 * ch(g) + 0.0722 * ch(bl)


def contrast_ratio(a: str, b: str) -> float:
    """WCAG コントラスト比（1.0〜21.0）。"""
    la, lb = _luminance(a), _luminance(b)
    return (max(la, lb) + 0.05) / (min(la, lb) + 0.05)


def readable_on(color: str, bg: str, min_ratio: float = 4.5) -> str:
    """color を bg 上で min_ratio 以上のコントラストになるまで白/黒側へ寄せた色を返す。
    暗い背景なら白へ、明るい背景なら黒へ 0.1 刻みで混ぜる（既に十分ならそのまま）。"""
    if contrast_ratio(color, bg) >= min_ratio:
        return color
    toward = "#FFFFFF" if _luminance(bg) < 0.5 else "#000000"
    for i in range(1, 11):
        cand = mix(color, toward, i / 10)
        if contrast_ratio(cand, bg) >= min_ratio:
            return cand
    return toward


def _derive_colors(colors: dict, own_color_keys) -> None:
    """default 以外の階層で、その階層が自分の JSON で明示していないトークン（own_color_keys に
    無いもの）を、colors の現在値（親から継承した値に own_color_keys 分を重ねたもの＝
    「その階層自身の基本色」）から導出して colors を書き換える（in-place）。
    親テーマ（default 以外）がそのトークンを明示していても、この階層が明示していなければ
    継承値を採用せず、この階層自身の primary/accent 等から改めて導出し直す（themes.md 参照）。"""
    bg = colors.get("background", "#FFFFFF")
    text = colors.get("text", "#1A2233")
    primary = colors.get("primary", "#0F3D6E")
    accent = colors.get("accent", "#007D96")
    on_primary = colors.get("on_primary", "#FFFFFF")
    if "surface" not in own_color_keys:
        # M-面の視認性: 0.93 だと surface/background のコントラスト比が 1.1 程度しかなく
        # ディスプレイによって面が背景に溶けて見えなくなる。0.88 まで primary 側に寄せて
        # 面自体をはっきり視認できるようにする（check_theme.py の面チェック参照）。
        colors["surface"] = mix(primary, bg, 0.88)
    if "muted" not in own_color_keys:
        colors["muted"] = mix(text, bg, 0.35)
    surface = colors.get("surface", bg)
    derived = {
        "table_header_bg": primary,
        "table_header_text": on_primary,
        "table_row_alt": mix(surface, bg, 0.5),
        "on_primary_soft": mix(primary, "#FFFFFF", 0.90),
        "on_primary_muted": mix(primary, "#FFFFFF", 0.75),
        "code_bg": mix(primary, "#000000", 0.55),
        "code_text": "#E6EDF5",
        # M-面の視認性: highlight_fill / border も同様に濃い側へ調整（0.85→0.75 / 0.82→0.72）。
        "highlight_fill": mix(accent, bg, 0.75),
        "border": mix(text, bg, 0.72),
        "accent_on_primary": mix(accent, "#FFFFFF", 0.45),
        # 見出し文字色: primary が背景上で読めなければ白/黒側へ寄せる（暗色テーマ対策）
        "heading_text": readable_on(primary, bg),
    }
    for key, value in derived.items():
        if key not in own_color_keys:
            colors[key] = value


def _load_theme_chain(name: str, dirs, _seen):
    """テーマ JSON を extends チェーンに沿って読み込み、マージ済み dict を返す。
    default 以外の各階層は、親（既にこの関数で解決済み＝派生も適用済み）を deep_merge した直後、
    自分の JSON ファイルが明示していない colors キーだけをその階層自身の基本色から再導出する。
    階層ごとに導出をやり直すことで、`extends` 先が default 以外のテーマ（例: accenture-purple）
    であっても、子テーマが primary/accent だけを上書きすれば派生トークン（table_header_bg 等）
    が親の固定値のままにならず新しい基本色に追随する。"""
    if name in _seen:
        fail(f"テーマの extends が循環しています: {' -> '.join(_seen + [name])}")
    path = next((d / f"{name}.json" for d in dirs if (d / f"{name}.json").exists()), None)
    if path is None:
        searched = ", ".join(str(d) for d in dirs)
        fail(f"テーマ '{name}' が見つかりません（探索: {searched}）")
    own = load_json(path)
    own_color_keys = set(own.get("colors", {}).keys())
    parent = own.get("extends") or (None if name == "default" else "default")
    if parent:
        parent_data = _load_theme_chain(parent, dirs, _seen + [name])
        merged = deep_merge(parent_data, own)
    else:
        merged = copy.deepcopy(own)
    merged["name"] = name
    if name != "default":
        _derive_colors(merged.setdefault("colors", {}), own_color_keys)
    return merged


def load_theme(name: str, dirs, _seen=None):
    """テーマを解決して返す。default.json（または extends 先）をベースにマージするため、
    新しいテーマは上書きしたいトークンだけ書けばよい。`extends` を指定すると任意のテーマを
    親にできる（未指定なら default が親）。探索は theme_dirs() の順。
    default 以外のテーマが明示していない派生トークンは、default の固定 hex ではなく
    そのテーマ自身の基本色から自動導出する（themes.md 参照）。"""
    return _load_theme_chain(name, dirs, _seen or [])


def load_deck(deck_dir: Path):
    deck = load_json(deck_dir / "deck.json")
    if not isinstance(deck, dict):
        fail("deck.json の形式が不正です（トップレベルはオブジェクトである必要があります）")
    meta = deck.get("meta", {})
    if not isinstance(meta, dict):
        fail("deck.json の形式が不正です（meta はオブジェクトである必要があります）")
    theme = load_theme(meta.get("theme", "default"), theme_dirs(deck_dir))
    layout = load_json(ROOT / "templates" / "layouts" / f"{meta.get('layout', 'default')}.json")
    return deck, theme, layout


def deep_merge(base, over):
    if not isinstance(base, dict) or not isinstance(over, dict):
        return copy.deepcopy(over)
    out = copy.deepcopy(base)
    for k, v in over.items():
        out[k] = deep_merge(out[k], v) if k in out else copy.deepcopy(v)
    return out


def resolve_style(layout, deck, slide):
    """スタイル解決順: layout["common"]（全タイプ共通の領域。brand 等）
       ← layout["types"][t]（タイプ既定値）
       ← meta.layout_overrides.common（全タイプ共通のデッキ全体調整）
       ← meta.layout_overrides[t]（タイプ単位のデッキ全体調整）
       ← slide.style（スライド個別の微修正）"""
    t = slide.get("type")
    types = layout.get("types", {})
    if not t:
        fail("スライドに type がありません")
    if t not in types:
        fail(f"未知のスライドタイプ: {t}")
    st = deep_merge(layout.get("common", {}), types[t])
    overrides = deck.get("meta", {}).get("layout_overrides", {})
    st = deep_merge(st, overrides.get("common", {}))
    st = deep_merge(st, overrides.get(t, {}))
    st = deep_merge(st, slide.get("style", {}))
    return st


_HEX_RE = re.compile(r"^#([0-9A-Fa-f]{3}|[0-9A-Fa-f]{6})$")


def is_color_value(theme, c) -> bool:
    """テーマ colors のキー名、または #RGB / #RRGGBB 形式の hex であれば True。"""
    if not isinstance(c, str):
        return False
    return c in theme.get("colors", {}) or bool(_HEX_RE.match(c))


def col(theme, c):
    """テーマトークン名 → hex。hex 直書きはそのまま返す。不正な値は fail する
    （通常は validate_deck が事前に検出するので、これは最後の防波堤）。"""
    colors = theme.get("colors", {})
    if isinstance(c, str) and c in colors:
        return colors[c]
    if isinstance(c, str) and _HEX_RE.match(c):
        return c
    fail(f"不明な色指定 '{c}' です。テーマトークン名か #RRGGBB 形式のhexを指定してください")


# ---------------------------------------------------------------------------
# フォント（fonts.heading / body / code は文字列 または 文字列配列）
# ---------------------------------------------------------------------------

FONT_FALLBACK_SANS = '"Yu Gothic UI", "Meiryo", "Hiragino Sans", "Hiragino Kaku Gothic ProN", "Noto Sans JP", "Noto Sans CJK JP", sans-serif'
FONT_FALLBACK_MONO = '"Consolas", "Menlo", "Courier New", monospace'


def _font_names(value) -> list:
    if isinstance(value, (list, tuple)):
        return [str(v) for v in value if v]
    if value:
        return [str(value)]
    return []


def font_css(value, mono: bool = False) -> str:
    """fonts.* の値（文字列 or 配列）から CSS font-family 用の文字列を作る。
    末尾に組込フォールバック（和文 or 等幅）を必ず付ける。"""
    names = _font_names(value)
    quoted = ", ".join(f'"{n}"' for n in names)
    fallback = FONT_FALLBACK_MONO if mono else FONT_FALLBACK_SANS
    return f"{quoted}, {fallback}" if quoted else fallback


def font_pptx(value):
    """fonts.* の値（文字列 or 配列）から PPTX に設定する先頭の1フォント名を返す。"""
    names = _font_names(value)
    return names[0] if names else None


def norm_items(items):
    """bullets 配列を正規化。要素は文字列 or {text, children, size, color, bold}。"""
    out = []
    for it in items or []:
        if isinstance(it, str):
            out.append({"text": it, "children": []})
        else:
            d = dict(it)
            d.setdefault("text", "")
            kids = []
            for c in d.get("children") or []:
                kids.append({"text": c} if isinstance(c, str) else dict(c))
            d["children"] = kids
            out.append(d)
    return out


def _route_edge(A, B):
    """2ノード間の直交ポリライン経路と、ラベル配置座標を返す。"""
    dx = B["cx"] - A["cx"]; dy = B["cy"] - A["cy"]
    if abs(A["cx"] - B["cx"]) < 2 and abs(dy) >= 1:
        # ほぼ同列：垂直に接続
        if dy > 0:
            p = [(A["cx"], A["y"] + A["h"]), (B["cx"], B["y"])]
        else:
            p = [(A["cx"], A["y"]), (B["cx"], B["y"] + B["h"])]
        return p, (A["cx"] + 8, (A["cy"] + B["cy"]) / 2)
    if abs(dy) < 1:
        # 同一レーン：水平に接続
        if dx >= 0:
            p = [(A["x"] + A["w"], A["cy"]), (B["x"], B["cy"])]
        else:
            p = [(A["x"], A["cy"]), (B["x"] + B["w"], B["cy"])]
        return p, ((p[0][0] + p[1][0]) / 2, A["cy"] - 12)
    # レーンをまたぐ：H-V-H の Z 経路
    if dx >= 0:
        sx = A["x"] + A["w"]; ex = B["x"]
    else:
        sx = A["x"]; ex = B["x"] + B["w"]
    midx = (sx + ex) / 2
    p = [(sx, A["cy"]), (midx, A["cy"]), (midx, B["cy"]), (ex, B["cy"])]
    return p, (midx + 8, (A["cy"] + B["cy"]) / 2)


def swimlane_geometry(slide, st, slide_no=None):
    """スイムレーン図の全要素（レーン・グループ・工程・ノード・エッジ経路）の座標を計算する。
    HTML と PPTX で同一のジオメトリを共有し、見た目を一致させる。
    slide_no（1始まりのスライド番号）を渡すと、自動拡張の warning にどのスライドかを含められる。"""
    f = st["flow"]
    lanes = slide.get("lanes", [])
    n = max(1, len(lanes))
    has_group = any(l.get("group") for l in lanes)
    g1 = f.get("g1_w", 34) if has_group else 0
    g2 = f.get("g2_w", 96)
    label_w = g1 + g2
    phases = slide.get("phases") or []
    phase_h = f.get("phase_h", 30) if any(phases) else 0
    content_x = f["x"] + label_w
    content_w = f["w"] - label_w
    top = f["y"] + phase_h
    avail_h = f["h"] - phase_h
    row_h = avail_h / n
    maxcol = 0
    for nd in slide.get("nodes", []):
        maxcol = max(maxcol, nd.get("col", 0))
    explicit_cols = slide.get("cols")
    cols = max(1, explicit_cols or (maxcol + 1))
    if explicit_cols and maxcol >= cols:
        new_cols = maxcol + 1
        where = f"{slide_no}枚目（type=swimlane）: " if slide_no else ""
        warn(f"{where}cols（{cols}）を超える col={maxcol} のノードがあるため cols を {new_cols} に自動拡張しました")
        cols = new_cols
    col_w = content_w / cols
    lane_boxes = []
    for i, l in enumerate(lanes):
        y = top + i * row_h
        lane_boxes.append({"label": l.get("name", l.get("label", "")),
                           "x": f["x"] + g1, "y": y, "w": g2, "h": row_h,
                           "band_x": content_x, "band_y": y, "band_w": content_w, "band_h": row_h,
                           "alt": i % 2})
    groups = []
    if has_group:
        i = 0
        while i < n:
            gname = lanes[i].get("group")
            j = i
            while j + 1 < n and lanes[j + 1].get("group") == gname and gname is not None:
                j += 1
            if gname:
                groups.append({"label": gname, "x": f["x"], "y": top + i * row_h,
                               "w": g1, "h": (j - i + 1) * row_h})
            i = j + 1
    phase_boxes = []
    for c, ph in enumerate(phases):
        if ph:
            phase_boxes.append({"label": ph, "x": content_x + c * col_w, "y": f["y"], "w": col_w, "h": phase_h})
    node_map = {}
    for nd in slide.get("nodes", []):
        c = nd.get("col", 0); L = nd.get("lane", 0)
        cx = content_x + (c + 0.5) * col_w
        cy = top + (L + 0.5) * row_h
        shape = nd.get("shape", "task")
        if shape == "decision":
            # ひし形は箱に内接する形で描く（HTML は SVG polygon、PPTX は DIAMOND プリセット）。
            # 旧実装（正方形を 45° 回転）と同じ見た目の大きさになるよう外接矩形を 0.93 倍にとる。
            # 当たり判定（配線・ラベル衝突）もこの外接矩形で行うので、菱形の角がすり抜けない。
            w = min(col_w, row_h) * 0.93; h = w
        elif shape in ("terminal", "connector"):
            w = col_w * f.get("node_wr", 0.72); h = row_h * 0.42
        elif shape == "marker":
            w = min(col_w, row_h) * 0.34; h = w
        elif shape == "mail":
            w = min(col_w, row_h) * 0.44; h = w * 0.72
        elif shape == "io":
            w = col_w * 0.88; h = row_h * 0.92
        else:  # task, system
            w = col_w * f.get("node_wr", 0.72); h = row_h * f.get("node_hr", 0.52)
        node_map[nd["id"]] = {"shape": shape, "variant": nd.get("variant"), "kind": nd.get("kind"),
                              "text": nd.get("text", nd.get("label", "")),  # label は text の別名として受理
                              "cx": cx, "cy": cy, "x": cx - w / 2, "y": cy - h / 2, "w": w, "h": h,
                              "input": nd.get("input"), "output": nd.get("output"), "loop": nd.get("loop")}
    # エッジ配線: 共通の格子ルーター（diagram_engine）。列境界・レーン境界のガターをトラックにし、
    # 途中のノードを横切らない直交経路を選ぶ。同じ辺から出る複数線・同じ回廊を通る複数線は
    # 等間隔にずらし、ラベルは線の脇に背景ピル付きで置く（Archify の workflow レンダラに倣う）。
    my = row_h * (1 - f.get("node_hr", 0.64)) / 2   # レーン内の上下マージン（片側）
    mx = col_w * (1 - f.get("node_wr", 0.74)) / 2
    h_tracks = [top + my * 0.5, top + avail_h - my * 0.5]
    for i in range(1, n):
        yb = top + i * row_h
        # レーン境界線そのものの上を線が走らないよう、境界の上下マージンの中央を通す
        h_tracks += [yb - my * 0.5, yb + my * 0.5]
    v_tracks = [content_x + mx * 0.5, content_x + content_w - mx * 0.5] + \
               [content_x + c * col_w for c in range(1, cols)]
    grid = de.Grid(content_x, top, content_w, avail_h, n, cols,
                   node_wr=f.get("node_wr", 0.74), node_hr=f.get("node_hr", 0.64),
                   v_tracks=v_tracks, h_tracks=h_tracks)
    boxes = {nid: dict(nd, label=nd.get("text", "")) for nid, nd in node_map.items()}
    diags = []
    routed = de.route_all(grid, boxes, slide.get("edges", []), f.get("label_size", 13), diags,
                          subject_prefix="")
    edges = []
    for e in routed:
        lb = e.get("label_box")
        edges.append({"points": e["points"], "label": e.get("label", ""),
                      "label_pos": ((lb["cx"], lb["cy"]) if lb else None), "label_box": lb,
                      "style": e.get("style", "solid"), "variant": e.get("variant", "default"),
                      "arrow": e.get("arrow", "end"), "from": e.get("from"), "to": e.get("to")})
    return {"flow": f, "content_x": content_x, "content_w": content_w, "top": top,
            "row_h": row_h, "col_w": col_w, "cols": cols, "lanes": lane_boxes,
            "groups": groups, "phases": phase_boxes, "nodes": node_map, "edges": edges,
            "diagnostics": diags}


# スイムレーンの凡例エントリ（業務フロー標準記号）: (shape, variant/kind, ラベル, 説明)
SWIMLANE_LEGEND = [
    ("task", "onpf", "作業（オンラインPF）", "PF 上で実施する作業"),
    ("task", "onother", "作業（PF以外）", "他システムで実施する作業"),
    ("task", "offline", "作業（オフライン）", "オフラインで実施する作業"),
    ("system", None, "システム/サービス", "利用するシステム・サービス名"),
    ("decision", None, "分岐条件", "フロー内で分岐する条件"),
    ("io", None, "成果物（input/output）", "作業の入力・出力ドキュメント"),
    ("marker", "start", "フロー開始", "フローの起点"),
    ("marker", "end", "フロー終了", "フローの終点"),
    ("marker", "mid", "フロー途中", "フロー途中の状態"),
    ("terminal", "prev", "前ページより", "前ページからの続き"),
    ("terminal", "next", "次ページへ", "後続ページへ続く"),
    ("connector", None, "遷移先（分岐）", "複数の後続業務への遷移先"),
    ("mail", None, "メール配信", "メール通知の発生"),
    ("flow", "solid", "作業の流れ", "実線＝作業の流れ"),
    ("flow", "dashed", "システム操作", "破線＝システム操作"),
]


def _sw_node_html(nd, theme, f):
    """スイムレーンのノード1個を HTML で描く（フロー本体・凡例ページ共通）。"""
    def C(tok):
        return col(theme, tok)
    x, y, w, h = nd["x"], nd["y"], nd["w"], nd["h"]
    box = f'position:absolute;left:{x:.1f}px;top:{y:.1f}px;width:{w:.1f}px;height:{h:.1f}px;'
    shape = nd["shape"]; text = esc(nd.get("text", "")); ns = f.get("node_size", 14)
    center = ('display:flex;align-items:center;justify-content:center;text-align:center;'
              'box-sizing:border-box;line-height:1.15;padding:2px;')
    if shape == "decision":
        # 箱に内接するひし形（PPTX の DIAMOND プリセットと同じ輪郭。当たり判定も箱＝外接矩形で一致する）
        return (f'<div style="{box}">'
                f'<svg viewBox="0 0 {w:.1f} {h:.1f}" width="{w:.1f}" height="{h:.1f}" style="position:absolute;inset:0;overflow:visible;">'
                f'<polygon points="{w/2:.1f},1 {w-1:.1f},{h/2:.1f} {w/2:.1f},{h-1:.1f} 1,{h/2:.1f}" '
                f'fill="{C(f.get("decision_fill","on_primary_soft"))}" stroke="{C(f.get("decision_border","accent"))}" '
                f'stroke-width="2" stroke-linejoin="round"/></svg>'
                f'<div style="position:absolute;inset:{h*0.2:.1f}px {w*0.18:.1f}px;{center}font-size:{ns-2}px;'
                f'color:{C(f.get("decision_color","primary"))};">{text}</div></div>')
    if shape == "terminal":
        return (f'<div style="{box}{center}font-size:{ns-2}px;background:{C(f.get("terminal_fill","surface"))};'
                f'border:1.5px solid {C(f.get("terminal_border","muted"))};border-radius:{h/2:.0f}px;'
                f'color:{C(f.get("terminal_color","muted"))};">{text}</div>')
    if shape == "connector":
        # M-90: PPTX（PENTAGON + connector_border 罫線）と見た目を揃える。
        # div+clip-path は塗り・背景しか切り抜かず斜辺に罫線を描けない（周囲と同系色の背景だと
        # 輪郭が消えてただの四角に見える）ため、SVG の polygon（fill+stroke）で5辺すべてに
        # 一貫した罫線を描く。テキストは別レイヤーの div で中央寄せする。
        pts = f"0,0 {w * 0.72:.1f},0 {w:.1f},{h / 2:.1f} {w * 0.72:.1f},{h:.1f} 0,{h:.1f}"
        return (f'<div style="{box}">'
                f'<svg viewBox="0 0 {w:.1f} {h:.1f}" style="position:absolute;inset:0;'
                f'width:100%;height:100%;overflow:visible;" preserveAspectRatio="none">'
                f'<polygon points="{pts}" fill="{C(f.get("connector_fill","surface"))}" '
                f'stroke="{C(f.get("connector_border","muted"))}" stroke-width="1"/></svg>'
                f'<div style="position:absolute;inset:0;{center}font-size:{ns-3}px;'
                f'color:{C(f.get("connector_color","text"))};">{text}</div></div>')
    if shape == "marker":
        kind = nd.get("kind", "mid")
        mc = {"start": f.get("marker_start", "accent"), "end": f.get("marker_end", "primary"),
              "mid": f.get("marker_mid", "#E6B800")}.get(kind, f.get("marker_mid", "#E6B800"))
        return f'<div style="{box}background:{C(mc)};border-radius:50%;border:1.5px solid {C("background")};"></div>'
    if shape == "mail":
        return (f'<svg viewBox="0 0 {w:.0f} {h:.0f}" style="{box}">'
                f'<rect x="1" y="1" width="{w-2:.0f}" height="{h-2:.0f}" fill="{C("background")}" '
                f'stroke="{C(f.get("edge_color","text"))}" stroke-width="1.5"/>'
                f'<path d="M1,1 L{w/2:.0f},{h*0.55:.0f} L{w-1:.0f},1" fill="none" '
                f'stroke="{C(f.get("edge_color","text"))}" stroke-width="1.5"/></svg>')
    if shape == "io":
        half = (w - 2) / 2
        cells = []
        for key, lbl in (("input", "input"), ("output", "output")):
            vals = nd.get(key) or []
            lis = "".join(f'<div style="margin-top:3px;">{esc(v)}</div>' for v in vals)
            cells.append(f'<div style="width:{half:.0f}px;padding:6px;box-sizing:border-box;'
                         f'font-size:{f.get("io_size",11)}px;color:{C(f.get("io_color","text"))};line-height:1.2;">'
                         f'<div style="font-weight:700;color:{C(f.get("io_head_color","accent"))};">{lbl}</div>{lis}</div>')
        return (f'<div style="{box}display:flex;background:{C(f.get("io_fill","on_primary_soft"))};'
                f'border:1px solid {C(f.get("io_border","accent"))};border-radius:8px;box-sizing:border-box;overflow:hidden;">'
                f'{cells[0]}<div style="width:1px;background:{C("border")};"></div>{cells[1]}</div>')
    if shape == "system":
        return (f'<div style="{box}{center}font-size:{ns}px;background:{C(f.get("system_fill","surface"))};'
                f'border:1.5px solid {C(f.get("system_border","muted"))};color:{C(f.get("system_color","text"))};">{text}</div>')
    # task (+variant)
    variant = nd.get("variant") or "onother"
    vmap = {
        "onpf": (f.get("task_onpf_fill", "on_primary_soft"), f.get("task_border", "accent"), f.get("task_onpf_color", "primary")),
        "onother": (f.get("task_onother_fill", "background"), f.get("task_border", "accent"), f.get("task_color", "text")),
        "offline": (f.get("task_offline_fill", "surface"), f.get("task_offline_border", "muted"), f.get("task_color", "text")),
    }
    fill_t, bord_t, col_t = vmap.get(variant, vmap["onother"])
    loop = ""
    if nd.get("loop"):
        # 反復記号: ノード右上に小さな円弧矢印（SVG）。旧実装の「↺」1 文字は 16px でも見落とされやすかった。
        loop = (f'<svg width="18" height="18" viewBox="0 0 18 18" style="position:absolute;right:-7px;top:-8px;" aria-hidden="true">'
                f'<circle cx="9" cy="9" r="8" fill="{C("background")}" stroke="{C("accent")}" stroke-width="1.5"/>'
                f'<path d="M13 9a4 4 0 1 1-1.2-2.9 M12 4.5v2.6h-2.6" fill="none" stroke="{C("accent")}" stroke-width="1.6" '
                f'stroke-linecap="round" stroke-linejoin="round"/></svg>')
    return (f'<div style="{box}{center}font-size:{ns}px;background:{C(fill_t)};'
            f'border:2px solid {C(bord_t)};border-radius:{f.get("node_radius",8)}px;color:{C(col_t)};">{text}{loop}</div>')


def _legend_sym_size(shape):
    return {"decision": (26, 26), "marker": (16, 16), "mail": (30, 20), "terminal": (74, 24),
            "connector": (66, 26), "io": (78, 30), "flow": (54, 0)}.get(shape, (60, 26))


def _legend_symbol_html(shape, vk, x, cy, sym_w, theme, f):
    """凡例の1シンボルを HTML で描く（中心 x..x+sym_w に配置）。"""
    w, h = _legend_sym_size(shape)
    cx = x + sym_w / 2
    if shape == "flow":
        dash = ' stroke-dasharray="7,5"' if vk == "dashed" else ""
        ec = col(theme, f.get("edge_color", "text"))
        return (f'<svg style="position:absolute;left:0;top:0;width:{CANVAS_W}px;height:{CANVAS_H}px;'
                f'pointer-events:none;" viewBox="0 0 {CANVAS_W} {CANVAS_H}">'
                f'<defs><marker id="lg-arrow" markerWidth="10" markerHeight="10" refX="8" refY="4" '
                f'orient="auto" markerUnits="userSpaceOnUse"><path d="M0,0 L9,4 L0,8 z" fill="{ec}"/></marker></defs>'
                f'<line x1="{cx-w/2:.0f}" y1="{cy:.0f}" x2="{cx+w/2:.0f}" y2="{cy:.0f}" stroke="{ec}" '
                f'stroke-width="2"{dash} marker-end="url(#lg-arrow)"/></svg>')
    nd = {"shape": shape, "variant": (vk if shape == "task" else None),
          "kind": (vk if shape == "marker" else None), "text": "",
          "x": cx - w / 2, "y": cy - h / 2, "w": w, "h": h,
          "input": (["…"] if shape == "io" else None), "output": (["…"] if shape == "io" else None),
          "loop": None}
    return _sw_node_html(nd, theme, f)


def norm_legend_items(items):
    """swimlane の legend_items を (shape, variant_or_kind, label, desc) の4要素タプル配列に正規化する。
    4要素配列（リスト/タプル）はそのまま、辞書は {shape, variant|kind, label, desc} から変換する。"""
    out = []
    for it in items or []:
        if isinstance(it, dict):
            shape = it.get("shape") or it.get("kind")
            variant = it.get("variant", it.get("kind"))
            label = it.get("label", "")
            desc = it.get("desc", "")
            out.append((shape, variant, label, desc))
        else:
            out.append(tuple(it))
    return out


# ---------------------------------------------------------------------------
# ネイティブ図解タイプ（architecture / dataflow / lifecycle / sequence）の見た目解決と HTML 描画
# ジオメトリと診断は diagram_engine が担当し、ここではテーマトークンへの解決と描画だけを行う。
# PPTX 側（build_pptx 内の _diag_* 関数）も同じ _node_visual / _edge_visual を使う。
# ---------------------------------------------------------------------------

GRID_DIAGRAM_TYPES = ("architecture", "dataflow", "lifecycle")
DIAGRAM_TYPES = GRID_DIAGRAM_TYPES + ("sequence",)

# type アイコン（16×16、stroke ベース）。HTML はインライン SVG、PPTX はプリセット図形で近似する。
ICON_PATHS = {
    "frontend": "M1.5 2.5h13v11h-13z M1.5 5.5h13 M3.6 4h.01 M5.6 4h.01",
    "backend": "M2 2.5h12v4.5H2z M2 9h12v4.5H2z M4.6 4.75h.01 M4.6 11.25h.01",
    "database": ("M2.5 4.5c0-1.4 2.5-2.5 5.5-2.5s5.5 1.1 5.5 2.5v7c0 1.4-2.5 2.5-5.5 2.5S2.5 12.9 2.5 11.5z "
                 "M2.5 4.5c0 1.4 2.5 2.5 5.5 2.5s5.5-1.1 5.5-2.5 M2.5 8c0 1.4 2.5 2.5 5.5 2.5S13.5 9.4 13.5 8"),
    "cloud": "M5 13.5h6.5a3 3 0 0 0 .5-5.96A4 4 0 0 0 4.3 7.2 3.2 3.2 0 0 0 5 13.5z",
    "security": "M8 1.8l5.5 2.2v3.9c0 3.3-2.3 5.6-5.5 6.6-3.2-1-5.5-3.3-5.5-6.6V4z M5.6 8l1.7 1.7 3.2-3.2",
    "messagebus": "M2 4.5h12 M2 8h9 M11.5 6.5L13 8l-1.5 1.5 M2 11.5h12",
    "external": "M12 9v4.5H2.5V4H7 M9.5 2H14v4.5 M14 2L7.5 8.5",
}
# lifecycle の kind 記号（右上の小さなグリフ）
KIND_GLYPHS = {
    "start": "M5 3l7 5-7 5z",
    "waiting": "M4 2h8 M4 14h8 M5 2c0 4 3 5 3 6s-3 2-3 6 M11 2c0 4-3 5-3 6s3 2 3 6",
    "success": "M3 8.5l3.2 3.2L13 5",
    "failure": "M4 4l8 8 M12 4l-8 8",
}
EDGE_DASH_SVG = {"dashed": "7,5", "security": "8,4,2,4", "return": "6,4"}
EDGE_DASH_PPTX = {"dashed": "dash", "security": "dashDot", "return": "sysDash"}


def _node_visual(box: dict, d: dict, kind: str) -> dict:
    """ノード 1 個の描画属性（テーマトークン名）を variant / lifecycle kind から解決する。
    返り値: fill, border, border_w, dash(bool), color, sub_color, radius_mode(rect|pill|diamond)"""
    v = box.get("variant") or "default"
    vis = {"fill": d.get("node_fill", "background"), "border": d.get("node_border", "border"),
           "border_w": d.get("node_border_w", 1.5), "dash": False,
           "color": d.get("node_color", "heading_text"), "sub_color": d.get("sub_color", "muted"),
           "shape": box.get("shape", "rect")}
    if kind == "lifecycle":
        k = box.get("kind") or "active"
        vis["fill"] = d.get(f"{k}_fill", vis["fill"])
        vis["border"] = d.get(f"{k}_border", vis["border"])
        if k in ("waiting", "external"):
            vis["dash"] = True
        elif k == "failure":
            vis["border_w"] = 2
            vis["dash"] = True
        elif k == "success":
            vis["color"] = d.get("success_color", "on_primary")
            vis["sub_color"] = d.get("success_sub_color", "on_primary_soft")
            vis["border_w"] = 2
        elif k in ("start", "decision"):
            vis["border_w"] = 2
    if v == "emphasis":
        vis["fill"] = d.get("emphasis_fill", "highlight_fill")
        vis["border"] = d.get("emphasis_border", "accent")
        vis["border_w"] = d.get("emphasis_border_w", 2.5)
    elif v == "security":
        vis["border"] = d.get("security_border", "primary")
        vis["border_w"] = 2
        vis["dash"] = True
    elif v == "dashed":
        vis["border"] = d.get("dashed_border", "muted")
        vis["dash"] = True
    elif v == "muted":
        vis["color"] = d.get("muted_color", "muted")
        vis["sub_color"] = d.get("muted_color", "muted")
    return vis


def _edge_visual(e: dict, d: dict) -> dict:
    """エッジの描画属性: color(token), w(px), dash(svg pattern or None), pptx_dash, open_arrow(bool)。"""
    v = e.get("variant") or "default"
    if v == "emphasis":
        return {"color": d.get("edge_emphasis_color", "accent"), "w": d.get("edge_emphasis_w", 3),
                "dash": None, "pptx_dash": None, "open": False}
    if v == "security":
        return {"color": d.get("edge_security_color", "primary"), "w": d.get("edge_w", 2),
                "dash": EDGE_DASH_SVG["security"], "pptx_dash": EDGE_DASH_PPTX["security"], "open": False}
    if v == "dashed":
        return {"color": d.get("edge_dashed_color", "muted"), "w": d.get("edge_w", 2),
                "dash": EDGE_DASH_SVG["dashed"], "pptx_dash": EDGE_DASH_PPTX["dashed"], "open": False}
    if v == "return":
        return {"color": d.get("edge_return_color", "muted"), "w": max(1.5, d.get("edge_w", 2) - 0.5),
                "dash": EDGE_DASH_SVG["return"], "pptx_dash": EDGE_DASH_PPTX["return"], "open": True}
    return {"color": d.get("edge_color", "text"), "w": d.get("edge_w", 2), "dash": None, "pptx_dash": None, "open": False}


def _icon_svg(kind: str, size: float, color_hex: str, extra_style: str = "") -> str:
    path = ICON_PATHS.get(kind)
    if not path:
        return ""
    return (f'<svg width="{size:.0f}" height="{size:.0f}" viewBox="0 0 16 16" style="flex:none;{extra_style}" '
            f'aria-hidden="true"><path d="{path}" fill="none" stroke="{color_hex}" stroke-width="1.5" '
            f'stroke-linecap="round" stroke-linejoin="round"/></svg>')


def _glyph_svg(kind: str, size: float, color_hex: str, x: float, y: float) -> str:
    path = KIND_GLYPHS.get(kind)
    if not path:
        return ""
    fill = color_hex if kind == "start" else "none"
    return (f'<svg width="{size:.0f}" height="{size:.0f}" viewBox="0 0 16 16" style="position:absolute;'
            f'left:{x:.1f}px;top:{y:.1f}px;" aria-hidden="true"><path d="{path}" fill="{fill}" stroke="{color_hex}" '
            f'stroke-width="1.6" stroke-linecap="round" stroke-linejoin="round"/></svg>')


def _diag_node_html(b: dict, fit: dict, d: dict, theme, kind: str, heading_font_css: str) -> str:
    """ノードカード 1 個（箱・アイコン・label/sublabel・tag・step）を HTML で描く。"""
    vis = _node_visual(b, d, kind)
    x, y, w, h = b["x"], b["y"], b["w"], b["h"]
    pad = d.get("node_pad", 8)
    radius = d.get("node_radius", 10)
    fill = col(theme, vis["fill"]); border = col(theme, vis["border"])
    bstyle = "dashed" if vis["dash"] else "solid"
    parts = []
    box_css = f'position:absolute;left:{x:.1f}px;top:{y:.1f}px;width:{w:.1f}px;height:{h:.1f}px;box-sizing:border-box;'
    if vis["shape"] == "diamond":
        # PPTX の DIAMOND プリセットと同じ輪郭（箱の 4 辺の中点を結ぶ）を SVG で描く
        dash_attr = ' stroke-dasharray="6,4"' if vis["dash"] else ""
        parts.append(f'<div style="{box_css}"><svg viewBox="0 0 {w:.1f} {h:.1f}" width="{w:.1f}" height="{h:.1f}" '
                     f'style="position:absolute;inset:0;overflow:visible;"><polygon points="{w/2:.1f},1 {w-1:.1f},{h/2:.1f} '
                     f'{w/2:.1f},{h-1:.1f} 1,{h/2:.1f}" fill="{fill}" stroke="{border}" stroke-width="{vis["border_w"]}"'
                     f'{dash_attr} stroke-linejoin="round"/></svg></div>')
        icon_w = 0.0
    else:
        rad = f"{h/2:.0f}px" if vis["shape"] == "pill" else f"{radius}px"
        parts.append(f'<div style="{box_css}background:{fill};border:{vis["border_w"]}px {bstyle} {border};border-radius:{rad};"></div>')
        icon_w = 0.0
    # type アイコン（左・縦中央）
    icon_size = d.get("icon_size", 16)
    ntype = b.get("type") or "generic"
    if kind != "lifecycle" and ntype in ICON_PATHS and vis["shape"] == "rect":
        # muted variant はアイコンも muted に落とす（type は形で表す・色は variant が担う、の唯一の例外）
        icon_tok = d.get("muted_color", "muted") if (b.get("variant") == "muted") else d.get("icon_color", "accent")
        parts.append(f'<div style="position:absolute;left:{x+pad:.1f}px;top:{y + h/2 - icon_size/2:.1f}px;'
                     f'width:{icon_size}px;height:{icon_size}px;">'
                     + _icon_svg(ntype, icon_size, col(theme, icon_tok)) + "</div>")
        icon_w = icon_size + 6
    # テキストブロック（アイコンの右側の残り幅で中央寄せ）
    tx = x + pad + icon_w
    tw = max(10.0, w - 2 * pad - icon_w)
    if vis["shape"] == "diamond":
        tx, tw = x + w * 0.12, w * 0.76
    size = fit.get("size", d.get("node_size", 14)) if fit else d.get("node_size", 14)
    lines = fit.get("lines") if fit else [b.get("label", "")]
    label_html = "<br>".join(esc(ln) for ln in (lines or [b.get("label", "")]))
    sub_html = ""
    if b.get("sublabel"):
        sub_lines = fit.get("sub_lines") if fit else [b["sublabel"]]
        sub_html = (f'<div style="font-size:{d.get("sub_size", 11)}px;color:{col(theme, vis["sub_color"])};'
                    f'line-height:1.2;margin-top:2px;">{"<br>".join(esc(l) for l in sub_lines)}</div>')
    parts.append(f'<div style="position:absolute;left:{tx:.1f}px;top:{y:.1f}px;width:{tw:.1f}px;height:{h:.1f}px;'
                 f'display:flex;flex-direction:column;align-items:center;justify-content:center;text-align:center;'
                 f'box-sizing:border-box;padding:{pad}px 0;">'
                 f'<div style="font-size:{size}px;font-weight:700;color:{col(theme, vis["color"])};line-height:1.15;'
                 f'font-family:{esc(heading_font_css)};">{label_html}</div>{sub_html}</div>')
    # tag（右上に小さなピル。枠線に半分かぶせる）
    if b.get("tag"):
        ts = d.get("tag_size", 10)
        tw_ = de.text_width(b["tag"], ts) + 12
        parts.append(f'<div style="position:absolute;left:{x + w - tw_ - 8:.1f}px;top:{y - ts*0.8:.1f}px;'
                     f'width:{tw_:.1f}px;height:{ts*1.6:.1f}px;border-radius:{ts*0.8:.0f}px;'
                     f'background:{col(theme, d.get("tag_fill", "highlight_fill"))};color:{col(theme, d.get("tag_color", "text"))};'
                     f'font-size:{ts}px;font-weight:700;display:flex;align-items:center;justify-content:center;'
                     f'line-height:1;">{esc(b["tag"])}</div>')
    if kind == "lifecycle":
        if b.get("step"):
            ss = d.get("step_size", 11)
            sx = (x + w / 2 + 8) if vis["shape"] == "diamond" else x + 8
            parts.append(f'<div style="position:absolute;left:{sx:.1f}px;top:{y + 4:.1f}px;font-size:{ss}px;'
                         f'font-weight:700;color:{col(theme, d.get("step_color", "accent"))};line-height:1;">{esc(b["step"])}</div>')
        k = b.get("kind") or "active"
        if k in KIND_GLYPHS and vis["shape"] == "rect":
            gc = col(theme, vis["color"]) if k == "success" else col(theme, d.get("icon_color", "accent"))
            parts.append(_glyph_svg(k, 13, gc, x + w - 21, y + 5))
    return "".join(parts)


def _arrow_defs(colors, uid: str) -> str:
    """色ごとの矢印マーカー（塗り三角）と開き矢印（return 用）を定義する。"""
    defs = []
    for hexc, is_open in colors:
        key = f"{uid}-{hexc.lstrip('#')}-{'o' if is_open else 'c'}"
        if is_open:
            defs.append(f'<marker id="ae-{key}" markerWidth="12" markerHeight="12" refX="9" refY="5" orient="auto-start-reverse" '
                        f'markerUnits="userSpaceOnUse"><path d="M1,1 L9,5 L1,9" fill="none" stroke="{hexc}" stroke-width="1.6"/></marker>')
        else:
            defs.append(f'<marker id="ae-{key}" markerWidth="10" markerHeight="10" refX="8" refY="4" orient="auto-start-reverse" '
                        f'markerUnits="userSpaceOnUse"><path d="M0,0 L9,4 L0,8 z" fill="{hexc}"/></marker>')
    return "<defs>" + "".join(defs) + "</defs>"


def _edges_svg_html(edges, d: dict, theme, uid: str, heading_font_css: str, layer: str = "both") -> str:
    """エッジ（折れ線＋矢印）とラベルピルを SVG オーバーレイに描く。
    layer="lines" で線だけ、"labels" でラベルピルだけ（ノードの上に重ねるため 2 層に分けて呼ぶ）。"""
    colors = set()
    lines = []
    pills = []
    lc = col(theme, d.get("label_color", "accent"))
    lf = col(theme, d.get("label_fill", "background"))
    lsize = d.get("label_size", 13)
    for e in edges:
        vis = _edge_visual(e, d)
        hexc = col(theme, vis["color"])
        colors.add((hexc, vis["open"]))
        key = f"{uid}-{hexc.lstrip('#')}-{'o' if vis['open'] else 'c'}"
        pts = " ".join(f"{x:.1f},{y:.1f}" for x, y in e["points"])
        dash = f' stroke-dasharray="{vis["dash"]}"' if vis["dash"] else ""
        arrow = e.get("arrow", "end")
        mk = ""
        if arrow in ("end", "both"):
            mk += f' marker-end="url(#ae-{key})"'
        if arrow == "both":
            mk += f' marker-start="url(#ae-{key})"'
        lines.append(f'<polyline points="{pts}" fill="none" stroke="{hexc}" stroke-width="{vis["w"]}"{dash}{mk} '
                     f'stroke-linejoin="round"/>')
        lb = e.get("label_box")
        if lb and e.get("label"):
            pills.append(f'<rect x="{lb["x"]:.1f}" y="{lb["y"]:.1f}" width="{lb["w"]:.1f}" height="{lb["h"]:.1f}" rx="4" fill="{lf}"/>')
            classification = e.get("classification")
            if classification:
                # ピルは 2 行分の高さ（engine が class_h を足している）。上段 label・下段 classification
                cs = d.get("class_size", 11)
                ch = e.get("class_h", cs * 1.2)
                pills.append(f'<text x="{lb["cx"]:.1f}" y="{lb["cy"] - ch / 2:.1f}" fill="{lc}" font-size="{lsize}" font-weight="700" '
                             f'text-anchor="middle" dominant-baseline="central">{esc(e["label"])}</text>')
                pills.append(f'<text x="{lb["cx"]:.1f}" y="{lb["cy"] + (lb["h"] - ch) / 2 - 1:.1f}" fill="{col(theme, d.get("class_color", "muted"))}" '
                             f'font-size="{cs}" text-anchor="middle" dominant-baseline="central">{esc(classification)}</text>')
            else:
                pills.append(f'<text x="{lb["cx"]:.1f}" y="{lb["cy"]:.1f}" fill="{lc}" font-size="{lsize}" font-weight="700" '
                             f'text-anchor="middle" dominant-baseline="central">{esc(e["label"])}</text>')
    body = ""
    if layer in ("both", "lines"):
        body += _arrow_defs(sorted(colors), uid) + "".join(lines)
    if layer in ("both", "labels"):
        body += "".join(pills)
    return (f'<svg style="position:absolute;left:0;top:0;width:{CANVAS_W}px;height:{CANVAS_H}px;pointer-events:none;" '
            f'viewBox="0 0 {CANVAS_W} {CANVAS_H}">' + body + "</svg>")


def _legend_html(items, y: float, x: float, d: dict, theme) -> str:
    if not items:
        return ""
    ls = d.get("legend_size", 12)
    lc = col(theme, d.get("legend_color", "muted"))
    ic = col(theme, d.get("icon_color", "accent"))
    cells = []
    for it in items:
        cells.append(f'<div style="display:flex;align-items:center;gap:5px;">{_icon_svg(it["type"], 14, ic)}'
                     f'<span>{esc(it["label"])}</span></div>')
    return (f'<div style="position:absolute;left:{x:.1f}px;top:{y:.1f}px;height:{d.get("legend_h", 26)}px;'
            f'display:flex;align-items:center;gap:22px;font-size:{ls}px;color:{lc};">' + "".join(cells) + "</div>")


def _group_html(g: dict, d: dict, theme) -> str:
    kind = g.get("kind") or "generic"
    radius = d.get("group_radius", 10)
    gc = col(theme, d.get("group_color", "muted"))
    if kind == "zone":
        box = (f'background:{col(theme, d.get("zone_fill", "surface"))};border:1px solid {col(theme, d.get("group_border", "muted"))};')
    elif kind == "security":
        box = f'border:1.5px dashed {col(theme, d.get("security_group_border", "primary"))};'
        gc = col(theme, d.get("security_group_border", "primary"))
    else:
        box = f'border:1.5px dashed {col(theme, d.get("group_border", "muted"))};'
    out = (f'<div style="position:absolute;left:{g["x"]:.1f}px;top:{g["y"]:.1f}px;width:{g["w"]:.1f}px;height:{g["h"]:.1f}px;'
           f'box-sizing:border-box;border-radius:{radius}px;{box}"></div>')
    if g.get("label"):
        gs = d.get("group_size", 12)
        out += (f'<div style="position:absolute;left:{g["x"] + 10:.1f}px;top:{g["y"] + 3:.1f}px;font-size:{gs}px;'
                f'font-weight:700;color:{gc};line-height:1.2;">{esc(g["label"])}</div>')
    return out


def grid_diagram_html(slide, st, theme, kind: str, page: int, heading_font_css: str) -> str:
    """architecture / dataflow / lifecycle の図領域を HTML で描く（chrome 以外）。"""
    d = st["diagram"]
    geo = de.layout_grid_diagram(slide, kind, d, [], subject_prefix="")
    parts = []
    # 行帯（lifecycle の lanes）。区切り線は帯 div の border で描く（swimlane と同じ。独立した 1px div に
    # すると check_layout がラベル文字との交差を「罫線」誤検出するため）。
    for i, rh in enumerate(geo["row_headers"]):
        bg = col(theme, d.get("row_band_b", "surface")) if rh["alt"] else col(theme, "background")
        bottom = f'border-bottom:1px solid {col(theme, d.get("row_border", "border"))};' if i == len(geo["row_headers"]) - 1 else ""
        parts.append(f'<div style="{pos({"x": rh["band_x"], "y": rh["y"], "w": rh["band_w"], "h": rh["h"]})}'
                     f'background:{bg};border-top:1px solid {col(theme, d.get("row_border", "border"))};{bottom}'
                     f'box-sizing:border-box;"></div>')
        parts.append(f'<div style="{pos({"x": rh["x"], "y": rh["y"], "w": rh["w"], "h": rh["h"]})}'
                     f'display:flex;align-items:center;justify-content:center;text-align:center;'
                     f'background:{col(theme, d.get("row_fill", "surface"))};border:1px solid {col(theme, d.get("row_border", "border"))};'
                     f'box-sizing:border-box;color:{col(theme, d.get("row_color", "text"))};font-size:{d.get("row_size", 13)}px;'
                     f'padding:2px;line-height:1.2;">{esc(rh["label"])}</div>')
    # 列見出し帯（dataflow の stages / col_headers）
    for ch in geo["col_headers"]:
        parts.append(f'<div style="{pos(ch)}display:flex;align-items:center;justify-content:center;'
                     f'background:{col(theme, d.get("header_fill", "highlight_fill"))};box-sizing:border-box;'
                     f'color:{col(theme, d.get("header_color", "text"))};font-weight:700;font-size:{d.get("header_size", 14)}px;'
                     f'border-left:2px solid {col(theme, "background")};">{esc(ch["label"])}</div>')
    for g in geo["groups"]:
        parts.append(_group_html(g, d, theme))
    parts.append(_edges_svg_html(geo["edges"], d, theme, f"s{page}", heading_font_css, layer="lines"))
    for nid, b in geo["nodes"].items():
        parts.append(_diag_node_html(b, geo["text"].get(nid, {}), d, theme, kind, heading_font_css))
    parts.append(_edges_svg_html(geo["edges"], d, theme, f"s{page}", heading_font_css, layer="labels"))
    if geo["legend"]:
        parts.append(_legend_html(geo["legend"], geo["legend_y"], d["x"] + 4, d, theme))
    return "".join(parts)


def sequence_html(slide, st, theme, page: int, heading_font_css: str) -> str:
    d = st["diagram"]
    geo = de.layout_sequence(slide, d, [], subject_prefix="")
    parts = []
    # segments（区間の薄い枠）
    for sg in geo["segments"]:
        parts.append(f'<div style="{pos(sg)}box-sizing:border-box;border:1px dashed {col(theme, d.get("segment_border", "muted"))};'
                     f'border-radius:6px;"></div>')
        if sg.get("label"):
            ss = d.get("segment_size", 12)
            parts.append(f'<div style="position:absolute;left:{sg["x"] + 8:.1f}px;top:{sg["y"] - ss * 0.75:.1f}px;padding:0 6px;'
                         f'background:{col(theme, "background")};font-size:{ss}px;font-weight:700;'
                         f'color:{col(theme, d.get("segment_color", "muted"))};line-height:1.3;">{esc(sg["label"])}</div>')
    # lifelines
    svg = [f'<svg style="position:absolute;left:0;top:0;width:{CANVAS_W}px;height:{CANVAS_H}px;pointer-events:none;" '
           f'viewBox="0 0 {CANVAS_W} {CANVAS_H}">']
    llc = col(theme, d.get("lifeline_color", "muted"))
    for pid in geo["order"]:
        pb = geo["participants"].get(pid)
        if not pb:
            continue
        svg.append(f'<line x1="{pb["cx"]:.1f}" y1="{geo["lifeline_top"]:.1f}" x2="{pb["cx"]:.1f}" y2="{geo["lifeline_bottom"]:.1f}" '
                   f'stroke="{llc}" stroke-width="{d.get("lifeline_w", 1.5)}" stroke-dasharray="6,5"/>')
    svg.append("</svg>")
    parts.append("".join(svg))
    # activations
    for a in geo["activations"]:
        parts.append(f'<div style="{pos(a)}background:{col(theme, d.get("activation_fill", "highlight_fill"))};'
                     f'border:1.5px solid {col(theme, d.get("activation_border", "accent"))};border-radius:3px;box-sizing:border-box;"></div>')
    # messages
    parts.append(_edges_svg_html(geo["messages"], d, theme, f"s{page}", heading_font_css))
    # participants
    fit = de.check_node_text(geo["participants"], d.get("node_size", 14), d.get("sub_size", 11), d.get("node_pad", 8), [],
                             icon_w=d.get("icon_size", 16) + 6)
    for pid, pb in geo["participants"].items():
        parts.append(_diag_node_html(dict(pb, shape="rect"), fit.get(pid, {}), d, theme, "sequence", heading_font_css))
    return "".join(parts)


def expand_slides(deck, layout):
    """ビルド前処理: agenda のあふれを複数ページに分割し、swimlane の直前に凡例ページを挿入する。"""
    out = []
    for s in deck.get("slides", []):
        t = s.get("type")
        if t == "swimlane" and s.get("legend", True):
            leg = {"type": "swimlane_legend",
                   "title": s.get("legend_title", "凡例（スイムレーン記号）"),
                   "eyebrow": s.get("legend_eyebrow", s.get("eyebrow"))}
            if "legend_items" in s:
                leg["items"] = norm_legend_items(s["legend_items"])
            out.append(leg)
            out.append(s)
        elif t == "agenda":
            items = s.get("items", [])
            st = resolve_style(layout, deck, s)
            b = st.get("body", {})
            pitch = b.get("row_h", 44) + b.get("gap", 10)
            R = max(1, int(b.get("h", 452) // pitch))
            cap = R * b.get("max_cols", 2)
            if len(items) > cap:
                for p in range(0, len(items), cap):
                    ns = dict(s)
                    ns["items"] = items[p:p + cap]
                    ns["num_start"] = p
                    if p > 0:
                        ns["title"] = s.get("title", "") + "（続き）"
                    out.append(ns)
            else:
                out.append(s)
        else:
            out.append(s)
    d = dict(deck); d["slides"] = out
    return d


# ---------------------------------------------------------------------------
# 文字あふれ推定（M-74）: title / lead / bullets 系 / table のオーバーフローを warn する
# ---------------------------------------------------------------------------

def _char_width_ratio(ch: str) -> float:
    """全角=1.0、半角=0.55（東アジアの幅プロパティ W/F を全角とみなす）。"""
    return 1.0 if unicodedata.east_asian_width(ch) in ("W", "F") else 0.55


def estimate_lines(text, size_px: float, width_px: float) -> int:
    """テキストが width_px の枠に折り返して何行になるかを概算する。
    全角=size×1.0、半角=size×0.55 として1文字ずつ幅を積算する簡易実装。"""
    text = str(text or "")
    if not text:
        return 1
    width_px = max(1.0, width_px)
    total_lines = 0
    for raw_line in text.split("\n"):
        if raw_line == "":
            total_lines += 1
            continue
        w = 0.0
        cur = 1
        for ch in raw_line:
            cw = size_px * _char_width_ratio(ch)
            if w > 0 and w + cw > width_px:
                cur += 1
                w = 0.0
            w += cw
        total_lines += cur
    return max(1, total_lines)


def estimate_bullets_height(items, r: dict) -> float:
    """bullets 領域（親＋子）が実際に描画されたときのおおよその高さを推定する（px）。
    items は norm_items() 前後どちらの形式でも受け付ける。"""
    items = norm_items(items)
    size = r.get("size", 16)
    lh = r.get("line_height", 1.4)
    gap = r.get("gap", 10)
    indent = r.get("indent", 30)
    width = max(10.0, r.get("w", 400) - indent)
    child_size = r.get("child_size", max(10, size - 3))
    child_gap = r.get("child_gap", 6)
    child_width = max(10.0, width - 20)
    total = 0.0
    for it in items:
        isize = it.get("size", size)
        lines = estimate_lines(it.get("text", ""), isize, width)
        total += lines * isize * lh + gap
        for c in it.get("children", []):
            csize = c.get("size", child_size)
            clines = estimate_lines(c.get("text", ""), csize, child_width)
            total += clines * csize * lh + child_gap
    return total


def steps_card_height(b: dict, steps, card_w: float) -> float:
    """steps のカード高さ。既定（fit="content"）では最も長いカードの内容に合わせて縮め、
    短い内容で下半分が空くのを防ぐ。fit="fill" で従来どおり body.h いっぱいに描く。
    下限 min_h（既定 260px）、上限 body.h。HTML/PPTX 共通。"""
    if b.get("fit", "content") != "content" or not steps:
        return b["h"]
    header_h = b.get("header_h", 60); pad = b.get("pad", 16)
    proto = {"w": card_w - 2*pad, "size": b.get("item_size", 15), "gap": b.get("item_gap", 8),
             "line_height": 1.4, "child_size": 13, "child_gap": 4, "indent": 16}
    est = max(estimate_bullets_height(norm_items(st.get("items", [])), proto) for st in steps)
    return max(min(b["h"], header_h + 2*pad + est + 12), min(b["h"], b.get("min_h", 260)))


def _matrix_content_layout(body_text, g: dict, cw: float, ch: float, pad: float, valign: str) -> dict:
    """matrix の1象限（heading + body の縦積み）の各要素の y オフセット（象限上端からの相対値）
    と body の高さを計算する（HTML/PPTX 共通）。valign:middle のときは heading を含めた
    コンテンツ全体の実高さを見積もり、その全体を象限の中央に配置できるよう heading の y から
    計算し直す（heading だけを上端固定にして body だけを残り高さでセンタリングすると、
    見出し直下に不自然な空白ができるため）。valign:top（既定）は従来どおり heading 上端固定。"""
    body_w = cw - 2 * pad
    head_h = 26
    head_gap = 4
    body_size = g.get("body_size", 14)
    if valign == "middle":
        lines = estimate_lines(body_text, body_size, body_w) if body_text else 0
        body_h = lines * body_size * 1.4
        content_h = head_h + (head_gap + body_h if body_text else 0)
        head_y = pad + max(0.0, (ch - 2 * pad - content_h) / 2)
    else:
        head_y = pad
        body_h = ch - 2 * pad - (head_h + head_gap)
    body_y = head_y + head_h + head_gap
    return {"head_h": head_h, "head_y": head_y, "body_y": body_y, "body_h": body_h}


def _cards_content_layout(card: dict, g: dict, cw: float, chh: float, pad: float, valign: str) -> dict:
    """cards の1枚（heading + body + items の縦積み）の各要素の y オフセット（カード上端からの
    相対値）と items の高さを計算する（HTML/PPTX 共通）。valign:middle のときは heading を含めた
    コンテンツ全体（実際に描画される段だけ）の実高さを見積もり、その全体をカードの中央に
    配置できるよう heading の y から計算し直す。valign:top（既定）は従来どおり heading 上端固定。"""
    body_w = cw - 2 * pad
    has_body = bool(card.get("body"))
    has_items = bool(card.get("items"))
    head_h = 28
    head_gap = 4
    body_h_fixed = 44
    body_gap = 4
    items_h = None
    if has_items:
        probe = {"size": g.get("item_size", 14), "gap": g.get("item_gap", 6), "line_height": 1.35,
                 "indent": 14, "child_size": 12, "child_gap": 4, "w": body_w}
        items_h = estimate_bullets_height(card.get("items"), probe)
    if valign == "middle":
        content_h = head_h
        if has_body:
            content_h += head_gap + body_h_fixed
        if has_items:
            content_h += body_gap + items_h
        head_y = pad + max(0.0, (chh - 2 * pad - content_h) / 2)
    else:
        head_y = pad
    yy = head_y + head_h + head_gap
    body_y = None
    if has_body:
        body_y = yy
        yy += body_h_fixed + body_gap
    items_y = None
    if has_items:
        items_y = yy
        if valign != "middle":
            items_h = chh - yy - pad
    return {"head_h": head_h, "head_y": head_y, "body_h": body_h_fixed, "body_y": body_y,
            "items_h": items_h, "items_y": items_y}


def _two_column_side_layout(colc: dict, box: dict, hd: dict, body_style: dict, valign: str) -> dict:
    """two_column の1サイド（heading + bullets の縦積み）の各要素の y オフセット（箱上端からの
    相対値）と bullets 領域の高さを計算する（PPTX 用。HTML 側は flexbox で同じ効果を実際の
    描画高さから得られるためこの見積もりを使わない）。style.left.valign / style.right.valign
    が "middle" のときは heading を含めたコンテンツ全体を箱の中央に配置できるよう heading の y
    から計算し直す。既定の "top" は従来どおり heading 上端固定。"""
    pad = box.get("pad", 0)
    inner_w = box.get("w", 400) - 2 * pad
    has_heading = bool(colc.get("heading"))
    head_h = (hd.get("size", 22) * 1.4) if has_heading else 0.0
    head_gap = hd.get("gap_below", 12) if has_heading else 0.0
    probe = dict(body_style)
    probe["w"] = inner_w
    body_h = estimate_bullets_height(colc.get("bullets"), probe) if colc.get("bullets") else 0.0
    if valign == "middle":
        content_h = head_h + head_gap + body_h
        avail_h = box.get("h", 300) - 2 * pad
        head_y = pad + max(0.0, (avail_h - content_h) / 2)
    else:
        head_y = pad
    body_y = head_y + head_h + head_gap
    if valign == "middle":
        final_body_h = body_h
    else:
        final_body_h = box.get("h", 300) - pad - body_y
    return {"head_y": head_y, "head_h": head_h, "body_y": body_y, "body_h": final_body_h}


def closing_message_style(st: dict, slide: dict) -> dict:
    """closing.message の描画スタイルを解決する（M-38: HTML/PPTX 共通の計算式）。
    message.follow_body が真で、かつスライドの style.message.y が明示されていなければ、
    message.y を「body の実描画高さの直後 + follow_gap」に自動配置する
    （上限は layout 既定の message.y、下限は min_y）。follow_body が無い、または
    style.message.y が明示されている場合は st["message"] をそのまま返す。"""
    msg = dict(st["message"])
    style_msg = (slide.get("style") or {}).get("message", {})
    if msg.get("follow_body") and "y" not in style_msg:
        body = st["body"]
        bullets = slide.get("bullets") or []
        body_bottom = body["y"] + (estimate_bullets_height(bullets, body) if bullets else 0)
        max_y = msg["y"]
        min_y = msg.get("min_y", max_y)
        msg["y"] = min(max_y, max(min_y, body_bottom + msg.get("follow_gap", 48)))
    return msg


def _check_overflow(slide: dict, st: dict, i: int, t: str):
    """title / lead / bullets 系 body / table のあふれを見積もって warn する（warn のみ、fail しない）。"""
    footer = st.get("footer_l") or st.get("footer_r")
    bottom = footer["y"] if footer else CANVAS_H - 20
    threshold = 24  # このくらいの超過までは丸め誤差として無視する

    def region_h(region: dict) -> float:
        if "h" in region:
            return region["h"]
        return max(0.0, bottom - region.get("y", 0))

    def warn_text(label: str, region: dict, text) -> None:
        if not region or not text:
            return
        h = region_h(region)
        lines = estimate_lines(text, region.get("size", 16), region.get("w", 400))
        est = lines * region.get("size", 16) * region.get("line_height", 1.2)
        over = est - h
        if over > threshold:
            warn(f"{i + 1}枚目（type={t}）: {label} が枠の高さを約{int(over)}px超える見込みです。"
                 f"文字数を減らすか style.{label}.size を下げてください")

    def warn_bullets(label: str, region: dict, items) -> None:
        if not region or not items:
            return
        h = region_h(region)
        est = estimate_bullets_height(items, region)
        over = est - h
        if over > threshold:
            warn(f"{i + 1}枚目（type={t}）: {label} が枠の高さを約{int(over)}px超える見込みです。"
                 f"項目を減らすか style.{label}.size を下げてください")

    warn_text("title", st.get("title"), slide.get("title"))
    warn_text("lead", st.get("lead"), slide.get("lead"))

    if t == "bullets":
        warn_bullets("body", st.get("body"), slide.get("bullets"))
    elif t == "two_column":
        for side in ("left", "right"):
            box = st.get(side)
            colc = slide.get(side, {})
            if not box or not colc.get("bullets"):
                continue
            pad = box.get("pad", 0)
            inner = dict(st.get("col_body", {}))
            inner["w"] = box.get("w", 400) - 2 * pad
            inner["h"] = box.get("h", 300) - 2 * pad
            warn_bullets(f"{side}.bullets", inner, colc.get("bullets"))
    elif t == "closing":
        warn_bullets("body", st.get("body"), slide.get("bullets"))
    elif t == "image_text":
        warn_bullets("body", st.get("body"), slide.get("bullets"))
    elif t == "table":
        tr = st.get("table")
        rows = slide.get("rows", [])
        if tr and rows:
            est = tr.get("header_h", 0) + tr.get("row_h", 0) * len(rows)
            h = region_h(tr)
            over = est - h
            if over > threshold:
                warn(f"{i + 1}枚目（type=table）: 表が枠の高さを約{int(over)}px超える見込みです。"
                     f"行数を減らすか style.table.row_h を下げてください")


# ---------------------------------------------------------------------------
# 色指定の検証（M-68 等）: テーマトークン名 or #RRGGBB/#RGB のみを許可する
# ---------------------------------------------------------------------------

_COLOR_KEY_RE = re.compile(r"(color|fill|fills)$")


def _color_issue(theme, value, path: str, i: int, t: str):
    if not isinstance(value, str):
        return None
    if is_color_value(theme, value):
        return None
    return (f"不明な色指定 '{value}'（{i + 1}枚目 type={t} の {path}）。"
            f"テーマトークン名か #RRGGBB 形式のhexを指定してください")


def _walk_colors(obj, path: str, theme, i: int, t: str, issues: list):
    if isinstance(obj, dict):
        for k, v in obj.items():
            npath = f"{path}.{k}" if path else k
            if _COLOR_KEY_RE.search(k):
                if isinstance(v, dict):
                    for kk, vv in v.items():
                        msg = _color_issue(theme, vv, f"{npath}[{kk}]", i, t)
                        if msg:
                            issues.append(msg)
                elif isinstance(v, list):
                    for idx, vv in enumerate(v):
                        msg = _color_issue(theme, vv, f"{npath}[{idx}]", i, t)
                        if msg:
                            issues.append(msg)
                else:
                    msg = _color_issue(theme, v, npath, i, t)
                    if msg:
                        issues.append(msg)
            elif isinstance(v, (dict, list)):
                _walk_colors(v, npath, theme, i, t, issues)
    elif isinstance(obj, list):
        for idx, v in enumerate(obj):
            if isinstance(v, (dict, list)):
                _walk_colors(v, f"{path}[{idx}]", theme, i, t, issues)


# ---------------------------------------------------------------------------
# 画像パスの検証（X-1 パストラバーサル対策 / M-72, M-96 存在チェック）
# ---------------------------------------------------------------------------

ALLOWED_IMAGE_EXT = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".svg", ".webp", ".tiff"}


def _validate_image_path(slide: dict, i: int, t: str, deck_dir: Path):
    rel = slide.get("path")
    if not rel:
        fail(f"{i + 1}枚目（type={t}）: path が指定されていません")
    try:
        base = Path(deck_dir).resolve()
        resolved = (base / rel).resolve()
    except (OSError, ValueError):
        fail(f"{i + 1}枚目（type={t}）: path '{rel}' を解決できません")
    if resolved != base and not resolved.is_relative_to(base):
        fail(f"{i + 1}枚目（type={t}）: path '{rel}' がデッキフォルダの外を指しています"
             f"（パストラバーサルは禁止されています。assets/ 配下の相対パスを指定してください）")
    ext = resolved.suffix.lower()
    if ext not in ALLOWED_IMAGE_EXT:
        fail(f"{i + 1}枚目（type={t}）: 画像形式 '{ext or rel}' には対応していません"
             f"（対応拡張子: {', '.join(sorted(e.lstrip('.') for e in ALLOWED_IMAGE_EXT))}）")
    if not resolved.exists():
        fail(f"{i + 1}枚目（type={t}）: 画像ファイルが見つかりません: {rel}")


# ---------------------------------------------------------------------------
# swimlane の検証（M-77 node id 重複 / edges 未知 id / lanes 空、M-53 legend_items の形）
# ---------------------------------------------------------------------------

def _validate_legend_items(items, i: int, t: str):
    if not isinstance(items, list):
        fail(f"{i + 1}枚目（type={t}）: legend_items はリスト（配列）で指定してください")
    for idx, item in enumerate(items):
        if isinstance(item, (list, tuple)):
            if len(item) != 4:
                fail(f"{i + 1}枚目（type={t}）: legend_items[{idx}] は4要素の配列"
                     f"（[shape, variant_or_kind, label, description]）である必要があります"
                     f"（{len(item)}要素でした）")
        elif isinstance(item, dict):
            # shape は必須ではなく、kind だけでも図形を解決できる（norm_legend_items 参照）。
            has_shape = "shape" in item or "kind" in item
            has_text = "label" in item and "desc" in item
            if not (has_shape and has_text):
                fail(f"{i + 1}枚目（type={t}）: legend_items[{idx}] の辞書形式には "
                     f"shape（または kind）/ label / desc がすべて必要です")
        else:
            fail(f"{i + 1}枚目（type={t}）: legend_items[{idx}] の形式が不正です"
                 f"（4要素配列、または {{shape, variant, label, desc}} の辞書を指定してください）")


SWIMLANE_TASK_VARIANTS = ("onpf", "onother", "offline")
SWIMLANE_SHAPES = ("task", "system", "decision", "terminal", "marker", "connector", "mail", "io")


def _validate_swimlane(slide: dict, i: int, t: str):
    if not slide.get("lanes"):
        fail(f"{i + 1}枚目（type=swimlane）: lanes が空です")
    lanes = slide["lanes"]
    if not isinstance(lanes, list):
        fail(f"{i + 1}枚目（type=swimlane）: lanes は配列（リスト）で指定してください")
    n_lanes = len(lanes)

    nodes = slide.get("nodes", [])
    if not isinstance(nodes, list):
        fail(f"{i + 1}枚目（type=swimlane）: nodes は配列（リスト）で指定してください")
    edges = slide.get("edges", [])
    if not isinstance(edges, list):
        fail(f"{i + 1}枚目（type=swimlane）: edges は配列（リスト）で指定してください")

    seen_ids = set()
    for nd in nodes:
        if not isinstance(nd, dict):
            fail(f"{i + 1}枚目（type=swimlane）: nodes の要素が正しい形式（オブジェクト）ではありません")
        nid = nd.get("id")
        if nid is None:
            fail(f"{i + 1}枚目（type=swimlane）: ノードに id が指定されていません")
        if nid in seen_ids:
            fail(f"{i + 1}枚目（type=swimlane）: ノード id '{nid}' が重複しています")
        seen_ids.add(nid)
        lane = nd.get("lane", 0)
        if isinstance(lane, bool) or not isinstance(lane, int) or not (0 <= lane < n_lanes):
            fail(f"{i + 1}枚目（type=swimlane）: ノード '{nid}' の lane は 0〜{n_lanes - 1} の整数で"
                 f"指定してください（lanes の行 index。id文字列ではありません）")
        col = nd.get("col", 0)
        if isinstance(col, bool) or not isinstance(col, int) or col < 0:
            fail(f"{i + 1}枚目（type=swimlane）: ノード '{nid}' の col は0以上の整数で指定してください")
        # variant は task ノード専用の語彙（onpf / onother / offline）。新図解タイプの variant
        # （emphasis / security …）とは別の名前空間なので、混同した値は warning で知らせる
        # （未知の値は無警告で onother 扱いになり、typo が成果物に紛れ込むため）。
        variant = nd.get("variant")
        if variant is not None and variant not in SWIMLANE_TASK_VARIANTS:
            hint = difflib.get_close_matches(str(variant), SWIMLANE_TASK_VARIANTS, n=1)
            hint_s = f"（候補: {hint[0]}）" if hint else ""
            warn(f"{i + 1}枚目（type=swimlane）: ノード '{nid}' の variant '{variant}' は swimlane では "
                 f"onpf / onother / offline のいずれかです{hint_s}（それ以外は onother として描かれます。"
                 f"emphasis 等の強調は architecture 系タイプの語彙で、swimlane では使えません）")
        shape = nd.get("shape", "task")
        if shape not in SWIMLANE_SHAPES:
            hint = difflib.get_close_matches(str(shape), SWIMLANE_SHAPES, n=1)
            hint_s = f"（候補: {hint[0]}）" if hint else ""
            warn(f"{i + 1}枚目（type=swimlane）: ノード '{nid}' の shape '{shape}' は未対応です{hint_s}"
                 f"（task として描かれます）")
        if "text" not in nd and "label" not in nd and shape not in ("marker", "mail", "io"):
            warn(f"{i + 1}枚目（type=swimlane）: ノード '{nid}' に text がありません（表示名は text に書く。label も受理）")
    for ei, e in enumerate(edges):
        if not isinstance(e, dict):
            fail(f"{i + 1}枚目（type=swimlane）: edges の要素が正しい形式（オブジェクト）ではありません")
        for key in ("from", "to"):
            eid = e.get(key)
            if eid not in seen_ids:
                fail(f"{i + 1}枚目（type=swimlane）: edges の {key} '{eid}' は未知のノード id です")
        ev = e.get("variant")
        if ev is not None and ev not in de.EDGE_VARIANTS:
            warn(f"{i + 1}枚目（type=swimlane）: edges[{ei}] の variant '{ev}' は未対応です"
                 f"（default / emphasis / security / dashed のいずれか）")
        es = e.get("style")
        if es is not None and es not in ("solid", "dashed"):
            warn(f"{i + 1}枚目（type=swimlane）: edges[{ei}] の style '{es}' は solid / dashed のいずれかです"
                 f"（強調などは variant で指定する）")
    if "legend_items" in slide:
        _validate_legend_items(slide["legend_items"], i, t)


# ---------------------------------------------------------------------------
# デッキ全体の検証（ビルド前に一括実行）
# ---------------------------------------------------------------------------

_ID_FORBIDDEN_RE = re.compile(r'[\\/:*?"<>|\x00-\x1f]')


def _report_diagnostics(diags, i: int, t: str, geometry_only: bool = False):
    """diagram_engine の診断を build の error / warning / note に振り分けて表示する。
    - スキーマ的な誤り（id 重複・未知 id・セル衝突など）の error はビルドを止める（fail）
    - geometry_only=True（配線・ラベル・文字あふれなど幾何の診断）のときは error でもビルドは
      止めず warning として報告する。厳密に通したいときは check_diagram.py --strict を使う
      （Archify の validate に相当。修理指示 fixes を「→」以降に表示する）"""
    for dg in diags:
        msg = f"{i + 1}枚目（type={t}）: {de.format_diagnostic(dg)}"
        level = dg.get("level")
        if level == "error" and not geometry_only:
            fail(msg)
        elif level in ("error", "warning"):
            warn(msg)
        else:
            note(msg)


def validate_deck(deck: dict, layout: dict, theme: dict, deck_dir: Path):
    """deck.json をビルド前に一括検証する。問題があれば fail()（即終了）または warn()/note()。"""
    if not isinstance(deck, dict):
        fail("deck.json の形式が不正です（トップレベルはオブジェクトである必要があります）")
    meta = deck.get("meta", {})
    if not isinstance(meta, dict):
        fail("deck.json の形式が不正です（meta はオブジェクトである必要があります）")

    # meta.id: 省略可。禁止文字は fail、フォルダ名と異なる場合は note。
    raw_id = meta.get("id")
    if raw_id is not None:
        if _ID_FORBIDDEN_RE.search(str(raw_id)):
            fail(f"meta.id に使用できない文字が含まれています: '{raw_id}'"
                 f"（\\ / : * ? \" < > | と制御文字は使用できません）")
        if deck_dir is not None and str(raw_id) != Path(deck_dir).name:
            note(f"meta.id（{raw_id}）とフォルダ名（{Path(deck_dir).name}）が一致していません")

    types_known = set(layout.get("types", {}).keys())
    common_keys = set(layout.get("common", {}).keys())

    # meta.layout_overrides のトップキー検証（"common" は特別扱い）
    overrides = meta.get("layout_overrides", {})
    if not isinstance(overrides, dict):
        fail("meta.layout_overrides はオブジェクトである必要があります")
    for ov_type, ov_val in overrides.items():
        if ov_type == "common":
            for key in (ov_val or {}):
                if key not in common_keys:
                    hint = difflib.get_close_matches(key, common_keys, n=1)
                    hint_s = f"（候補: {hint[0]}）" if hint else ""
                    warn(f"meta.layout_overrides.common: '{key}' は共通領域名ではありません{hint_s}")
            continue
        if ov_type not in types_known:
            hint = difflib.get_close_matches(ov_type, types_known, n=1)
            hint_s = f"（候補: {hint[0]}）" if hint else ""
            fail(f"meta.layout_overrides に未知のスライドタイプ '{ov_type}' があります{hint_s}")
        allowed = common_keys | set(layout["types"][ov_type].keys())
        for key in (ov_val or {}):
            if key not in allowed:
                hint = difflib.get_close_matches(key, allowed, n=1)
                hint_s = f"（候補: {hint[0]}）" if hint else ""
                warn(f"meta.layout_overrides.{ov_type}: '{key}' は {ov_type} の領域名ではありません{hint_s}")

    slides = deck.get("slides", [])
    if not isinstance(slides, list):
        fail("slides は配列（リスト）である必要があります")
    if not slides:
        warn("slides が空です。1枚もスライドがありません")

    for i, slide in enumerate(slides):
        if not isinstance(slide, dict):
            fail(f"{i + 1}枚目のスライドが正しい形式（オブジェクト）ではありません")

        t = slide.get("type")
        if not t:
            fail(f"{i + 1}枚目のスライドに type がありません")
        if t not in types_known:
            hint = difflib.get_close_matches(t, types_known, n=1)
            hint_s = f"（候補: {hint[0]}）" if hint else ""
            fail(f"{i + 1}枚目: 未知のスライドタイプ: {t}{hint_s}")

        # bullets は配列（リスト）以外だと norm_items() が文字列を1文字ずつのイテレーションとして
        # 扱ってしまい、クラッシュせず「1文字ずつの箇条書き」という無音の誤動作になる。
        bullets_field = slide.get("bullets")
        if bullets_field is not None and not isinstance(bullets_field, list):
            fail(f"{i + 1}枚目（type={t}）: bullets は配列（リスト）で指定してください")

        # slide.style は必ずオブジェクト（またはキー欠落）。文字列等だと deep_merge が
        # 解決済みスタイルを丸ごと非dict値に置き換えてしまい、後段の全処理がクラッシュする。
        slide_style = slide.get("style", {})
        if slide_style and not isinstance(slide_style, dict):
            fail(f"{i + 1}枚目（type={t}）: style はオブジェクトである必要があります")

        # slide.style のキー名タイポ検証（M-75）
        allowed = common_keys | set(layout["types"][t].keys())
        for key in slide_style:
            if key not in allowed:
                hint = difflib.get_close_matches(key, allowed, n=1)
                hint_s = f"（候補: {hint[0]}）" if hint else ""
                warn(f"{i + 1}枚目（type={t}）: style の '{key}' は {t} の領域名ではありません{hint_s}")

        st = resolve_style(layout, deck, slide)

        # 色検証: 解決済みスタイル辞書と、スライド生データ（bullets の color 等）の両方を走査
        issues: list = []
        _walk_colors(st, "", theme, i, t, issues)
        _walk_colors(slide, "", theme, i, t, issues)
        if issues:
            fail(issues[0])

        if t == "table":
            cols = slide.get("columns", [])
            for ridx, row in enumerate(slide.get("rows", [])):
                if len(row) != len(cols):
                    fail(f"{i + 1}枚目（type=table）: {ridx + 1}行目のセル数（{len(row)}）が "
                         f"columns 数（{len(cols)}）と一致しません")
            widths = st.get("table", {}).get("col_widths")
            if widths and len(widths) != len(cols):
                fail(f"{i + 1}枚目（type=table）: col_widths の要素数（{len(widths)}）が "
                     f"columns 数（{len(cols)}）と一致しません")

        if t in ("image", "image_text"):
            _validate_image_path(slide, i, t, deck_dir)

        if t == "matrix":
            quadrants = slide.get("quadrants", [])
            if not isinstance(quadrants, list):
                fail(f"{i + 1}枚目（type=matrix）: quadrants は配列（リスト）で指定してください")
            n = len(quadrants)
            if n != 4:
                warn(f"{i + 1}枚目（type=matrix）: quadrants は4つの想定ですが{n}個指定されています"
                     f"（5個以上は無視され、4個未満は空の象限になります）")

        if t == "cards":
            columns = slide.get("columns")
            if columns is not None and not isinstance(columns, int):
                fail(f"{i + 1}枚目（type=cards）: columns は数値で指定してください")

        if t == "swimlane":
            _validate_swimlane(slide, i, t)
            _report_diagnostics(swimlane_geometry(slide, st, i + 1).get("diagnostics", []), i, t,
                                geometry_only=True)

        if t in GRID_DIAGRAM_TYPES:
            _report_diagnostics(de.validate_grid_diagram(slide, t), i, t)
            geo = de.layout_grid_diagram(slide, t, st["diagram"], [])
            _report_diagnostics(geo["diagnostics"], i, t, geometry_only=True)

        if t == "sequence":
            diags = []
            de.layout_sequence(slide, st["diagram"], diags)
            _report_diagnostics(diags, i, t)

        _check_overflow(slide, st, i, t)


def fit_code_size(code: str, r: dict):
    """コード領域に収まるフォントサイズを計算（HTML/PPTX 共通ロジック）。"""
    lines = code.splitlines() or [""]
    size = r["size"]
    lh = r["line_height"]
    avail = r["h"] - 2 * r["pad"]
    while size > r.get("min_size", 12) and len(lines) * size * lh > avail:
        size -= 1
    return size, lines


def is_http_url(value) -> bool:
    """Return True only when the entire table-cell value is an HTTP(S) URL."""
    return bool(re.fullmatch(r"https?://\S+", str(value)))


def table_cell_parts(value):
    """Normalize a table cell into visible text and an optional external URL."""
    if value is None:
        return "", None
    if isinstance(value, dict):
        text = str(value.get("text", ""))
        url = value.get("url")
        return text, str(url) if url and is_http_url(url) else None
    text = str(value)
    return text, text if is_http_url(text) else None


def svg_size(svg_text: str):
    """SVG の描画サイズを viewBox / width / height 属性から取得する。"""
    m = re.search(
        r'viewBox\s*=\s*"\s*[-\d.]+[\s,]+[-\d.]+[\s,]+([\d.]+)[\s,]+([\d.]+)\s*"',
        svg_text,
    )
    if m:
        return float(m.group(1)), float(m.group(2))
    wm = re.search(r'<svg[^>]*?\swidth\s*=\s*"([\d.]+)(?:px)?"', svg_text)
    hm = re.search(r'<svg[^>]*?\sheight\s*=\s*"([\d.]+)(?:px)?"', svg_text)
    if wm and hm:
        return float(wm.group(1)), float(hm.group(1))
    return 1136.0, 440.0


def svg_to_png(svg_path: Path, out_png: Path, browser=None):
    """SVG を透過 PNG にラスタライズ（PPTX 埋め込み用、2x 解像度）。
    browser に既存の Playwright Browser を渡すと、それを再利用する
    （M-25: デッキ内に複数 SVG があっても Chromium の起動は1回で済む）。
    省略時は従来どおり自分で Playwright/Chromium を起動して使い終わったら閉じる。"""
    svg_text = svg_path.read_text(encoding="utf-8")
    w, h = svg_size(svg_text)

    def _render(b):
        page = b.new_page(
            viewport={"width": max(1, int(w)), "height": max(1, int(h))},
            device_scale_factor=2,
        )
        try:
            page.set_content(
                "<!doctype html><style>html,body{margin:0;background:transparent}"
                "svg{display:block}</style>" + svg_text
            )
            page.wait_for_timeout(120)
            page.screenshot(path=str(out_png), omit_background=True,
                            clip={"x": 0, "y": 0, "width": w, "height": h})
        finally:
            page.close()

    if browser is not None:
        _render(browser)
        return
    from playwright.sync_api import sync_playwright
    with sync_playwright() as pw:
        b = pw.chromium.launch()
        try:
            _render(b)
        finally:
            b.close()


def flip_regions(st: dict, keys) -> dict:
    """image_side: "left" 用に、指定領域の x をキャンバス中心で左右反転する。"""
    out = copy.deepcopy(st)
    for k in keys:
        r = out.get(k)
        if isinstance(r, dict) and "x" in r and "w" in r:
            r["x"] = CANVAS_W - r["x"] - r["w"]
    return out


# webp/tiff/bmp は python-pptx が扱えない、または PowerPoint との互換性が低いため
# PPTX 埋め込み時のみ PNG に変換する（M-72）。HTML はブラウザがそのまま表示できるので対象外。
PPTX_PNG_CONVERT_EXTS = (".webp", ".tiff", ".tif", ".bmp")


def _image_has_alpha(im) -> bool:
    if im.mode in ("RGBA", "LA"):
        return True
    if im.mode == "P":
        return "transparency" in im.info
    return False


def _resample_image_cache(deck_dir: Path, img_path: Path, target_w, target_h, force_convert: bool):
    """Pillow で画像をリサイズ・再圧縮／形式変換してキャッシュする
    （X-4: リサイズ・再圧縮キャッシュ、M-72: webp/tiff/bmp の PPTX 用 PNG 変換）。
    - target_w/target_h（表示枠 px）が指定されていれば、長辺が表示枠の2倍を超える場合のみ
      縮小する（拡大はしない）。アニメーション画像（複数フレームの gif/webp 等）はアニメーションが
      壊れるためリサイズ対象から除外する。
    - force_convert が真なら、サイズに関わらず PNG（または JPEG）へ変換したキャッシュを作る
      （webp/tiff/bmp を PPTX 埋め込み可能な形式にするため。アニメーションがある場合は先頭フレーム
      のみになる）。
    - アルファチャンネルが無ければ JPEG（品質85）、あれば PNG で保存する。
    - 条件を満たさない（リサイズも変換も不要）場合、または Pillow が無い／変換に失敗した場合は
      None を返す（呼び出し元は元ファイルをそのまま使う）。"""
    try:
        from PIL import Image as PILImage
    except ImportError:
        warn("Pillow が見つからないため画像のリサイズ・再圧縮・形式変換をスキップします"
             "（PowerPoint 非対応形式の画像がある場合は PPTX 生成に失敗する可能性があります。"
             "python-pptx を導入すれば通常 Pillow も同梱されます）")
        return None
    try:
        with PILImage.open(img_path) as im:
            is_animated = bool(getattr(im, "is_animated", False))
            im.load()
            w, h = im.size
            cap_w = target_w * 2 if target_w else None
            cap_h = target_h * 2 if target_h else None
            need_resize = (not is_animated) and cap_w and cap_h and (w > cap_w or h > cap_h)
            if not need_resize and not force_convert:
                return None
            has_alpha = _image_has_alpha(im)
            out_ext = ".png" if has_alpha else ".jpg"
            cache_dir = deck_dir / "build" / ".img-cache"
            cache_dir.mkdir(parents=True, exist_ok=True)
            tag = f"{int(cap_w or 0)}x{int(cap_h or 0)}" if need_resize else "conv"
            # 元のファイル名（拡張子込み）とパスのハッシュをキャッシュ名に含める。stem のみだと
            # 同一フォルダ内の拡張子違い（foo.png/foo.jpg）や別フォルダの同名ファイルが
            # 衝突してキャッシュを取り違えるため。
            path_hash = hashlib.sha1(str(img_path.resolve()).encode("utf-8")).hexdigest()[:8]
            out_path = cache_dir / f"{img_path.stem}-{path_hash}-{tag}{out_ext}"
            if out_path.exists() and out_path.stat().st_mtime >= img_path.stat().st_mtime:
                return out_path
            frame = im
            if need_resize:
                ratio = min(cap_w / w, cap_h / h)
                if ratio < 1.0:
                    frame = im.resize((max(1, round(w * ratio)), max(1, round(h * ratio))), PILImage.LANCZOS)
            if out_ext == ".jpg":
                if frame.mode != "RGB":
                    frame = frame.convert("RGB")
                frame.save(out_path, "JPEG", quality=85)
            else:
                if frame.mode not in ("RGBA", "RGB", "LA", "P"):
                    frame = frame.convert("RGBA" if has_alpha else "RGB")
                frame.save(out_path, "PNG")
            return out_path
    except Exception as e:
        warn(f"画像 '{img_path.name}' のリサイズ・再圧縮・形式変換に失敗したため原本を使用します（{e}）")
        return None


def image_source(deck_dir: Path, rel_path: str, for_pptx: bool, browser=None,
                  target_w=None, target_h=None) -> Path:
    """image/image_text スライドの実ファイルを解決する。HTML の data URI・PPTX 埋め込みの
    両方から呼ばれる共通の入口（X-4）。
    - SVG: HTML はそのまま返す。PPTX はラスタライズして PNG 化しキャッシュする
      （browser に既存の Playwright Browser を渡すと svg_to_png() がそれを再利用する。M-25）。
    - webp/tiff/bmp: for_pptx のときだけ PNG 化する（M-72）。
    - それ以外（png/jpg/jpeg/gif 等）: target_w/target_h（表示枠 px）が指定されていれば、
      deck_dir/build/.img-cache/ に「長辺が表示枠の2倍」を上限にリサイズ・再圧縮した版を作り
      それを返す（X-4）。HTML と PPTX の両方が同じキャッシュを共有するため、リサイズは一度だけ行う。
      Pillow が無い、またはリサイズ・変換が不要なときは元ファイルをそのまま返す。"""
    img_path = deck_dir / rel_path
    ext = img_path.suffix.lower()
    if ext == ".svg":
        if not for_pptx:
            return img_path
        cache_dir = deck_dir / "build" / ".img-cache"
        cache_dir.mkdir(parents=True, exist_ok=True)
        path_hash = hashlib.sha1(str(img_path.resolve()).encode("utf-8")).hexdigest()[:8]
        png = cache_dir / f"{img_path.stem}-{path_hash}.svg.png"
        if not png.exists() or png.stat().st_mtime < img_path.stat().st_mtime:
            svg_to_png(img_path, png, browser=browser)
        return png
    force_convert = for_pptx and ext in PPTX_PNG_CONVERT_EXTS
    cached = _resample_image_cache(deck_dir, img_path, target_w, target_h, force_convert)
    return cached if cached is not None else img_path


# ---------------------------------------------------------------------------
# HTML ビルド
# ---------------------------------------------------------------------------

def esc(s):
    return html_mod.escape(str(s), quote=True)


def pos(r, extra=""):
    s = (f"position:absolute;left:{r['x']}px;top:{r['y']}px;width:{r['w']}px;")
    if "h" in r:
        s += f"height:{r['h']}px;"
    return s + extra


def text_style(r, theme):
    s = f"font-size:{r['size']}px;color:{col(theme, r.get('color', 'text'))};"
    if r.get("bold"):
        s += "font-weight:700;"
    if r.get("align"):
        s += f"text-align:{r['align']};"
    if r.get("line_height"):
        s += f"line-height:{r['line_height']};"
    return s


def div(r, theme, inner, extra="", font_family=None, heading=False):
    """テキストの絶対配置 div を1つ作る。
    font_family: font_css() 済みの CSS 値を渡すと font-family を明示する（M-21/M-64/M-83: 見出し系に
    fonts.heading を適用する際に使う。未指定なら .slide の既定フォント（fonts.body）を継承する）。
    heading=True でその要素をスクリーンリーダー向けの見出しとしてマークする
    （role="heading" aria-level="1"。X-3: タイトル要素のアクセシビリティ対応）。
    font_family は font_css() が返す `"名前1", "名前2", ...` という値（二重引用符を含む）を
    そのまま style 属性（二重引用符区切り）に埋めると属性値がそこで打ち切られ、以降の
    プロパティが丸ごと無効になる（P0: code スライドの本文が判読不能になる不具合の原因）。
    esc() で `"` を `&quot;` にエスケープしてから埋め込む（ブラウザは属性値を文字参照
    展開してから CSS として使うため、フォント名自体は変わらない）。
    inner中の生の改行(\n)は<br>に変換する（P2: title等の明示的な改行がPPTXでは反映されるのに
    HTMLのwhite-space:normalではスペースに畳まれて消える食い違いの修正。呼び出し側は
    esc()済みの文字列を渡す前提なので、ここでの変換対象はエスケープ後に残る生の\nのみで、
    HTMLタグを誤って壊す心配はない）。"""
    ff = f"font-family:{esc(font_family)};" if font_family else ""
    attrs = ' role="heading" aria-level="1"' if heading else ""
    if isinstance(inner, str) and "\n" in inner:
        inner = inner.replace("\n", "<br>")
    return f'<div{attrs} style="{pos(r)}{text_style(r, theme)}{ff}{extra}">{inner}</div>'


def rect(r, theme, extra=""):
    fill = col(theme, r.get("fill", "accent"))
    radius = f"border-radius:{r['radius']}px;" if r.get("radius") else ""
    return f'<div style="{pos(r)}background:{fill};{radius}{extra}"></div>'


def bullets_html(items, r, theme):
    mk = col(theme, r.get("marker_color", "accent"))
    parts = [
        f'<ul style="margin:0;padding-left:{r["indent"]}px;list-style:disc;'
        f'font-size:{r["size"]}px;line-height:{r["line_height"]};'
        f'color:{col(theme, r["color"])};--mk:{mk};">'
    ]
    for it in items:
        li = f'font-size:{it["size"]}px;' if it.get("size") else ""
        li += f'color:{col(theme, it["color"])};' if it.get("color") else ""
        li += "font-weight:700;" if it.get("bold") else ""
        parts.append(f'<li style="margin-bottom:{r["gap"]}px;{li}">{esc(it["text"])}')
        if it["children"]:
            parts.append(
                f'<ul style="margin:{r["child_gap"]}px 0 0;padding-left:{r["indent"] - 6}px;'
                f'list-style:circle;font-size:{r["child_size"]}px;'
                f'color:{col(theme, r["child_color"])};">'
            )
            for c in it["children"]:
                cs = f'font-size:{c["size"]}px;' if c.get("size") else ""
                cs += f'color:{col(theme, c["color"])};' if c.get("color") else ""
                parts.append(f'<li style="margin-bottom:{r["child_gap"]}px;{cs}">{esc(c["text"])}</li>')
            parts.append("</ul>")
        parts.append("</li>")
    parts.append("</ul>")
    return "".join(parts)


def footer_html(st, deck, theme, page, total):
    out = ""
    if "footer_l" in st:
        out += div(st["footer_l"], theme, esc(deck.get("meta", {}).get("title", "")))
    if "footer_r" in st:
        out += div(st["footer_r"], theme, esc(st["footer_r"].get("text") or f"{page} / {total}"))
    return out


def html_slide_body(slide, st, deck, theme, deck_dir, page, total):
    t = slide["type"]
    fonts = theme["fonts"]
    parts = []
    # M-21/M-64/M-83: 見出し系要素（title/eyebrow/section.number/quote.text/col_heading/
    # cards・matrix の heading/steps の label/image_text の punch/closing.message）には
    # fonts.heading（未指定なら body）を適用する。PPTX 側の HEADING_FONT と同じ考え方。
    HEADING_FONT_CSS = font_css(fonts.get("heading") or fonts["body"])

    def chrome():
        if slide.get("eyebrow") and "eyebrow" in st:
            parts.append(div(st["eyebrow"], theme, esc(slide["eyebrow"]), font_family=HEADING_FONT_CSS))
        parts.append(div(st["title"], theme, esc(slide.get("title", "")), font_family=HEADING_FONT_CSS, heading=True))
        parts.append(rect(st["rule"], theme))
        if slide.get("lead") and "lead" in st:
            parts.append(div(st["lead"], theme, esc(slide["lead"])))
        parts.append(footer_html(st, deck, theme, page, total))

    if t == "title":
        parts.append(rect(st["bar"], theme))
        parts.append(div(st["title"], theme, esc(slide.get("title", "")), font_family=HEADING_FONT_CSS, heading=True))
        if slide.get("subtitle"):
            parts.append(div(st["subtitle"], theme, esc(slide["subtitle"])))
        if slide.get("meta"):
            parts.append(div(st["meta"], theme, esc(slide["meta"])))

    elif t == "section":
        if slide.get("number"):
            parts.append(div(st["number"], theme, esc(slide["number"]), font_family=HEADING_FONT_CSS))
        parts.append(div(st["title"], theme, esc(slide.get("title", "")), font_family=HEADING_FONT_CSS, heading=True))
        parts.append(rect(st["rule"], theme))
        if slide.get("subtitle"):
            parts.append(div(st["subtitle"], theme, esc(slide["subtitle"])))

    elif t == "bullets":
        chrome()
        # style.body.valign（既定 top）: "middle" で項目群を body 枠の縦中央に寄せる（項目数が
        # 少ないデッキでキャンバス下半分だけ空白になる問題への対応。two_column/matrix/cards と
        # 同じ考え方。PPTX 側は text_frame.vertical_anchor=MIDDLE で対応、こちらは実際の描画高さ
        # で自動的に中央寄せされる flex を使うため見積もり計算は不要）。
        body_valign = st["body"].get("valign", "top")
        body_extra = "display:flex;flex-direction:column;justify-content:center;" if body_valign == "middle" else ""
        parts.append(
            f'<div style="{pos(st["body"])}{body_extra}">'
            + bullets_html(norm_items(slide.get("bullets")), st["body"], theme)
            + "</div>"
        )

    elif t == "two_column":
        chrome()
        for side in ("left", "right"):
            box = st[side]
            colc = slide.get(side, {})
            radius = f"border-radius:{box.get('radius', 0)}px;"
            inner = ""
            hd = st["col_heading"]
            if colc.get("heading"):
                inner += (
                    f'<div style="font-size:{hd["size"]}px;color:{col(theme, hd["color"])};'
                    # P0: font_css() の戻り値は二重引用符を含むため esc() でエスケープしてから
                    # style 属性（二重引用符区切り）に埋め込む（そのまま埋めると属性値がここで
                    # 打ち切られ、font-weight/margin-bottom が丸ごと無効になる）。
                    f'font-family:{esc(HEADING_FONT_CSS)};'
                    f'font-weight:700;margin-bottom:{hd["gap_below"]}px;">{esc(colc["heading"])}</div>'
                )
            inner += bullets_html(norm_items(colc.get("bullets")), st["col_body"], theme)
            # style.left.valign / style.right.valign（deck-schema.md 記載）: "middle" で
            # heading+bullets をまとめて箱の縦中央に寄せる。flex column なら実際に描画された
            # 高さで自動的に中央寄せされるため、estimate_lines() 等の見積もりが不要。
            valign = box.get("valign", "top")
            justify = "justify-content:center;" if valign == "middle" else ""
            border = f'border:1px solid {col(theme, box["border_color"])};' if box.get("border_color") else ""
            parts.append(
                f'<div style="{pos(box)}background:{col(theme, box["fill"])};{radius}{border}'
                f'padding:{box["pad"]}px;box-sizing:border-box;overflow:hidden;'
                f'display:flex;flex-direction:column;{justify}">{inner}</div>'
            )

    elif t == "table":
        chrome()
        tr = st["table"]
        cols = slide.get("columns", [])
        widths = tr.get("col_widths")
        colgroup = ""
        if widths:
            total_w = sum(widths)
            colgroup = "<colgroup>" + "".join(
                f'<col style="width:{w / total_w * 100:.2f}%;">' for w in widths
            ) + "</colgroup>"
        cells_pad = f"padding:{tr['pad_y']}px {tr['pad_x']}px;"
        html = [
            f'<div style="{pos(tr)}"><table style="width:100%;border-collapse:collapse;'
            f'line-height:{tr["line_height"]};">{colgroup}<thead><tr>'
        ]
        for c in cols:
            html.append(
                f'<th style="background:{col(theme, tr["header_fill"])};'
                f'color:{col(theme, tr["header_color"])};font-size:{tr["header_size"]}px;'
                f'text-align:left;font-weight:600;height:{tr["header_h"]}px;'
                f'{cells_pad}box-sizing:border-box;">{esc(c)}</th>'
            )
        html.append("</tr></thead><tbody>")
        row_fills = tr.get("row_fills", {})
        for i, row in enumerate(slide.get("rows", [])):
            if str(i) in row_fills:
                bg = col(theme, row_fills[str(i)])
            else:
                bg = col(theme, tr["row_alt_fill"]) if i % 2 else col(theme, "background")
            html.append("<tr>")
            for cell in row:
                visible_text, cell_url = table_cell_parts(cell)
                cell_text = esc(visible_text)
                if cell_url:
                    cell_text = (
                        f'<a href="{esc(cell_url)}" target="_blank" rel="noopener noreferrer" '
                        'style="color:inherit;text-decoration:underline;'
                        f'text-underline-offset:2px;">{cell_text}</a>'
                    )
                html.append(
                    f'<td style="background:{bg};color:{col(theme, tr["cell_color"])};'
                    f'font-size:{tr["cell_size"]}px;height:{tr["row_h"]}px;{cells_pad}'
                    f'box-sizing:border-box;border-bottom:1px solid {col(theme, tr.get("border_color", "border"))};">{cell_text}</td>'
                )
            html.append("</tr>")
        html.append("</tbody></table></div>")
        parts.append("".join(html))

    elif t == "code":
        chrome()
        r = st["code"]
        size, _lines = fit_code_size(slide.get("code", ""), r)
        radius = f"border-radius:{r.get('radius', 0)}px;"
        parts.append(
            f'<pre style="{pos(r)}margin:0;background:{col(theme, r["fill"])};{radius}'
            f'padding:{r["pad"]}px;box-sizing:border-box;overflow:hidden;">'
            # P0: font_css() は二重引用符を含む値を返すため esc() でエスケープしてから style
            # 属性に埋める（そのまま埋めると属性値がここで打ち切られ、font-size/line-height/
            # color が丸ごと無効になり、コード本文が背景とほぼ同色で判読不能になっていた）。
            f'<code style="font-family:{esc(font_css(fonts["code"], mono=True))};font-size:{size}px;'
            f'line-height:{r["line_height"]};color:{col(theme, r["color"])};">'
            f"{esc(slide.get('code', ''))}</code></pre>"
        )

    elif t == "quote":
        parts.append(div(st["mark"], theme, "“"))
        parts.append(div(st["text"], theme, esc(slide.get("text", "")), font_family=HEADING_FONT_CSS))
        if slide.get("attribution"):
            parts.append(div(st["attribution"], theme, "— " + esc(slide["attribution"])))

    elif t == "image":
        if slide.get("eyebrow") and "eyebrow" in st:
            parts.append(div(st["eyebrow"], theme, esc(slide["eyebrow"]), font_family=HEADING_FONT_CSS))
        if slide.get("title"):
            parts.append(div(st["title"], theme, esc(slide["title"]), font_family=HEADING_FONT_CSS, heading=True))
            parts.append(rect(st["rule"], theme))
        r = st["img"]
        # X-4/M-72: image_source() を通す（HTML/PPTX 共通のリサイズ・再圧縮キャッシュ）。
        src = image_source(deck_dir, slide["path"], False, target_w=r.get("w"), target_h=r.get("h"))
        mime = mimetypes.guess_type(src.name)[0] or "image/png"
        data = base64.b64encode(src.read_bytes()).decode()
        alt = esc(slide.get("caption") or slide.get("title") or "")  # X-3
        parts.append(
            f'<div style="{pos(r)}"><img src="data:{mime};base64,{data}" alt="{alt}" '
            f'style="width:100%;height:100%;object-fit:contain;"></div>'
        )
        if slide.get("caption"):
            parts.append(div(st["caption"], theme, esc(slide["caption"])))
        parts.append(footer_html(st, deck, theme, page, total))

    elif t == "image_text":
        if slide.get("image_side") == "left":
            st = flip_regions(st, ("body", "img", "caption"))
        if slide.get("eyebrow") and "eyebrow" in st:
            parts.append(div(st["eyebrow"], theme, esc(slide["eyebrow"]), font_family=HEADING_FONT_CSS))
        parts.append(div(st["title"], theme, esc(slide.get("title", "")), font_family=HEADING_FONT_CSS, heading=True))
        parts.append(rect(st["rule"], theme))
        if slide.get("punch"):
            parts.append(div(st["punch"], theme, esc(slide["punch"]), font_family=HEADING_FONT_CSS))
        parts.append(
            f'<div style="{pos(st["body"])}">'
            + bullets_html(norm_items(slide.get("bullets")), st["body"], theme)
            + "</div>"
        )
        r = st["img"]
        src = image_source(deck_dir, slide["path"], False, target_w=r.get("w"), target_h=r.get("h"))
        mime = mimetypes.guess_type(src.name)[0] or "image/png"
        data = base64.b64encode(src.read_bytes()).decode()
        alt = esc(slide.get("caption") or slide.get("title") or "")  # X-3
        parts.append(
            f'<div style="{pos(r)}"><img src="data:{mime};base64,{data}" alt="{alt}" '
            f'style="width:100%;height:100%;object-fit:contain;"></div>'
        )
        if slide.get("caption"):
            parts.append(div(st["caption"], theme, esc(slide["caption"])))
        parts.append(footer_html(st, deck, theme, page, total))

    elif t == "closing":
        parts.append(div(st["title"], theme, esc(slide.get("title", "")), font_family=HEADING_FONT_CSS, heading=True))
        parts.append(rect(st["rule"], theme))
        if slide.get("bullets"):
            parts.append(
                f'<div style="{pos(st["body"])}">'
                + bullets_html(norm_items(slide.get("bullets")), st["body"], theme)
                + "</div>"
            )
        if slide.get("message"):
            msg = closing_message_style(st, slide)  # M-38: body 直後に自動配置（PPTX と共通の計算式）
            parts.append(div(msg, theme, esc(slide["message"]), font_family=HEADING_FONT_CSS))

    elif t == "agenda":
        chrome()
        b = st["body"]
        items = norm_items(slide.get("items", []))
        num_start = slide.get("num_start", 0)
        row_h = b.get("row_h", 44); gap = b.get("gap", 10); num_w = b.get("num_w", 52)
        col_gap = b.get("col_gap", 56); pitch = row_h + gap
        R = max(1, int(b["h"] // pitch))
        N = len(items)
        cols = 1 if N <= R else min(b.get("max_cols", 2), -(-N // R))
        rows = -(-N // cols)
        col_w = (b["w"] - (cols - 1) * col_gap) / cols
        noff = (row_h - b.get("num_size", 24)) / 2
        toff = (row_h - b.get("text_size", 20)) / 2
        for i, it in enumerate(items):
            c = i // rows; r = i % rows
            x = b["x"] + c * (col_w + col_gap); y = b["y"] + r * pitch
            active = it.get("active")
            if active:
                a_border = (f'border:1px solid {col(theme, b["active_border"])};box-sizing:border-box;'
                            if b.get("active_border") else "")
                parts.append(f'<div style="{pos({"x": x-14, "y": y-3, "w": col_w+18, "h": row_h+6})}'
                             f'background:{col(theme, b.get("active_fill","surface"))};'
                             f'border-radius:8px;{a_border}"></div>')
            parts.append(div({"x": x, "y": y+noff, "w": num_w, "h": row_h, "size": b.get("num_size", 24),
                              "color": b.get("num_color", "accent"), "bold": True, "line_height": 1.1},
                             theme, f"{num_start+i+1:02d}"))
            parts.append(div({"x": x+num_w, "y": y+toff, "w": col_w-num_w, "h": row_h, "size": b.get("text_size", 20),
                              "color": (b.get("active_color", "primary") if active else b.get("text_color", "text")),
                              "bold": bool(active), "line_height": 1.2}, theme, esc(it["text"])))

    elif t == "steps":
        chrome()
        b = st["body"]; steps = slide.get("steps", [])
        n = max(1, len(steps)); gap = b.get("gap", 20)
        card_w = (b["w"] - (n-1)*gap) / n
        header_h = b.get("header_h", 60); pad = b.get("pad", 16); radius = b.get("radius", 12)
        card_h = steps_card_height(b, steps, card_w)
        # fit="content" で縮めた分は行全体を少し上寄りの中央（40%）に置き、上下の余白を均す
        row_y = b["y"] + (b["h"] - card_h) * 0.4
        card_border = (f'border:1px solid {col(theme, b["card_border"])};box-sizing:border-box;'
                       if b.get("card_border") else "")
        for i, step in enumerate(steps):
            x = b["x"] + i * (card_w + gap)
            parts.append(f'<div style="{pos({"x": x, "y": row_y, "w": card_w, "h": card_h})}'
                         f'background:{col(theme, b.get("card_fill","surface"))};'
                         f'border-radius:{radius}px;{card_border}"></div>')
            parts.append(f'<div style="{pos({"x": x, "y": row_y, "w": card_w, "h": header_h})}'
                         f'background:{col(theme, b.get("header_fill","primary"))};'
                         f'border-radius:{radius}px {radius}px 0 0;"></div>')
            parts.append(div({"x": x+pad, "y": row_y+9, "w": card_w-2*pad, "h": 18, "size": 13,
                              "color": b.get("num_color", "accent"), "bold": True}, theme, f"STEP {i+1}"))
            parts.append(div({"x": x+pad, "y": row_y+28, "w": card_w-2*pad, "h": header_h-26,
                              "size": b.get("label_size", 18), "color": b.get("header_color", "on_primary"),
                              "bold": True, "line_height": 1.15}, theme, esc(step.get("label", "")),
                             font_family=HEADING_FONT_CSS))
            itemr = {"x": x+pad, "y": row_y+header_h+pad, "w": card_w-2*pad, "h": card_h-header_h-2*pad,
                     "size": b.get("item_size", 15), "gap": b.get("item_gap", 8), "color": b.get("item_color", "text"),
                     "line_height": 1.4, "marker_color": "accent", "child_size": 13, "child_color": "muted",
                     "indent": 16, "child_gap": 4}
            # style.body.valign（既定 top）: "middle" で items をカードの残り高さの中で縦中央寄せにする
            # （ヘッダー帯は動かさない。items だけが可変コンテンツのため flex で十分）。
            items_extra = "display:flex;flex-direction:column;justify-content:center;" if b.get("valign", "top") == "middle" else ""
            parts.append(f'<div style="{pos(itemr)}{items_extra}">'
                         + bullets_html(norm_items(step.get("items", [])), itemr, theme) + "</div>")
            if i < n-1:
                parts.append(div({"x": x+card_w, "y": row_y+card_h/2-18, "w": gap, "h": 36,
                                  "size": b.get("chevron_size", 28), "color": b.get("chevron_color", "accent"),
                                  "align": "center", "bold": True}, theme, "›"))

    elif t == "matrix":
        chrome()
        g = st["grid"]; ax = st.get("axis", {})
        gx, gy, gw, gh = g["x"], g["y"], g["w"], g["h"]
        gap = g.get("gap", 14); pad = g.get("pad", 18); radius = g.get("radius", 12)
        valign = g.get("valign", "top")  # M-37
        cw = (gw - gap) / 2; ch = (gh - gap) / 2
        quads = slide.get("quadrants", [])
        positions = [(gx, gy), (gx+cw+gap, gy), (gx, gy+ch+gap), (gx+cw+gap, gy+ch+gap)]
        g_border = (f'border:1px solid {col(theme, g["border_color"])};box-sizing:border-box;'
                    if g.get("border_color") else "")
        for qi, (qx, qy) in enumerate(positions):
            q = quads[qi] if qi < len(quads) else {}
            hi = q.get("highlight")
            fill = col(theme, g.get("hi_fill", "on_primary_soft") if hi else g.get("fill", "surface"))
            parts.append(f'<div style="{pos({"x": qx, "y": qy, "w": cw, "h": ch})}'
                         f'background:{fill};border-radius:{radius}px;{g_border}"></div>')
            # M-37: heading を含めたコンテンツ全体を象限中央に配置する（heading の y も
            # valign に応じて計算し直す。見出し固定＋本文だけセンタリングだと見出し直下に
            # 不自然な空白ができるため）。
            lay = _matrix_content_layout(q.get("body", ""), g, cw, ch, pad, valign)
            parts.append(div({"x": qx+pad, "y": qy+lay["head_y"], "w": cw-2*pad, "h": lay["head_h"],
                              "size": g.get("heading_size", 18),
                              "color": (g.get("hi_heading_color", "accent") if hi else g.get("heading_color", "primary")),
                              "bold": True, "line_height": 1.2}, theme, esc(q.get("heading", "")),
                             font_family=HEADING_FONT_CSS))
            parts.append(div({"x": qx+pad, "y": qy+lay["body_y"], "w": cw-2*pad, "h": lay["body_h"],
                              "size": g.get("body_size", 14), "color": g.get("body_color", "muted"),
                              "line_height": 1.4}, theme, esc(q.get("body", "")).replace("\n", "<br>")))
        xa = slide.get("x_axis", {}); ya = slide.get("y_axis", {})
        asize = ax.get("size", 15); acol = ax.get("color", "muted"); ncol = ax.get("name_color", "text")
        parts.append(div({"x": 72, "y": gy-30, "w": 360, "h": 22, "size": asize, "color": ncol, "bold": True},
                         theme, "▲ " + esc(ya.get("label", ""))))
        parts.append(div({"x": gx-96, "y": gy+6, "w": 88, "h": 22, "size": asize, "color": acol, "align": "right"}, theme, esc(ya.get("high", ""))))
        parts.append(div({"x": gx-96, "y": gy+gh-28, "w": 88, "h": 22, "size": asize, "color": acol, "align": "right"}, theme, esc(ya.get("low", ""))))
        parts.append(div({"x": gx, "y": gy+gh+30, "w": gw, "h": 22, "size": asize, "color": ncol, "align": "center", "bold": True},
                         theme, esc(xa.get("label", "")) + " ▶"))
        parts.append(div({"x": gx, "y": gy+gh+6, "w": 160, "h": 22, "size": asize, "color": acol}, theme, esc(xa.get("low", ""))))
        parts.append(div({"x": gx+gw-160, "y": gy+gh+6, "w": 160, "h": 22, "size": asize, "color": acol, "align": "right"}, theme, esc(xa.get("high", ""))))

    elif t == "cards":
        chrome()
        g = st["grid"]; cards = slide.get("cards", [])
        cols = max(1, slide.get("columns", g.get("cols", 3)))
        n = len(cards); rows = max(1, (n + cols - 1) // cols)
        gap = g.get("gap", 16); pad = g.get("pad", 18); radius = g.get("radius", 12)
        valign = g.get("valign", "top")  # M-37
        cw = (g["w"] - (cols-1)*gap) / cols; chh = (g["h"] - (rows-1)*gap) / rows
        g_border = (f'border:1px solid {col(theme, g["border_color"])};box-sizing:border-box;'
                    if g.get("border_color") else "")
        for idx, card in enumerate(cards):
            rr = idx // cols; cc = idx % cols
            x = g["x"] + cc * (cw + gap); y = g["y"] + rr * (chh + gap)
            parts.append(f'<div style="{pos({"x": x, "y": y, "w": cw, "h": chh})}'
                         f'background:{col(theme, g.get("fill","surface"))};border-radius:{radius}px;{g_border}"></div>')
            # M-37: heading を含めたコンテンツ全体（実際に描画される段だけ）をカード中央に配置する
            # （heading の y も valign に応じて計算し直す。見出し固定＋本文/itemsだけセンタリング
            # だと見出し直下に不自然な空白ができるため）。
            lay = _cards_content_layout(card, g, cw, chh, pad, valign)
            has_items = bool(card.get("items"))
            parts.append(div({"x": x+pad, "y": y+lay["head_y"], "w": cw-2*pad, "h": lay["head_h"],
                              "size": g.get("heading_size", 19),
                              "color": g.get("heading_color", "primary"), "bold": True, "line_height": 1.2},
                             theme, esc(card.get("heading", "")), font_family=HEADING_FONT_CSS))
            if card.get("body"):
                body_r = {"x": x+pad, "y": y+lay["body_y"], "w": cw-2*pad, "h": lay["body_h"],
                          "size": g.get("body_size", 15), "color": g.get("body_color", "text"), "line_height": 1.4}
                parts.append(div(body_r, theme, esc(card["body"]).replace("\n", "<br>")))
            if has_items:
                itemr = {"x": x+pad, "y": y+lay["items_y"], "w": cw-2*pad, "h": lay["items_h"],
                         "size": g.get("item_size", 14),
                         "gap": g.get("item_gap", 6), "color": g.get("item_color", "muted"), "line_height": 1.35,
                         "marker_color": "accent", "child_size": 12, "child_color": "muted", "indent": 14, "child_gap": 4}
                parts.append(f'<div style="{pos(itemr)}">'
                             + bullets_html(norm_items(card["items"]), itemr, theme) + "</div>")

    elif t == "swimlane":
        chrome()
        f = st["flow"]
        geo = swimlane_geometry(slide, st, page)
        def cc(key, default):
            return col(theme, f.get(key, default))
        # lane bands (subtle stripes across full width)
        for ln in geo["lanes"]:
            bg = cc("lane_band_b", "surface") if ln["alt"] else col(theme, "background")
            parts.append(f'<div style="{pos({"x": f["x"], "y": ln["y"], "w": f["w"], "h": ln["h"]})}'
                         f'background:{bg};border-top:1px solid {cc("lane_border", "border")};"></div>')
        # group boxes (Lv1, vertical text)
        for gp in geo["groups"]:
            parts.append(f'<div style="{pos(gp)}display:flex;align-items:center;justify-content:center;'
                         f'background:{cc("group_fill","primary")};box-sizing:border-box;">'
                         f'<span style="writing-mode:vertical-rl;text-orientation:mixed;'
                         f'color:{cc("group_color","on_primary")};font-weight:700;font-size:{f.get("group_size",13)}px;">'
                         f'{esc(gp["label"])}</span></div>')
        # lane labels (Lv2)
        for ln in geo["lanes"]:
            parts.append(f'<div style="{pos({"x": ln["x"], "y": ln["y"], "w": ln["w"], "h": ln["h"]})}'
                         f'display:flex;align-items:center;justify-content:center;text-align:center;'
                         f'background:{cc("lane_fill","surface")};border:1px solid {cc("lane_border","border")};'
                         f'box-sizing:border-box;color:{cc("lane_color","text")};font-size:{f.get("lane_size",13)}px;'
                         f'padding:2px;line-height:1.2;">{esc(ln["label"])}</div>')
        # phase band
        for ph in geo["phases"]:
            parts.append(f'<div style="{pos(ph)}display:flex;align-items:center;justify-content:center;'
                         f'background:{cc("phase_fill","on_primary_soft")};box-sizing:border-box;'
                         f'color:{cc("phase_color","primary")};font-weight:700;font-size:{f.get("phase_size",14)}px;'
                         f'border-left:2px solid {col(theme,"background")};">{esc(ph["label"])}</div>')
        # edges (SVG overlay). 実線=作業の流れ / 破線=システム操作。配線・ラベル位置は
        # diagram_engine（共通ルーター）が計算し、描画は新図解タイプと同じ _edges_svg_html を使う
        # （variant: emphasis / security / dashed、ラベルは背景ピル付き）。
        # swimlane の style:"dashed" は edge_dashed_color ではなく従来どおり edge_color で描く。
        ed = dict(f)
        ed.setdefault("edge_dashed_color", f.get("edge_color", "text"))
        ed.setdefault("label_fill", "background")
        parts.append(_edges_svg_html(geo["edges"], ed, theme, f"sw{page}", HEADING_FONT_CSS, layer="lines"))
        # nodes
        for nid, nd in geo["nodes"].items():
            parts.append(_sw_node_html(nd, theme, f))
        parts.append(_edges_svg_html(geo["edges"], ed, theme, f"sw{page}", HEADING_FONT_CSS, layer="labels"))

    elif t in GRID_DIAGRAM_TYPES:
        chrome()
        parts.append(grid_diagram_html(slide, st, theme, t, page, HEADING_FONT_CSS))

    elif t == "sequence":
        chrome()
        parts.append(sequence_html(slide, st, theme, page, HEADING_FONT_CSS))

    elif t == "swimlane_legend":
        if slide.get("eyebrow") and "eyebrow" in st:
            parts.append(div(st["eyebrow"], theme, esc(slide["eyebrow"]), font_family=HEADING_FONT_CSS))
        parts.append(div(st["title"], theme, esc(slide.get("title", "凡例")), font_family=HEADING_FONT_CSS, heading=True))
        parts.append(rect(st["rule"], theme))
        if slide.get("lead") and "lead" in st:
            parts.append(div(st["lead"], theme, esc(slide["lead"])))
        a = st["area"]; sym = st.get("sym", {})
        entries = slide.get("items") or SWIMLANE_LEGEND
        cols = a.get("cols", 2); rows = -(-len(entries) // cols)
        col_gap = a.get("col_gap", 60); row_h = a.get("row_h", 54); sym_w = a.get("sym_w", 120)
        col_w = (a["w"] - (cols - 1) * col_gap) / cols
        for i, (shape, vk, label, desc) in enumerate(entries):
            c = i // rows; r = i % rows
            x = a["x"] + c * (col_w + col_gap); y = a["y"] + r * row_h
            parts.append(_legend_symbol_html(shape, vk, x, y + row_h / 2, sym_w, theme, sym))
            tx = x + sym_w + 8
            parts.append(div({"x": tx, "y": y + 6, "w": col_w - sym_w - 8, "h": 22,
                              "size": a.get("label_size", 16), "color": "text", "bold": True}, theme, esc(label)))
            parts.append(div({"x": tx, "y": y + 29, "w": col_w - sym_w - 8, "h": 20,
                              "size": a.get("desc_size", 13), "color": "muted"}, theme, esc(desc)))
        parts.append(footer_html(st, deck, theme, page, total))

    return "".join(parts)


HTML_TEMPLATE = Template("""<!doctype html>
<html lang="ja">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>$TITLE</title>
<style>
* { margin: 0; padding: 0; box-sizing: border-box; }
html, body { height: 100%; background: #101014; overflow: hidden; }
#stage { position: absolute; top: 50%; left: 50%; width: ${W}px; height: ${H}px; }
.slide {
  position: absolute; inset: 0; width: ${W}px; height: ${H}px;
  font-family: $F_BODY;
  display: none; overflow: hidden;
}
.slide.active { display: block; }
.slide li::marker { color: var(--mk, inherit); }
#hud {
  position: fixed; right: 16px; bottom: 12px; color: #9a94ab;
  font: 13px/1 $F_BODY; user-select: none; z-index: 10;
}
#notes {
  position: fixed; left: 0; right: 0; bottom: 0; max-height: 32vh; overflow: auto;
  background: rgba(16, 10, 26, 0.94); color: #e6def3; padding: 16px 28px 20px;
  font: 15px/1.7 $F_BODY; z-index: 20; border-top: 2px solid $C_ACCENT;
}
#notes .label { color: $C_ACCENT; font-size: 12px; font-weight: 700; margin-bottom: 6px; }
@media print {
  /* X-2: #stage の JS 由来の transform（画面フィット用の translate/scale）と、
     body の overflow:hidden/height:100% は、印刷時に複数ページへ正しく流し込む妨げになる
     ため無効化する（.slide/#hud/#notes だけを print 用にしても、この2つが残ると
     スライドが縮小・クリップされて壊れた出力になることを Playwright の page.pdf() で確認済み）。*/
  html, body { height: auto !important; overflow: visible !important; background: #fff !important; }
  #stage { position: static !important; top: auto !important; left: auto !important; transform: none !important; width: auto !important; height: auto !important; }
  .slide { position: relative !important; display: block !important; page-break-after: always; transform: none !important; }
  #hud, #notes { display: none !important; }
  /* 主要ブラウザは印刷ダイアログの「背景のグラフィック」が既定でOFF（Playwright の
     page.pdf() も printBackground の既定値が false）で、これが無いと背景色で塗られた
     要素（section の全面背景・table のヘッダ色・cards/matrix の塗り・code 背景等、
     テーマの大半の視覚設計）が丸ごと消え、白背景に白文字になるスライドが出る。 */
  * { -webkit-print-color-adjust: exact !important; print-color-adjust: exact !important; color-adjust: exact !important; }
  @page { size: ${W}px ${H}px; margin: 0; }
}
</style>
</head>
<body>
<div id="stage">
$SLIDES
</div>
<div id="hud"><span id="pageno"></span>&nbsp;&nbsp;[&larr;/&rarr;] 移動 [N] ノート</div>
<div id="notes" hidden><div class="label">SPEAKER NOTES</div><div id="notes-body"></div></div>
<script>
var slides = Array.prototype.slice.call(document.querySelectorAll('.slide'));
var stage = document.getElementById('stage');
var pageno = document.getElementById('pageno');
var notes = document.getElementById('notes');
var notesBody = document.getElementById('notes-body');
var idx = 0;

function fit() {
  var s = Math.min(window.innerWidth / $W, window.innerHeight / $H);
  stage.style.transform = 'translate(-50%, -50%) scale(' + s + ')';
}
function show(i) {
  idx = Math.max(0, Math.min(slides.length - 1, i));
  slides.forEach(function (s, j) { s.classList.toggle('active', j === idx); });
  pageno.textContent = (idx + 1) + ' / ' + slides.length;
  location.replace('#' + (idx + 1));
  var n = slides[idx].getAttribute('data-notes') || '';
  notesBody.textContent = n || '（このスライドにノートはありません）';
}
function toggleNotes() { notes.hidden = !notes.hidden; }

window.addEventListener('resize', fit);
window.addEventListener('keydown', function (e) {
  if (e.key === 'ArrowRight' || e.key === 'PageDown' || e.key === ' ') { show(idx + 1); e.preventDefault(); }
  else if (e.key === 'ArrowLeft' || e.key === 'PageUp') { show(idx - 1); e.preventDefault(); }
  else if (e.key === 'Home') { show(0); }
  else if (e.key === 'End') { show(slides.length - 1); }
  else if (e.key === 'n' || e.key === 'N') { toggleNotes(); }
});
window.addEventListener('click', function (e) {
  if (notes.contains(e.target)) { return; }
  if (e.clientX > window.innerWidth / 2) { show(idx + 1); } else { show(idx - 1); }
});

var start = parseInt((location.hash || '#1').slice(1), 10);
fit();
show(isNaN(start) ? 0 : start - 1);
</script>
</body>
</html>
""")


def slide_bg(slide, st, theme):
    if "bg" in st:
        return col(theme, st["bg"].get("fill", "primary"))
    return col(theme, "background")


def build_html(deck, theme, layout, deck_dir: Path, out_path: Path):
    slides = deck.get("slides", [])
    total = len(slides)
    rendered = []
    for i, slide in enumerate(slides):
        st = resolve_style(layout, deck, slide)
        body = html_slide_body(slide, st, deck, theme, deck_dir, i + 1, total)
        if deck.get("meta", {}).get("brand") and "brand" in st:
            body += div(st["brand"], theme, esc(deck.get("meta", {}).get("brand", "")))
        notes_attr = f' data-notes="{esc(slide["notes"])}"' if slide.get("notes") else ""
        bg = slide_bg(slide, st, theme)
        # X-3: スクリーンリーダー向けにスライド番号とタイトルを aria-label で伝える
        title_txt = slide.get("title", "")
        aria_label = f"{i + 1} / {total}: {title_txt}" if title_txt else f"{i + 1} / {total}"
        rendered.append(
            f'<section class="slide" aria-label="{esc(aria_label)}" '
            f'style="background:{bg};color:{col(theme, "text")};"'
            f"{notes_attr}>{body}</section>"
        )
    doc = HTML_TEMPLATE.safe_substitute(
        TITLE=esc(deck.get("meta", {}).get("title", deck.get("meta", {}).get("id", ""))),
        SLIDES="\n".join(rendered),
        W=str(CANVAS_W), H=str(CANVAS_H),
        C_ACCENT=col(theme, "accent"),
        F_BODY=font_css(theme["fonts"]["body"]),
    )
    try:
        out_path.write_text(doc, encoding="utf-8")
    except PermissionError:
        fail(f"{out_path} に書き込めません。ファイル（ブラウザ等）を開いている場合は閉じてから再実行してください")


# ---------------------------------------------------------------------------
# PPTX ビルド
# ---------------------------------------------------------------------------

def _apply_ooxml_theme(prs, theme, fonts):
    """OOXML の theme1.xml（clrScheme/fontScheme）をデッキのテーマトークンで書き換える（M-49）。
    PowerPoint の「テーマの色/フォント」タブから見ても一致するようにするための付随機能であり、
    シェイプの塗り自体は従来どおり RGB 直書きのまま。失敗しても PPTX 生成は継続する。"""
    try:
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT
        from pptx.oxml import parse_xml
        from pptx.oxml.ns import qn as _qn
        from lxml import etree

        theme_part = prs.slide_masters[0].part.part_related_by(RT.THEME)
        root = parse_xml(theme_part.blob)

        def _set_solid(el, hex_color):
            if el is None:
                return
            for child in list(el):
                el.remove(child)
            el.append(el.makeelement(_qn("a:srgbClr"), {"val": hex_color}))

        clr = root.find(".//" + _qn("a:clrScheme"))
        if clr is not None:
            mapping = {
                "dk1": "text", "lt1": "background", "dk2": "primary", "lt2": "surface",
                "accent1": "accent", "accent2": "primary",
                "accent3": "highlight_fill", "accent4": "muted",
            }
            for tag, token in mapping.items():
                _set_solid(clr.find(_qn(f"a:{tag}")), col(theme, token).lstrip("#"))

        heading_name = font_pptx(fonts.get("heading")) or font_pptx(fonts["body"])
        body_name = font_pptx(fonts["body"])
        fs = root.find(".//" + _qn("a:fontScheme"))
        if fs is not None:
            for node, name in ((fs.find(_qn("a:majorFont")), heading_name),
                               (fs.find(_qn("a:minorFont")), body_name)):
                if node is None or not name:
                    continue
                for tag in ("a:latin", "a:ea"):
                    el = node.find(_qn(tag))
                    if el is not None:
                        el.set("typeface", name)

        theme_part.blob = etree.tostring(
            root, xml_declaration=True, encoding="UTF-8", standalone=True
        )
    except Exception:
        pass


def build_pptx(deck, theme, layout, deck_dir: Path, out_path: Path):
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
        from pptx.enum.text import PP_ALIGN
        from pptx.oxml.ns import qn
        from pptx.util import Emu, Inches, Pt
    except ImportError:
        sys.exit(
            "error: python-pptx が未インストールのため PPTX を生成できません。\n"
            "  初回セットアップを実行してください: /slide-deck:setup"
            "（または `pip install python-pptx`）。HTML だけなら --html を付けて再実行できます。"
        )

    fonts = theme["fonts"]
    # fonts.heading（見出し系トークン）。未指定なら body にフォールバック（M-21/M-64/M-83）。
    HEADING_FONT = fonts.get("heading") or fonts["body"]

    def IN(px):
        return Inches(px / 96.0)

    def PTS(px):
        return Pt(px * 0.75)

    def C(name):
        return RGBColor.from_string(col(theme, name).lstrip("#"))

    ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}

    prs = Presentation()
    prs.slide_width = IN(CANVAS_W)
    prs.slide_height = IN(CANVAS_H)
    sldsz = prs.element.find(qn("p:sldSz"))
    if sldsz is not None:
        sldsz.set("type", "screen16x9")
    _apply_ooxml_theme(prs, theme, fonts)
    blank = prs.slide_layouts[6]

    def set_run(run, text, size_px, color, bold=False, name=None):
        run.text = text
        resolved_name = font_pptx(name) if name is not None else font_pptx(fonts["body"])
        f = run.font
        f.name = resolved_name
        f.size = PTS(size_px)
        f.bold = bold
        f.color.rgb = color
        rPr = run._r.get_or_add_rPr()
        for tag in ("a:ea", "a:cs"):
            el = rPr.find(qn(tag))
            if el is None:
                el = rPr.makeelement(qn(tag), {})
                rPr.append(el)
            el.set("typeface", resolved_name)

    def add_box(slide, r):
        box = slide.shapes.add_textbox(IN(r["x"]), IN(r["y"]), IN(r["w"]), IN(r.get("h", 40)))
        tf = box.text_frame
        tf.word_wrap = True
        for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
            setattr(tf, m, Emu(0))
        return tf

    def para(tf, used):
        if not used[0]:
            used[0] = True
            return tf.paragraphs[0]
        return tf.add_paragraph()

    def text_region(slide, r, text, font=None):
        tf = add_box(slide, r)
        p = tf.paragraphs[0]
        if r.get("align"):
            p.alignment = ALIGN[r["align"]]
        if r.get("line_height"):
            p.line_spacing = r["line_height"]
        set_run(p.add_run(), text, r["size"], C(r.get("color", "text")),
                bold=r.get("bold", False), name=font)
        return tf

    def add_rect(slide, x, y, w, h, fill, radius=0, border=None):
        shape = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        sp = slide.shapes.add_shape(shape, IN(x), IN(y), IN(w), IN(h))
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
        if border is not None:
            # 面の視認性: surface 系の淡い塗りだけに頼らず 1px の枠線で輪郭を出す。
            sp.line.color.rgb = border
            sp.line.width = Pt(0.75)
        else:
            sp.line.fill.background()
        sp.shadow.inherit = False
        if radius:
            # M-46: 角丸半径（px）を短辺に対する比率に変換して明示指定する
            # （px と EMU は単位変換の一次比例なので px のまま比を取ってよい）。
            sp.adjustments[0] = min(0.5, radius / max(1, min(w, h)))
        return sp

    def rect_region(slide, r):
        add_rect(slide, r["x"], r["y"], r["w"], r["h"], C(r.get("fill", "accent")),
                 radius=r.get("radius", 0))

    def _set_bullet(p, char, color, marL_px, hang_px, size_pct=100):
        """段落をネイティブの PowerPoint 箇条書きにする（M-47）。
        マーカーを文字として書かず a:buChar/a:buClr/a:buSzPct/a:buFont を設定し、
        marL=インデント・indent=-ハング幅 のぶら下げインデントにする。"""
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", str(int(IN(marL_px))))
        pPr.set("indent", str(int(-IN(hang_px))))
        buClr = pPr.makeelement(qn("a:buClr"), {})
        buClr.append(buClr.makeelement(qn("a:srgbClr"), {"val": str(color)}))
        pPr.append(buClr)
        pPr.append(pPr.makeelement(qn("a:buSzPct"), {"val": str(int(size_pct * 1000))}))
        pPr.append(pPr.makeelement(qn("a:buFont"), {"typeface": font_pptx(fonts["body"])}))
        pPr.append(pPr.makeelement(qn("a:buChar"), {"char": char}))

    def add_bullets(slide, items, r):
        tf = add_box(slide, r)
        used = [False]
        marL = r["indent"]
        child_hang = 28  # marker（○）と本文の間に十分な余白を確保する
        for it in items:
            p = para(tf, used)
            p.space_after = PTS(r["gap"])
            p.line_spacing = r["line_height"]
            size = it.get("size", r["size"])
            color = C(it["color"]) if it.get("color") else C(r["color"])
            _set_bullet(p, "•", C(r["marker_color"]), marL, marL)
            set_run(p.add_run(), it["text"], size, color, bold=it.get("bold", False))
            for c in it["children"]:
                p = para(tf, used)
                p.space_after = PTS(r["child_gap"])
                p.line_spacing = max(1.0, r["line_height"] - 0.15)
                csize = c.get("size", r["child_size"])
                ccolor = C(c["color"]) if c.get("color") else C(r["child_color"])
                _set_bullet(p, "○", C(r["marker_color"]), marL + child_hang, child_hang)
                set_run(p.add_run(), c["text"], csize, ccolor)
        return tf

    def footer(slide, st, page, total):
        if "footer_l" in st:
            text_region(slide, st["footer_l"], deck.get("meta", {}).get("title", ""))
        if "footer_r" in st:
            text_region(slide, st["footer_r"], st["footer_r"].get("text") or f"{page} / {total}")

    def chrome(slide, data, st, page, total):
        if data.get("eyebrow") and "eyebrow" in st:
            text_region(slide, st["eyebrow"], data["eyebrow"], font=HEADING_FONT)
        text_region(slide, st["title"], data.get("title", ""), font=HEADING_FONT)
        rect_region(slide, st["rule"])
        if data.get("lead") and "lead" in st:
            text_region(slide, st["lead"], data["lead"])
        footer(slide, st, page, total)

    def s_title(slide, data, st, page, total):
        rect_region(slide, st["bar"])
        text_region(slide, st["title"], data.get("title", ""), font=HEADING_FONT)
        if data.get("subtitle"):
            text_region(slide, st["subtitle"], data["subtitle"])
        if data.get("meta"):
            text_region(slide, st["meta"], data["meta"])

    def s_section(slide, data, st, page, total):
        add_rect(slide, 0, 0, CANVAS_W, CANVAS_H, C(st["bg"]["fill"]))
        if data.get("number"):
            text_region(slide, st["number"], data["number"], font=HEADING_FONT)
        text_region(slide, st["title"], data.get("title", ""), font=HEADING_FONT)
        rect_region(slide, st["rule"])
        if data.get("subtitle"):
            text_region(slide, st["subtitle"], data["subtitle"])

    def s_bullets(slide, data, st, page, total):
        chrome(slide, data, st, page, total)
        tf_body = add_bullets(slide, norm_items(data.get("bullets")), st["body"])
        # style.body.valign（既定 top）: HTML 側の flex 中央寄せと同じ考え方（M-37 系と同一パターン）。
        if st["body"].get("valign", "top") == "middle":
            from pptx.enum.text import MSO_ANCHOR
            tf_body.vertical_anchor = MSO_ANCHOR.MIDDLE

    def s_two_column(slide, data, st, page, total):
        chrome(slide, data, st, page, total)
        for side in ("left", "right"):
            box = st[side]
            colc = data.get(side, {})
            add_rect(slide, box["x"], box["y"], box["w"], box["h"],
                     C(box["fill"]), radius=box.get("radius", 0),
                     border=C(box["border_color"]) if box.get("border_color") else None)
            pad = box["pad"]
            # style.left.valign / style.right.valign: "middle" で heading+bullets をまとめて
            # 箱の縦中央に寄せる（HTML と同じ考え方。PPTX は実描画高さを取れないため見積もる）。
            hd = st["col_heading"]
            valign = box.get("valign", "top")
            lay = _two_column_side_layout(colc, box, hd, st["col_body"], valign)
            if colc.get("heading"):
                hr = {"x": box["x"] + pad, "y": box["y"] + lay["head_y"], "w": box["w"] - 2 * pad,
                      "h": max(1.0, lay["head_h"]),
                      "size": hd["size"], "color": hd["color"], "bold": hd.get("bold", False)}
                text_region(slide, hr, colc["heading"], font=HEADING_FONT)
            body = dict(st["col_body"])
            body.update({"x": box["x"] + pad, "y": box["y"] + lay["body_y"], "w": box["w"] - 2 * pad,
                         "h": max(1.0, lay["body_h"])})
            add_bullets(slide, norm_items(colc.get("bullets")), body)

    def _cell_border_bottom(cell, hex_color, width_pt=1.0):
        """セルの下罫線を tr.border_color トークンで設定する（M-85）。
        既存の lnB を消してから schema 順（lnL/lnR/lnT/lnB → 塗り）を守るため先頭に挿入する。"""
        tcPr = cell._tc.get_or_add_tcPr()
        old = tcPr.find(qn("a:lnB"))
        if old is not None:
            tcPr.remove(old)
        lnB = tcPr.makeelement(qn("a:lnB"), {"w": str(int(Pt(width_pt))), "cap": "flat"})
        fill = lnB.makeelement(qn("a:solidFill"), {})
        fill.append(fill.makeelement(qn("a:srgbClr"), {"val": hex_color}))
        lnB.append(fill)
        tcPr.insert(0, lnB)

    def s_table(slide, data, st, page, total):
        from pptx.enum.text import MSO_ANCHOR
        chrome(slide, data, st, page, total)
        tr = st["table"]
        border_hex = str(C(tr.get("border_color", "border")))
        cols = data.get("columns", [])
        rows = data.get("rows", [])
        n_rows, n_cols = len(rows) + 1, len(cols)
        total_h = tr["header_h"] + tr["row_h"] * len(rows)
        gfx = slide.shapes.add_table(n_rows, n_cols, IN(tr["x"]), IN(tr["y"]),
                                     IN(tr["w"]), IN(total_h))
        table = gfx.table
        table.first_row = False
        table.horz_banding = False
        widths = tr.get("col_widths")
        if widths:
            unit = tr["w"] / sum(widths)
            for j, w in enumerate(widths):
                table.columns[j].width = IN(w * unit)
        table.rows[0].height = IN(tr["header_h"])
        for i in range(len(rows)):
            table.rows[i + 1].height = IN(tr["row_h"])
        for j, name in enumerate(cols):
            cell = table.cell(0, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = C(tr["header_fill"])
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = IN(tr["pad_x"])
            cell.margin_right = IN(tr["pad_x"])
            _cell_border_bottom(cell, border_hex, width_pt=1.5)
            p = cell.text_frame.paragraphs[0]
            set_run(p.add_run(), str(name), tr["header_size"], C(tr["header_color"]), bold=True)
        row_fills = tr.get("row_fills", {})
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                cell = table.cell(i + 1, j)
                cell.fill.solid()
                if str(i) in row_fills:
                    cell.fill.fore_color.rgb = C(row_fills[str(i)])
                else:
                    cell.fill.fore_color.rgb = C(tr["row_alt_fill"]) if i % 2 else C("background")
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.margin_left = IN(tr["pad_x"])
                cell.margin_right = IN(tr["pad_x"])
                _cell_border_bottom(cell, border_hex, width_pt=0.75)
                p = cell.text_frame.paragraphs[0]
                p.line_spacing = tr["line_height"]
                visible_text, cell_url = table_cell_parts(val)
                run = p.add_run()
                set_run(run, visible_text, tr["cell_size"], C(tr["cell_color"]))
                if cell_url:
                    run.hyperlink.address = cell_url

    def s_code(slide, data, st, page, total):
        chrome(slide, data, st, page, total)
        r = st["code"]
        size, lines = fit_code_size(data.get("code", ""), r)
        add_rect(slide, r["x"], r["y"], r["w"], r["h"], C(r["fill"]), radius=r.get("radius", 0))
        tf = add_box(slide, {"x": r["x"] + r["pad"], "y": r["y"] + r["pad"],
                             "w": r["w"] - 2 * r["pad"], "h": r["h"] - 2 * r["pad"]})
        tf.word_wrap = False  # M-24: HTML の <pre>（折り返しなし）に合わせる
        used = [False]
        for line in lines:
            p = para(tf, used)
            # M-24: spcPct（倍率）ではなく spcPts（px 起源の実寸 pt）で行送りを指定し、
            # HTML の line-height と同じ実測値に揃える。
            p.line_spacing = Pt(size * r["line_height"] * 0.75)
            set_run(p.add_run(), line if line else " ", size, C(r["color"]),
                    name=fonts["code"])

    def s_quote(slide, data, st, page, total):
        text_region(slide, st["mark"], "“")
        text_region(slide, st["text"], data.get("text", ""), font=HEADING_FONT)
        if data.get("attribution"):
            text_region(slide, st["attribution"], "— " + data["attribution"])

    def _load_pptx_picture(deck_dir_, rel_path, r, i, t):
        """image_source() 経由で PPTX 埋め込み用の画像パスを解決する（X-4: リサイズ・再圧縮
        キャッシュを HTML と共有、M-72: webp/tiff/bmp は PNG 化）。変換が必要なのに Pillow 不在等で
        変換できなかった場合は、生トレースバックの代わりに分かりやすいエラーで止める。"""
        src = image_source(deck_dir_, rel_path, True, browser, target_w=r.get("w"), target_h=r.get("h"))
        if src.suffix.lower() in PPTX_PNG_CONVERT_EXTS:
            fail(f"{i}枚目（type={t}）: 画像形式 '{src.suffix}' は PowerPoint 非対応のため PNG に変換が"
                 f"必要ですが失敗しました（Pillow を導入してください）")
        return src

    def _place_picture(slide, img_path, r, i, t):
        """画像を枠へ等比フィットして配置する（M-20: 縮小方向のクランプは外し拡大も行う）。
        枠に対して大きく拡大する場合（元解像度が枠の50%未満相当）は粗くなる旨を warn する。"""
        pic = slide.shapes.add_picture(str(img_path), IN(r["x"]), IN(r["y"]))
        ratio = min(IN(r["w"]) / pic.width, IN(r["h"]) / pic.height)
        if ratio > 2.0:
            warn(f"{i}枚目（type={t}）: 画像の解像度が表示枠に対して低く（約{ratio:.1f}倍に拡大）、"
                 f"粗く見える可能性があります。より高解像度の画像を用意してください")
        pic.width = int(pic.width * ratio)
        pic.height = int(pic.height * ratio)
        pic.left = int(IN(r["x"]) + (IN(r["w"]) - pic.width) / 2)
        pic.top = int(IN(r["y"]) + (IN(r["h"]) - pic.height) / 2)
        return pic

    def s_image(slide, data, st, page, total):
        if data.get("eyebrow") and "eyebrow" in st:
            text_region(slide, st["eyebrow"], data["eyebrow"], font=HEADING_FONT)
        if data.get("title"):
            text_region(slide, st["title"], data["title"], font=HEADING_FONT)
            rect_region(slide, st["rule"])
        r = st["img"]
        _place_picture(slide, _load_pptx_picture(deck_dir, data["path"], r, page, "image"), r, page, "image")
        if data.get("caption"):
            text_region(slide, st["caption"], data["caption"])
        footer(slide, st, page, total)

    def s_image_text(slide, data, st, page, total):
        if data.get("image_side") == "left":
            st = flip_regions(st, ("body", "img", "caption"))
        if data.get("eyebrow") and "eyebrow" in st:
            text_region(slide, st["eyebrow"], data["eyebrow"], font=HEADING_FONT)
        text_region(slide, st["title"], data.get("title", ""), font=HEADING_FONT)
        rect_region(slide, st["rule"])
        if data.get("punch"):
            text_region(slide, st["punch"], data["punch"], font=HEADING_FONT)
        add_bullets(slide, norm_items(data.get("bullets")), st["body"])
        r = st["img"]
        _place_picture(slide, _load_pptx_picture(deck_dir, data["path"], r, page, "image_text"), r, page, "image_text")
        if data.get("caption"):
            text_region(slide, st["caption"], data["caption"])
        footer(slide, st, page, total)

    def s_closing(slide, data, st, page, total):
        add_rect(slide, 0, 0, CANVAS_W, CANVAS_H, C(st["bg"]["fill"]))
        text_region(slide, st["title"], data.get("title", ""), font=HEADING_FONT)
        rect_region(slide, st["rule"])
        if data.get("bullets"):
            add_bullets(slide, norm_items(data.get("bullets")), st["body"])
        if data.get("message"):
            # M-38: body 直後に自動配置。HTML 側（closing_message_style）と共通の計算式。
            msg = closing_message_style(st, data)
            text_region(slide, msg, data["message"], font=HEADING_FONT)

    def s_agenda(slide, data, st, page, total):
        chrome(slide, data, st, page, total)
        b = st["body"]
        items = norm_items(data.get("items", []))
        num_start = data.get("num_start", 0)
        row_h = b.get("row_h", 44); gap = b.get("gap", 10); num_w = b.get("num_w", 52)
        col_gap = b.get("col_gap", 56); pitch = row_h + gap
        R = max(1, int(b["h"] // pitch))
        N = len(items)
        cols = 1 if N <= R else min(b.get("max_cols", 2), -(-N // R))
        rows = -(-N // cols)
        col_w = (b["w"] - (cols - 1) * col_gap) / cols
        noff = (row_h - b.get("num_size", 24)) / 2
        toff = (row_h - b.get("text_size", 20)) / 2
        for i, it in enumerate(items):
            c = i // rows; r = i % rows
            x = b["x"] + c * (col_w + col_gap); y = b["y"] + r * pitch
            active = it.get("active")
            if active:
                add_rect(slide, x-14, y-3, col_w+18, row_h+6, C(b.get("active_fill", "surface")), radius=8,
                         border=C(b["active_border"]) if b.get("active_border") else None)
            text_region(slide, {"x": x, "y": y+noff, "w": num_w, "h": row_h, "size": b.get("num_size", 24),
                                "color": b.get("num_color", "accent"), "bold": True, "line_height": 1.1}, f"{num_start+i+1:02d}")
            text_region(slide, {"x": x+num_w, "y": y+toff, "w": col_w-num_w, "h": row_h, "size": b.get("text_size", 20),
                                "color": (b.get("active_color", "primary") if active else b.get("text_color", "text")),
                                "bold": bool(active), "line_height": 1.2}, it["text"])

    def s_steps(slide, data, st, page, total):
        from pptx.enum.text import MSO_ANCHOR
        chrome(slide, data, st, page, total)
        b = st["body"]; steps = data.get("steps", [])
        n = max(1, len(steps)); gap = b.get("gap", 20)
        card_w = (b["w"] - (n-1)*gap) / n
        header_h = b.get("header_h", 60); pad = b.get("pad", 16); radius = b.get("radius", 12)
        card_h = steps_card_height(b, steps, card_w)
        row_y = b["y"] + (b["h"] - card_h) * 0.4
        for i, step in enumerate(steps):
            x = b["x"] + i * (card_w + gap)
            add_rect(slide, x, row_y, card_w, card_h, C(b.get("card_fill", "surface")), radius=radius,
                     border=C(b["card_border"]) if b.get("card_border") else None)
            add_rect(slide, x, row_y, card_w, header_h, C(b.get("header_fill", "primary")), radius=radius)
            text_region(slide, {"x": x+pad, "y": row_y+9, "w": card_w-2*pad, "h": 18, "size": 13,
                                "color": b.get("num_color", "accent"), "bold": True}, f"STEP {i+1}")
            text_region(slide, {"x": x+pad, "y": row_y+28, "w": card_w-2*pad, "h": header_h-26,
                                "size": b.get("label_size", 18), "color": b.get("header_color", "on_primary"),
                                "bold": True, "line_height": 1.15}, step.get("label", ""), font=HEADING_FONT)
            itemr = {"x": x+pad, "y": row_y+header_h+pad, "w": card_w-2*pad, "h": card_h-header_h-2*pad,
                     "size": b.get("item_size", 15), "gap": b.get("item_gap", 8), "color": b.get("item_color", "text"),
                     "line_height": 1.4, "marker_color": "accent", "child_size": 13, "child_color": "muted",
                     "indent": 16, "child_gap": 4}
            tf_items = add_bullets(slide, norm_items(step.get("items", [])), itemr)
            # style.body.valign（既定 top）: "middle" で items をカードの残り高さの中で縦中央寄せにする
            # （HTML と同じ考え方。ヘッダー帯は動かさない）。
            if b.get("valign", "top") == "middle":
                tf_items.vertical_anchor = MSO_ANCHOR.MIDDLE
            if i < n-1:
                text_region(slide, {"x": x+card_w, "y": row_y+card_h/2-18, "w": gap, "h": 36,
                                    "size": b.get("chevron_size", 28), "color": b.get("chevron_color", "accent"),
                                    "align": "center", "bold": True}, "›")

    def s_matrix(slide, data, st, page, total):
        from pptx.enum.text import MSO_ANCHOR
        chrome(slide, data, st, page, total)
        g = st["grid"]; ax = st.get("axis", {})
        gx, gy, gw, gh = g["x"], g["y"], g["w"], g["h"]
        gap = g.get("gap", 14); pad = g.get("pad", 18); radius = g.get("radius", 12)
        valign = g.get("valign", "top")
        cw = (gw - gap) / 2; ch = (gh - gap) / 2
        quads = data.get("quadrants", [])
        positions = [(gx, gy), (gx+cw+gap, gy), (gx, gy+ch+gap), (gx+cw+gap, gy+ch+gap)]
        for qi, (qx, qy) in enumerate(positions):
            q = quads[qi] if qi < len(quads) else {}
            hi = q.get("highlight")
            fill = C(g.get("hi_fill", "on_primary_soft")) if hi else C(g.get("fill", "surface"))
            add_rect(slide, qx, qy, cw, ch, fill, radius=radius,
                     border=C(g["border_color"]) if g.get("border_color") else None)
            # M-37: heading を含めたコンテンツ全体を象限中央に配置する（HTML と同じ計算式）。
            lay = _matrix_content_layout(q.get("body", ""), g, cw, ch, pad, valign)
            text_region(slide, {"x": qx+pad, "y": qy+lay["head_y"], "w": cw-2*pad, "h": lay["head_h"],
                                "size": g.get("heading_size", 18),
                                "color": (g.get("hi_heading_color", "accent") if hi else g.get("heading_color", "primary")),
                                "bold": True, "line_height": 1.2}, q.get("heading", ""), font=HEADING_FONT)
            tf_body = text_region(slide, {"x": qx+pad, "y": qy+lay["body_y"], "w": cw-2*pad, "h": lay["body_h"],
                                "size": g.get("body_size", 14), "color": g.get("body_color", "muted"),
                                "line_height": 1.4}, q.get("body", ""))
            if valign == "middle":
                tf_body.vertical_anchor = MSO_ANCHOR.MIDDLE
        xa = data.get("x_axis", {}); ya = data.get("y_axis", {})
        asize = ax.get("size", 15); acol = ax.get("color", "muted"); ncol = ax.get("name_color", "text")
        text_region(slide, {"x": 72, "y": gy-30, "w": 360, "h": 22, "size": asize, "color": ncol, "bold": True}, "▲ " + ya.get("label", ""))
        text_region(slide, {"x": gx-96, "y": gy+6, "w": 88, "h": 22, "size": asize, "color": acol, "align": "right"}, ya.get("high", ""))
        text_region(slide, {"x": gx-96, "y": gy+gh-28, "w": 88, "h": 22, "size": asize, "color": acol, "align": "right"}, ya.get("low", ""))
        text_region(slide, {"x": gx, "y": gy+gh+30, "w": gw, "h": 22, "size": asize, "color": ncol, "align": "center", "bold": True}, xa.get("label", "") + " ▶")
        text_region(slide, {"x": gx, "y": gy+gh+6, "w": 160, "h": 22, "size": asize, "color": acol}, xa.get("low", ""))
        text_region(slide, {"x": gx+gw-160, "y": gy+gh+6, "w": 160, "h": 22, "size": asize, "color": acol, "align": "right"}, xa.get("high", ""))

    def s_cards(slide, data, st, page, total):
        from pptx.enum.text import MSO_ANCHOR
        chrome(slide, data, st, page, total)
        g = st["grid"]; cards = data.get("cards", [])
        cols = max(1, data.get("columns", g.get("cols", 3)))
        n = len(cards); rows = max(1, (n + cols - 1) // cols)
        gap = g.get("gap", 16); pad = g.get("pad", 18); radius = g.get("radius", 12)
        valign = g.get("valign", "top")
        cw = (g["w"] - (cols-1)*gap) / cols; chh = (g["h"] - (rows-1)*gap) / rows
        for idx, card in enumerate(cards):
            rr = idx // cols; cc = idx % cols
            x = g["x"] + cc * (cw + gap); y = g["y"] + rr * (chh + gap)
            add_rect(slide, x, y, cw, chh, C(g.get("fill", "surface")), radius=radius,
                     border=C(g["border_color"]) if g.get("border_color") else None)
            # M-37: heading を含めたコンテンツ全体（実際に描画される段だけ）をカード中央に配置する
            # （HTML と同じ計算式）。
            lay = _cards_content_layout(card, g, cw, chh, pad, valign)
            has_items = bool(card.get("items"))
            text_region(slide, {"x": x+pad, "y": y+lay["head_y"], "w": cw-2*pad, "h": lay["head_h"],
                                "size": g.get("heading_size", 19),
                                "color": g.get("heading_color", "primary"), "bold": True, "line_height": 1.2},
                        card.get("heading", ""), font=HEADING_FONT)
            if card.get("body"):
                body_r = {"x": x+pad, "y": y+lay["body_y"], "w": cw-2*pad, "h": lay["body_h"],
                          "size": g.get("body_size", 15), "color": g.get("body_color", "text"), "line_height": 1.4}
                tf_body = text_region(slide, body_r, card["body"])
                if valign == "middle":
                    tf_body.vertical_anchor = MSO_ANCHOR.MIDDLE
            if has_items:
                itemr = {"x": x+pad, "y": y+lay["items_y"], "w": cw-2*pad, "h": lay["items_h"],
                         "size": g.get("item_size", 14),
                         "gap": g.get("item_gap", 6), "color": g.get("item_color", "muted"), "line_height": 1.35,
                         "marker_color": "accent", "child_size": 12, "child_color": "muted", "indent": 14, "child_gap": 4}
                tf_items = add_bullets(slide, norm_items(card["items"]), itemr)
                if valign == "middle":
                    tf_items.vertical_anchor = MSO_ANCHOR.MIDDLE

    def _shape_text(sp, text, size, color, bold=False):
        tf = sp.text_frame; tf.word_wrap = True
        tf.vertical_anchor = 3
        for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
            setattr(tf, m, Emu(0))
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.line_spacing = 1.1
        set_run(p.add_run(), text, size, color, bold=bold)

    def _sw_node_pptx(slide, nd, f):
        ns = f.get("node_size", 14)
        shape = nd["shape"]; x, y, w, h = nd["x"], nd["y"], nd["w"], nd["h"]
        if shape == "decision":
            sp = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, IN(x), IN(y), IN(w), IN(h))
            sp.fill.solid(); sp.fill.fore_color.rgb = C(f.get("decision_fill", "on_primary_soft"))
            sp.line.color.rgb = C(f.get("decision_border", "accent")); sp.line.width = Pt(1.5); sp.shadow.inherit = False
            if nd.get("text"):
                # ひし形自体の text_frame は PowerPoint が内接矩形へ極端に幅を絞るため、
                # word_wrap=True のままだと全角文字が1文字ずつ縦に折り返されてしまう
                # （旧実装はこれを避けるため word_wrap=False にしていたが、今度は3〜4文字
                # 程度のラベルで図形の外へテキストがはみ出していた）。HTML 側は箱全体の幅 w
                # で自然に複数行へ折り返しており、ここでも同じ幅の別レイヤーのテキストボックスを
                # ひし形の上に重ねることで、はみ出しと縦1文字折り返しの両方を避ける。
                tb = slide.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h))
                _shape_text(tb, nd["text"], ns - 2, C(f.get("decision_color", "primary")))
        elif shape == "terminal":
            # M-48: HTML のスタジアム型（border-radius:h/2）に合わせ、真円ではなく
            # 角丸長方形＋adjustments[0]=0.5 で描く。
            sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, IN(x), IN(y), IN(w), IN(h))
            sp.adjustments[0] = 0.5
            sp.fill.solid(); sp.fill.fore_color.rgb = C(f.get("terminal_fill", "surface"))
            sp.line.color.rgb = C(f.get("terminal_border", "muted")); sp.line.width = Pt(1.25); sp.shadow.inherit = False
            if nd.get("text"):
                _shape_text(sp, nd["text"], ns - 2, C(f.get("terminal_color", "muted")))
        elif shape == "connector":
            sp = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, IN(x), IN(y), IN(w), IN(h))
            sp.fill.solid(); sp.fill.fore_color.rgb = C(f.get("connector_fill", "surface"))
            sp.line.color.rgb = C(f.get("connector_border", "muted")); sp.line.width = Pt(1); sp.shadow.inherit = False
            if nd.get("text"):
                _shape_text(sp, nd["text"], ns - 3, C(f.get("connector_color", "text")))
        elif shape == "marker":
            kind = nd.get("kind", "mid")
            mc = {"start": f.get("marker_start", "accent"), "end": f.get("marker_end", "primary"),
                  "mid": f.get("marker_mid", "#E6B800")}.get(kind, f.get("marker_mid", "#E6B800"))
            sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, IN(x), IN(y), IN(w), IN(h))
            sp.fill.solid(); sp.fill.fore_color.rgb = C(mc)
            sp.line.color.rgb = C("background"); sp.line.width = Pt(1.25); sp.shadow.inherit = False
        elif shape == "mail":
            sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(x), IN(y), IN(w), IN(h))
            sp.fill.solid(); sp.fill.fore_color.rgb = C("background")
            sp.line.color.rgb = C(f.get("edge_color", "text")); sp.line.width = Pt(1.25); sp.shadow.inherit = False
            cxm = x + w / 2; my = y + h * 0.55
            for (x1, y1, x2, y2) in ((x, y, cxm, my), (cxm, my, x + w, y)):
                cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, IN(x1), IN(y1), IN(x2), IN(y2))
                cn.line.color.rgb = C(f.get("edge_color", "text")); cn.line.width = Pt(1.25); cn.shadow.inherit = False
        elif shape == "io":
            sp = add_rect(slide, x, y, w, h, C(f.get("io_fill", "on_primary_soft")), radius=8)
            sp.line.color.rgb = C(f.get("io_border", "accent")); sp.line.width = Pt(1)
            half = (w - 2) / 2
            for idx, (key, lbl) in enumerate((("input", "input"), ("output", "output"))):
                vals = nd.get(key) or []
                tf = add_box(slide, {"x": x + idx * half + 6, "y": y + 6, "w": half - 10, "h": h - 12})
                p0 = tf.paragraphs[0]
                set_run(p0.add_run(), lbl, f.get("io_size", 11), C(f.get("io_head_color", "accent")), bold=True)
                for v in vals:
                    pp = tf.add_paragraph()
                    set_run(pp.add_run(), v, f.get("io_size", 11), C(f.get("io_color", "text")))
        elif shape == "system":
            sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(x), IN(y), IN(w), IN(h))
            sp.fill.solid(); sp.fill.fore_color.rgb = C(f.get("system_fill", "surface"))
            sp.line.color.rgb = C(f.get("system_border", "muted")); sp.line.width = Pt(1.5); sp.shadow.inherit = False
            if nd.get("text"):
                _shape_text(sp, nd["text"], ns, C(f.get("system_color", "text")))
        else:  # task (+variant)
            variant = nd.get("variant") or "onother"
            vmap = {"onpf": (f.get("task_onpf_fill", "on_primary_soft"), f.get("task_border", "accent"), f.get("task_onpf_color", "primary")),
                    "onother": (f.get("task_onother_fill", "background"), f.get("task_border", "accent"), f.get("task_color", "text")),
                    "offline": (f.get("task_offline_fill", "surface"), f.get("task_offline_border", "muted"), f.get("task_color", "text"))}
            fl, bd, cl = vmap.get(variant, vmap["onother"])
            sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, IN(x), IN(y), IN(w), IN(h))
            sp.fill.solid(); sp.fill.fore_color.rgb = C(fl)
            sp.line.color.rgb = C(bd); sp.line.width = Pt(2); sp.shadow.inherit = False
            if nd.get("text"):
                _shape_text(sp, nd["text"], ns, C(cl))
            if nd.get("loop"):
                # 反復記号: 右上角に丸バッジ＋円弧矢印（HTML と同じ位置・見た目）
                bx, by = x + w - 9, y - 9
                badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, IN(bx), IN(by), IN(18), IN(18))
                badge.fill.solid(); badge.fill.fore_color.rgb = C("background")
                badge.line.color.rgb = C("accent"); badge.line.width = Pt(1.1); badge.shadow.inherit = False
                arr = slide.shapes.add_shape(MSO_SHAPE.CIRCULAR_ARROW, IN(bx + 3.5), IN(by + 3.5), IN(11), IN(11))
                arr.fill.solid(); arr.fill.fore_color.rgb = C("accent"); arr.line.fill.background(); arr.shadow.inherit = False

    def s_swimlane(slide, data, st, page, total):
        chrome(slide, data, st, page, total)
        f = st["flow"]
        geo = swimlane_geometry(data, st, page)
        start_idx = len(slide.shapes)  # M-50: この位置から後で作る図形だけを1グループにまとめる
        # lane bands (alt stripes)
        for ln in geo["lanes"]:
            if ln["alt"]:
                add_rect(slide, f["x"], ln["y"], f["w"], ln["h"], C(f.get("lane_band_b", "surface")))
        # group boxes (Lv1): 矩形自体は回転させず、bodyPr の vert=vert270 でテキストだけ縦書きにする
        # （M-19/M-45/M-89: sp.rotation を使うと矩形の外接位置がレーンからずれてしまうため）。
        for gp in geo["groups"]:
            sp = add_rect(slide, gp["x"], gp["y"], gp["w"], gp["h"], C(f.get("group_fill", "primary")))
            _shape_text(sp, gp["label"], f.get("group_size", 13), C(f.get("group_color", "on_primary")), bold=True)
            sp.text_frame.word_wrap = False
            sp.text_frame._txBody.bodyPr.set("vert", "vert270")
        # lane labels (Lv2)
        for ln in geo["lanes"]:
            sp = add_rect(slide, ln["x"], ln["y"], ln["w"], ln["h"], C(f.get("lane_fill", "surface")))
            sp.line.color.rgb = C(f.get("lane_border", "border")); sp.line.width = Pt(0.5)
            _shape_text(sp, ln["label"], f.get("lane_size", 13), C(f.get("lane_color", "text")))
        # phase band
        for ph in geo["phases"]:
            sp = add_rect(slide, ph["x"], ph["y"], ph["w"], ph["h"], C(f.get("phase_fill", "on_primary_soft")))
            _shape_text(sp, ph["label"], f.get("phase_size", 14), C(f.get("phase_color", "primary")), bold=True)
        # edges: 共通ルーターの経路を連結コネクタで描き、ラベルは背景ピル付きでノードの上に重ねる
        # （HTML 側と同じ _edge_visual。swimlane の style:"dashed" は従来どおり edge_color で描く）。
        ed = dict(f)
        ed.setdefault("edge_dashed_color", f.get("edge_color", "text"))
        ed.setdefault("label_fill", "background")
        for e in geo["edges"]:
            vis = _edge_visual(e, ed)
            arrow = e.get("arrow", "end")
            _poly_pptx(slide, e["points"], C(vis["color"]), vis["w"], vis["pptx_dash"],
                       arrow_end=arrow in ("end", "both"), arrow_start=(arrow == "both"), open_arrow=vis["open"])
        # nodes
        for nid, nd in geo["nodes"].items():
            _sw_node_pptx(slide, nd, f)
        for e in geo["edges"]:
            _label_pill_pptx(slide, e, ed, f.get("label_size", 13))
        # M-50: 図全体（タイトル等のチロームは除く）を1つの p:grpSp にまとめ、
        # PowerPoint 上でまとめて移動・拡大縮小できるようにする。
        group_shapes = list(slide.shapes)[start_idx:]
        if len(group_shapes) > 1:
            slide.shapes.add_group_shape(group_shapes)

    def s_swimlane_legend(slide, data, st, page, total):
        if data.get("eyebrow") and "eyebrow" in st:
            text_region(slide, st["eyebrow"], data["eyebrow"], font=HEADING_FONT)
        text_region(slide, st["title"], data.get("title", "凡例"), font=HEADING_FONT)
        rect_region(slide, st["rule"])
        if data.get("lead") and "lead" in st:
            text_region(slide, st["lead"], data["lead"])
        a = st["area"]; sym = st.get("sym", {})
        legend_start_idx = len(slide.shapes)
        entries = data.get("items") or SWIMLANE_LEGEND
        cols = a.get("cols", 2); rows = -(-len(entries) // cols)
        col_gap = a.get("col_gap", 60); row_h = a.get("row_h", 54); sym_w = a.get("sym_w", 120)
        col_w = (a["w"] - (cols - 1) * col_gap) / cols
        for i, (shape, vk, label, desc) in enumerate(entries):
            c = i // rows; r = i % rows
            x = a["x"] + c * (col_w + col_gap); y = a["y"] + r * row_h; cy = y + row_h / 2
            sw, sh = _legend_sym_size(shape)
            scx = x + sym_w / 2
            if shape == "flow":
                x1, x2 = scx - sw / 2, scx + sw / 2
                conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, IN(x1), IN(cy), IN(x2), IN(cy))
                conn.line.color.rgb = C(sym.get("edge_color", "text")); conn.line.width = Pt(2); conn.shadow.inherit = False
                lnx = conn.line._get_or_add_ln()
                if vk == "dashed":
                    lnx.append(lnx.makeelement(qn("a:prstDash"), {"val": "dash"}))
                lnx.append(lnx.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
            else:
                nd = {"shape": shape, "variant": (vk if shape == "task" else None),
                      "kind": (vk if shape == "marker" else None), "text": "",
                      "x": scx - sw / 2, "y": cy - sh / 2, "w": sw, "h": sh,
                      "input": (["…"] if shape == "io" else None), "output": (["…"] if shape == "io" else None), "loop": None}
                _sw_node_pptx(slide, nd, sym)
            tx = x + sym_w + 8
            text_region(slide, {"x": tx, "y": y + 6, "w": col_w - sym_w - 8, "h": 22,
                                "size": a.get("label_size", 16), "color": "text", "bold": True}, label)
            text_region(slide, {"x": tx, "y": y + 29, "w": col_w - sym_w - 8, "h": 20,
                                "size": a.get("desc_size", 13), "color": "muted"}, desc)
        legend_shapes = list(slide.shapes)[legend_start_idx:]
        if len(legend_shapes) > 1:
            slide.shapes.add_group_shape(legend_shapes)
        footer(slide, st, page, total)

    # -----------------------------------------------------------------------
    # ネイティブ図解タイプ（architecture / dataflow / lifecycle / sequence）の PPTX 描画。
    # ジオメトリは diagram_engine、見た目の解決は _node_visual / _edge_visual（HTML と共通）。
    # -----------------------------------------------------------------------

    def _set_dash(shape_or_conn, pptx_dash):
        if not pptx_dash:
            return
        lnx = shape_or_conn.line._get_or_add_ln()
        lnx.append(lnx.makeelement(qn("a:prstDash"), {"val": pptx_dash}))

    def _poly_pptx(slide, pts, color, width_px, pptx_dash=None, arrow_end=True, arrow_start=False, open_arrow=False):
        """直交折れ線を連結コネクタで描く。最後のセグメントに矢印（arrow_start で先頭にも）。"""
        made = []
        for k in range(len(pts) - 1):
            (x1, y1), (x2, y2) = pts[k], pts[k + 1]
            if abs(x1 - x2) < 0.05 and abs(y1 - y2) < 0.05:
                continue
            conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, IN(x1), IN(y1), IN(x2), IN(y2))
            conn.line.color.rgb = color
            conn.line.width = Pt(width_px * 0.75)
            conn.shadow.inherit = False
            lnx = conn.line._get_or_add_ln()
            if pptx_dash:
                lnx.append(lnx.makeelement(qn("a:prstDash"), {"val": pptx_dash}))
            head_type = "arrow" if open_arrow else "triangle"
            if k == len(pts) - 2 and arrow_end:
                lnx.append(lnx.makeelement(qn("a:tailEnd"), {"type": head_type, "w": "med", "len": "med"}))
            if k == 0 and arrow_start:
                lnx.append(lnx.makeelement(qn("a:headEnd"), {"type": head_type, "w": "med", "len": "med"}))
            made.append(conn)
        return made

    def _outline(sp, color, width_px, pptx_dash=None):
        sp.fill.background()
        sp.line.color.rgb = color
        sp.line.width = Pt(width_px * 0.75)
        sp.shadow.inherit = False
        _set_dash(sp, pptx_dash)
        return sp

    def _icon_pptx(slide, kind, x, y, s, color):
        """type アイコンをプリセット図形の組合せで近似する（HTML の 16×16 線画に対応）。"""
        lw = 1.4
        if kind == "database":
            _outline(slide.shapes.add_shape(MSO_SHAPE.CAN, IN(x + s * 0.12), IN(y), IN(s * 0.76), IN(s)), color, lw)
        elif kind == "cloud":
            _outline(slide.shapes.add_shape(MSO_SHAPE.CLOUD, IN(x), IN(y + s * 0.14), IN(s), IN(s * 0.72)), color, lw)
        elif kind == "security":
            _outline(slide.shapes.add_shape(MSO_SHAPE.FLOWCHART_OFFPAGE_CONNECTOR, IN(x + s * 0.1), IN(y), IN(s * 0.8), IN(s)), color, lw)
        elif kind == "frontend":
            _outline(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(x), IN(y + s * 0.12), IN(s), IN(s * 0.76)), color, lw)
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(x), IN(y + s * 0.12), IN(s), IN(s * 0.2))
            bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background(); bar.shadow.inherit = False
        elif kind == "backend":
            _outline(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(x), IN(y + s * 0.1), IN(s), IN(s * 0.32)), color, lw)
            _outline(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(x), IN(y + s * 0.58), IN(s), IN(s * 0.32)), color, lw)
        elif kind == "messagebus":
            for fy, fw in ((0.25, 1.0), (0.5, 0.72), (0.75, 1.0)):
                _poly_pptx(slide, [(x, y + s * fy), (x + s * fw, y + s * fy)], color, lw, arrow_end=(fy == 0.5), open_arrow=True)
        elif kind == "external":
            # 外部リンク（箱＋右上へ抜ける矢印）。箱は右上を欠いた形の代わりに全周矩形で近似する
            _outline(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(x + s * 0.15), IN(y + s * 0.25), IN(s * 0.6), IN(s * 0.6)), color, lw)
            _poly_pptx(slide, [(x + s * 0.47, y + s * 0.53), (x + s * 0.9, y + s * 0.1)], color, lw, arrow_end=True, open_arrow=True)

    def _glyph_pptx(slide, kind, x, y, s, color):
        """lifecycle kind の小グリフ（右上）。"""
        lw = 1.6
        if kind == "success":
            _poly_pptx(slide, [(x + s * 0.18, y + s * 0.55), (x + s * 0.42, y + s * 0.78)], color, lw, arrow_end=False)
            _poly_pptx(slide, [(x + s * 0.42, y + s * 0.78), (x + s * 0.85, y + s * 0.28)], color, lw, arrow_end=False)
        elif kind == "failure":
            _poly_pptx(slide, [(x + s * 0.22, y + s * 0.22), (x + s * 0.78, y + s * 0.78)], color, lw, arrow_end=False)
            _poly_pptx(slide, [(x + s * 0.78, y + s * 0.22), (x + s * 0.22, y + s * 0.78)], color, lw, arrow_end=False)
        elif kind == "waiting":
            _outline(slide.shapes.add_shape(MSO_SHAPE.FLOWCHART_COLLATE, IN(x + s * 0.2), IN(y + s * 0.08), IN(s * 0.6), IN(s * 0.84)), color, 1.2)
        elif kind == "start":
            sp = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, IN(x + s * 0.28), IN(y + s * 0.2), IN(s * 0.5), IN(s * 0.6))
            sp.rotation = 90
            sp.fill.solid(); sp.fill.fore_color.rgb = color; sp.line.fill.background(); sp.shadow.inherit = False

    def _diag_text_pptx(slide, r, lines, size, color, sub_lines=None, sub_size=11, sub_color=None, font=None):
        """箱の中央に label（太字・複数行）＋ sublabel（小）を描く（上下中央・左右中央）。"""
        from pptx.enum.text import MSO_ANCHOR
        tf = add_box(slide, r)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        used = [False]
        for ln in (lines or [""]):
            p = para(tf, used)
            p.alignment = PP_ALIGN.CENTER
            p.line_spacing = 1.1
            set_run(p.add_run(), ln, size, color, bold=True, name=font)
        for ln in (sub_lines or []):
            p = para(tf, used)
            p.alignment = PP_ALIGN.CENTER
            p.line_spacing = 1.1
            p.space_before = PTS(2)
            set_run(p.add_run(), ln, sub_size, sub_color if sub_color is not None else color)
        return tf

    def _diag_node_pptx(slide, b, fit, d, kind):
        vis = _node_visual(b, d, kind)
        x, y, w, h = b["x"], b["y"], b["w"], b["h"]
        pad = d.get("node_pad", 8)
        radius = d.get("node_radius", 10)
        fill = C(vis["fill"]); border = C(vis["border"])
        dash = "dash" if vis["dash"] else None
        icon_w = 0.0
        if vis["shape"] == "diamond":
            sp = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, IN(x), IN(y), IN(w), IN(h))
            sp.fill.solid(); sp.fill.fore_color.rgb = fill
            sp.line.color.rgb = border; sp.line.width = Pt(vis["border_w"] * 0.75); sp.shadow.inherit = False
            _set_dash(sp, dash)
            tx, tw = x + w * 0.12, w * 0.76
        else:
            sp = add_rect(slide, x, y, w, h, fill, radius=radius, border=border)
            sp.line.width = Pt(vis["border_w"] * 0.75)
            _set_dash(sp, dash)
            if vis["shape"] == "pill":
                sp.adjustments[0] = 0.5
            icon_size = d.get("icon_size", 16)
            ntype = b.get("type") or "generic"
            if kind != "lifecycle" and ntype in ICON_PATHS:
                icon_tok = d.get("muted_color", "muted") if (b.get("variant") == "muted") else d.get("icon_color", "accent")
                _icon_pptx(slide, ntype, x + pad, y + h / 2 - icon_size / 2, icon_size, C(icon_tok))
                icon_w = icon_size + 6
            tx = x + pad + icon_w
            tw = max(10.0, w - 2 * pad - icon_w)
        size = fit.get("size", d.get("node_size", 14)) if fit else d.get("node_size", 14)
        lines = (fit.get("lines") if fit else None) or [b.get("label", "")]
        sub_lines = (fit.get("sub_lines") if fit else None) or ([b["sublabel"]] if b.get("sublabel") else [])
        _diag_text_pptx(slide, {"x": tx, "y": y + 2, "w": tw, "h": h - 4}, lines, size, C(vis["color"]),
                        sub_lines, d.get("sub_size", 11), C(vis["sub_color"]), font=HEADING_FONT)
        if b.get("tag"):
            ts = d.get("tag_size", 10)
            tw_ = de.text_width(b["tag"], ts) + 12
            pill = add_rect(slide, x + w - tw_ - 8, y - ts * 0.8, tw_, ts * 1.6, C(d.get("tag_fill", "highlight_fill")), radius=ts)
            pill.adjustments[0] = 0.5
            _shape_text(pill, b["tag"], ts, C(d.get("tag_color", "text")), bold=True)
        if kind == "lifecycle":
            if b.get("step"):
                ss = d.get("step_size", 11)
                sx = (x + w / 2 - 20) if vis["shape"] == "diamond" else x + 8
                text_region(slide, {"x": sx, "y": y + 3, "w": 40, "h": ss * 1.4, "size": ss,
                                    "color": d.get("step_color", "accent"), "bold": True}, b["step"])
            k = b.get("kind") or "active"
            if k in KIND_GLYPHS and vis["shape"] == "rect":
                gc = C(vis["color"]) if k == "success" else C(d.get("icon_color", "accent"))
                _glyph_pptx(slide, k, x + w - 21, y + 5, 13, gc)

    def _label_pill_pptx(slide, e, d, lsize):
        lb = e.get("label_box")
        if not lb or not e.get("label"):
            return
        pill = add_rect(slide, lb["x"], lb["y"], lb["w"], lb["h"], C(d.get("label_fill", "background")), radius=4)
        classification = e.get("classification")
        if classification:
            cs = d.get("class_size", 11)
            _diag_text_pptx(slide, {"x": lb["x"] - 20, "y": lb["y"], "w": lb["w"] + 40, "h": lb["h"]},
                            [e["label"]], lsize, C(d.get("label_color", "accent")), [classification], cs,
                            C(d.get("class_color", "muted")))
        else:
            _shape_text(pill, e["label"], lsize, C(d.get("label_color", "accent")), bold=True)
            pill.text_frame.word_wrap = False

    def _edges_pptx(slide, edges, d):
        lsize = d.get("label_size", 13)
        for e in edges:
            vis = _edge_visual(e, d)
            arrow = e.get("arrow", "end")
            _poly_pptx(slide, e["points"], C(vis["color"]), vis["w"], vis["pptx_dash"],
                       arrow_end=arrow in ("end", "both"), arrow_start=(arrow == "both"), open_arrow=vis["open"])
        for e in edges:
            _label_pill_pptx(slide, e, d, lsize)

    def _group_pptx(slide, g, d):
        kind = g.get("kind") or "generic"
        radius = d.get("group_radius", 10)
        gc = d.get("group_color", "muted")
        if kind == "zone":
            sp = add_rect(slide, g["x"], g["y"], g["w"], g["h"], C(d.get("zone_fill", "surface")), radius=radius,
                          border=C(d.get("group_border", "muted")))
        else:
            sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, IN(g["x"]), IN(g["y"]), IN(g["w"]), IN(g["h"]))
            sp.adjustments[0] = min(0.5, radius / max(1, min(g["w"], g["h"])))
            bc = d.get("security_group_border", "primary") if kind == "security" else d.get("group_border", "muted")
            _outline(sp, C(bc), 1.5, "dash")
            if kind == "security":
                gc = bc
        if g.get("label"):
            gs = d.get("group_size", 12)
            text_region(slide, {"x": g["x"] + 10, "y": g["y"] + 3, "w": max(40, g["w"] - 20), "h": gs * 1.5,
                                "size": gs, "color": gc, "bold": True}, g["label"])

    def _legend_pptx(slide, items, y, x, d):
        ls = d.get("legend_size", 12)
        cx = x
        for it in items:
            _icon_pptx(slide, it["type"], cx, y + (d.get("legend_h", 26) - 14) / 2, 14, C(d.get("icon_color", "accent")))
            lw = de.text_width(it["label"], ls) + 6
            text_region(slide, {"x": cx + 19, "y": y + 4, "w": lw + 8, "h": ls * 1.5, "size": ls,
                                "color": d.get("legend_color", "muted")}, it["label"])
            cx += 19 + lw + 22

    def _grid_diagram_pptx(slide, data, st, kind):
        d = st["diagram"]
        geo = de.layout_grid_diagram(data, kind, d, [])
        start_idx = len(slide.shapes)
        for rh in geo["row_headers"]:
            if rh["alt"]:
                add_rect(slide, rh["band_x"], rh["y"], rh["band_w"], rh["h"], C(d.get("row_band_b", "surface")))
            sp = add_rect(slide, rh["x"], rh["y"], rh["w"], rh["h"], C(d.get("row_fill", "surface")))
            sp.line.color.rgb = C(d.get("row_border", "border")); sp.line.width = Pt(0.5)
            _shape_text(sp, rh["label"], d.get("row_size", 13), C(d.get("row_color", "text")))
        for rh in geo["row_headers"]:
            _poly_pptx(slide, [(rh["band_x"], rh["y"]), (rh["band_x"] + rh["band_w"], rh["y"])],
                       C(d.get("row_border", "border")), 1, arrow_end=False)
        if geo["row_headers"]:
            last = geo["row_headers"][-1]
            _poly_pptx(slide, [(last["band_x"], last["y"] + last["h"]), (last["band_x"] + last["band_w"], last["y"] + last["h"])],
                       C(d.get("row_border", "border")), 1, arrow_end=False)
        for ch in geo["col_headers"]:
            sp = add_rect(slide, ch["x"] + 1, ch["y"], ch["w"] - 2, ch["h"], C(d.get("header_fill", "highlight_fill")))
            _shape_text(sp, ch["label"], d.get("header_size", 14), C(d.get("header_color", "text")), bold=True)
        for g in geo["groups"]:
            _group_pptx(slide, g, d)
        lsize = d.get("label_size", 13)
        for e in geo["edges"]:
            vis = _edge_visual(e, d)
            arrow = e.get("arrow", "end")
            _poly_pptx(slide, e["points"], C(vis["color"]), vis["w"], vis["pptx_dash"],
                       arrow_end=arrow in ("end", "both"), arrow_start=(arrow == "both"), open_arrow=vis["open"])
        for nid, b in geo["nodes"].items():
            _diag_node_pptx(slide, b, geo["text"].get(nid, {}), d, kind)
        for e in geo["edges"]:
            _label_pill_pptx(slide, e, d, lsize)
        if geo["legend"]:
            _legend_pptx(slide, geo["legend"], geo["legend_y"], d["x"] + 4, d)
        group_shapes = list(slide.shapes)[start_idx:]
        if len(group_shapes) > 1:
            slide.shapes.add_group_shape(group_shapes)

    def s_architecture(slide, data, st, page, total):
        chrome(slide, data, st, page, total)
        _grid_diagram_pptx(slide, data, st, "architecture")

    def s_dataflow(slide, data, st, page, total):
        chrome(slide, data, st, page, total)
        _grid_diagram_pptx(slide, data, st, "dataflow")

    def s_lifecycle(slide, data, st, page, total):
        chrome(slide, data, st, page, total)
        _grid_diagram_pptx(slide, data, st, "lifecycle")

    def s_sequence(slide, data, st, page, total):
        chrome(slide, data, st, page, total)
        d = st["diagram"]
        geo = de.layout_sequence(data, d, [])
        start_idx = len(slide.shapes)
        for sg in geo["segments"]:
            sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, IN(sg["x"]), IN(sg["y"]), IN(sg["w"]), IN(sg["h"]))
            sp.adjustments[0] = min(0.5, 6 / max(1, min(sg["w"], sg["h"])))
            _outline(sp, C(d.get("segment_border", "muted")), 1, "dash")
            if sg.get("label"):
                ss = d.get("segment_size", 12)
                lw = de.text_width(sg["label"], ss) + 12
                bg = add_rect(slide, sg["x"] + 8, sg["y"] - ss * 0.75, lw, ss * 1.5, C("background"))
                _shape_text(bg, sg["label"], ss, C(d.get("segment_color", "muted")), bold=True)
                bg.text_frame.word_wrap = False
        for pid in geo["order"]:
            pb = geo["participants"].get(pid)
            if not pb:
                continue
            _poly_pptx(slide, [(pb["cx"], geo["lifeline_top"]), (pb["cx"], geo["lifeline_bottom"])],
                       C(d.get("lifeline_color", "muted")), d.get("lifeline_w", 1.5), "dash", arrow_end=False)
        for a in geo["activations"]:
            sp = add_rect(slide, a["x"], a["y"], a["w"], a["h"], C(d.get("activation_fill", "highlight_fill")), radius=3,
                          border=C(d.get("activation_border", "accent")))
            sp.line.width = Pt(1.1)
        _edges_pptx(slide, geo["messages"], d)
        fit = de.check_node_text(geo["participants"], d.get("node_size", 14), d.get("sub_size", 11), d.get("node_pad", 8), [],
                                 icon_w=d.get("icon_size", 16) + 6)
        for pid, pb in geo["participants"].items():
            _diag_node_pptx(slide, dict(pb, shape="rect"), fit.get(pid, {}), d, "sequence")
        group_shapes = list(slide.shapes)[start_idx:]
        if len(group_shapes) > 1:
            slide.shapes.add_group_shape(group_shapes)

    builders = {
        "title": s_title, "section": s_section, "bullets": s_bullets,
        "two_column": s_two_column, "table": s_table, "code": s_code,
        "quote": s_quote, "image": s_image, "image_text": s_image_text,
        "closing": s_closing,
        "agenda": s_agenda, "steps": s_steps, "matrix": s_matrix, "cards": s_cards,
        "swimlane": s_swimlane, "swimlane_legend": s_swimlane_legend,
        "architecture": s_architecture, "dataflow": s_dataflow, "lifecycle": s_lifecycle,
        "sequence": s_sequence,
    }

    slides = deck.get("slides", [])
    total = len(slides)

    # M-25: デッキ内に SVG 画像があるときだけ Playwright/Chromium を1回起動し、
    # 全スライド分の svg_to_png() で使い回す（1点ごとの再起動を防ぐ）。
    needs_svg = any(
        d.get("type") in ("image", "image_text")
        and str(d.get("path", "")).lower().endswith(".svg")
        for d in slides
    )
    browser = None
    _pw_cm = None
    if needs_svg:
        from playwright.sync_api import sync_playwright
        _pw_cm = sync_playwright()
        pw = _pw_cm.__enter__()
        browser = pw.chromium.launch()

    try:
        for i, data in enumerate(slides):
            st = resolve_style(layout, deck, data)
            slide = prs.slides.add_slide(blank)
            if "bg" not in st:
                add_rect(slide, 0, 0, CANVAS_W, CANVAS_H, C("background"))
            builders[data["type"]](slide, data, st, i + 1, total)
            if deck.get("meta", {}).get("brand") and "brand" in st:
                text_region(slide, st["brand"], deck.get("meta", {}).get("brand", ""))
            # 注意: notes は意図的に PPTX へ出力しない（Keynote 互換性問題のため）
    finally:
        if browser is not None:
            browser.close()
        if _pw_cm is not None:
            _pw_cm.__exit__(None, None, None)

    try:
        prs.save(str(out_path))
    except PermissionError:
        fail(f"{out_path} に書き込めません。PowerPoint 等で開いている場合は閉じてから再実行してください")


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main():
    setup_console()
    ap = argparse.ArgumentParser(description="deck.json → HTML / PPTX ビルダー v2")
    ap.add_argument("deck_dir", help="デッキのディレクトリ（deck.json を含む）")
    ap.add_argument("--html", action="store_true", help="HTML のみ生成")
    ap.add_argument("--pptx", action="store_true", help="PPTX のみ生成")
    args = ap.parse_args()

    deck_dir = Path(args.deck_dir).resolve()
    if not deck_dir.is_dir():
        fail(f"{deck_dir} はディレクトリではありません")
    deck, theme, layout = load_deck(deck_dir)
    validate_deck(deck, layout, theme, deck_dir)
    deck = expand_slides(deck, layout)
    deck_id = deck.get("meta", {}).get("id") or deck_dir.name

    out_dir = deck_dir / "build"
    out_dir.mkdir(exist_ok=True)

    if args.html or not args.pptx:
        out = out_dir / f"{deck_id}.html"
        build_html(deck, theme, layout, deck_dir, out)
        print(f"HTML: {out}")
    if args.pptx or not args.html:
        out = out_dir / f"{deck_id}.pptx"
        build_pptx(deck, theme, layout, deck_dir, out)
        print(f"PPTX: {out}")


if __name__ == "__main__":
    main()
