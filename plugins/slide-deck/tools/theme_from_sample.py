#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""サンプルスライド（PPTX / HTML / PDF / 画像）から配色・フォントを抽出してテーマ JSON を生成する。

使い方:
    python tools/theme_from_sample.py <file> --name <name> [--dir DIR] [--pages N] [--colors N]

対応形式と抽出方法:
    .pptx/.pptm  埋め込みテーマ（clrScheme/fontScheme）を直接読む（最良品質・追加依存なし）
    .png/.jpg…   主要色を量子化抽出（要 pillow）
    .pdf         先頭ページを描画して主要色を抽出＋埋め込みフォント名（要 pymupdf）
    .html/.htm   宣言された色（#hex / rgb()）と font-family を静的抽出

抽出は自動推定なので、生成後は必ずレポートを見て `colors`/`fonts` を微調整すること。
出力先の優先順: --dir → 環境変数 SLIDE_DECK_THEMES の先頭 → カレントディレクトリの ./themes/
（同梱の templates/themes/ には --dir で明示したときだけ書く）
トークンの意味は references/themes.md を参照。
"""
from __future__ import annotations

import argparse
import json
import os
import re
import sys
import zipfile
import xml.etree.ElementTree as ET
from collections import Counter
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
TOOLS_DIR = Path(__file__).resolve().parent
BUNDLED_THEMES = ROOT / "templates" / "themes"
NAME_RE = re.compile(r"^[a-z0-9]+(?:-[a-z0-9]+)*$")

sys.path.insert(0, str(TOOLS_DIR))
try:
    import build_deck  # type: ignore
except Exception:
    build_deck = None


def _fallback_setup_console():
    for stream in (sys.stdout, sys.stderr):
        try:
            stream.reconfigure(encoding="utf-8", errors="replace")
        except Exception:
            pass


def _fallback_fail(msg):
    print(f"error: {msg}", file=sys.stderr)
    sys.exit(1)


def _fallback_warn(msg):
    print(f"warning: {msg}", file=sys.stderr)


def _fallback_note(msg):
    print(f"note: {msg}", file=sys.stderr)


setup_console = getattr(build_deck, "setup_console", None) or _fallback_setup_console
fail = getattr(build_deck, "fail", None) or _fallback_fail
warn = getattr(build_deck, "warn", None) or _fallback_warn
note = getattr(build_deck, "note", None) or _fallback_note

# 文字色候補の彩度ガード: 輝度だけで「暗い色」を選ぶと、彩度の高いブランド色
# （例: 純青紫は緑チャンネル重みの小さいsRGB輝度式では暗く見える）を本文の文字色に
# 誤選定しがちなので、まず低彩度から選び、無ければ最も暗い色にフォールバックする。
TEXT_SAT_MAX = 0.30
# primary/accent候補の上位2件の重み比がこの値未満なら「僅差につき要確認」警告を出す。
AMBIGUOUS_RATIO = 1.5

# ---------------------------------------------------------------------------
# 色ユーティリティ（純Python・依存なし）
# ---------------------------------------------------------------------------

def hex2rgb(h):
    h = h.lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def rgb2hex(rgb):
    return "#%02X%02X%02X" % tuple(max(0, min(255, int(round(v)))) for v in rgb)


def mix(c1, c2, t):
    """c1 と c2 を線形補間。t は c2 の重み(0..1)。"""
    return tuple(c1[i] * (1 - t) + c2[i] * t for i in range(3))


def lighten(c, t):
    return mix(c, (255, 255, 255), t)


def darken(c, t):
    return mix(c, (0, 0, 0), t)


def lum(rgb):
    """相対輝度 0..1（sRGB）。"""
    def ch(v):
        v /= 255.0
        return v / 12.92 if v <= 0.03928 else ((v + 0.055) / 1.055) ** 2.4
    r, g, b = (ch(x) for x in rgb)
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def sat(rgb):
    """HSV 彩度 0..1。"""
    mx, mn = max(rgb), min(rgb)
    return 0.0 if mx == 0 else (mx - mn) / mx


def dist(a, b):
    return sum((a[i] - b[i]) ** 2 for i in range(3)) ** 0.5


# ---------------------------------------------------------------------------
# 4キー(bg/text/primary/accent) + フォント → 14+3 トークンへ導出
# ---------------------------------------------------------------------------

def derive_theme(name, bg, text, primary, accent, fonts):
    on_primary = (255, 255, 255) if lum(primary) < 0.55 else (27, 27, 36)
    colors = {
        "background": bg,
        "surface": lighten(primary, 0.92),
        "text": text,
        "muted": mix(text, bg, 0.45),
        "primary": primary,
        "accent": accent,
        "on_primary": on_primary,
        "on_primary_soft": lighten(primary, 0.90),
        "on_primary_muted": lighten(primary, 0.72),
        "code_bg": darken(primary, 0.62),
        "code_text": lighten(primary, 0.88),
        "table_header_bg": primary,
        "table_header_text": on_primary,
        "table_row_alt": lighten(primary, 0.92),
    }
    return {
        "name": name,
        "colors": {k: rgb2hex(v) for k, v in colors.items()},
        "fonts": fonts,
    }


# ---------------------------------------------------------------------------
# パレット分類（画像 / PDF / HTML 共通）: [(rgb, weight)] → bg,text,primary,accent
# ---------------------------------------------------------------------------

def classify_palette(colors, default_accent):
    """頻度付き色リストから 4 キーを推定。colors: [(rgb, weight)] weight 降順が望ましい。
    戻り値の最後は (colored, ambiguous)。ambiguous は primary/accent 候補が僅差で
    自動判定の確度が低いことを示す（要目視確認の警告に使う）。"""
    if not colors:
        raise ValueError("色を抽出できませんでした")
    by_weight = sorted(colors, key=lambda cw: -cw[1])
    rgbs = [c for c, _ in by_weight]
    # 背景: 明るい色で最も使われているもの（なければ最も明るい）
    light = [c for c in rgbs if lum(c) > 0.8]
    bg = light[0] if light else max(rgbs, key=lum)
    # 文字: 暗い色のうち低彩度（本文らしい無彩色〜低彩度）を優先して選ぶ。
    # 純度の高いブランド色は緑チャンネル重みの小さいsRGB輝度式では「暗い」と
    # 判定されることがあるため、彩度ガードを掛けてから輝度で絞り込む。
    dark_all = [c for c in rgbs if lum(c) < 0.25]
    dark_neutral = [c for c in dark_all if sat(c) <= TEXT_SAT_MAX]
    dark = dark_neutral or dark_all
    text = dark[0] if dark else min(rgbs, key=lum)
    # ブランド色候補: 彩度があり、bg/text から離れた色
    colored = [(c, w) for c, w in by_weight
               if sat(c) > 0.15 and dist(c, bg) > 40 and dist(c, text) > 40]
    ambiguous = False
    if colored:
        # primary/accent は単純な出現量ランキングでは入れ替わりやすい
        # （accentは装飾・小さい強調文字など「あちこちに少しずつ」使われがちで、
        #  出現量だけで比べると primary より優勢に見えることがある）。
        # 同梱2テーマ(default/accenture-purple)はいずれも primary が accent より
        # 暗い（輝度が低い）ため、出現量の多い上位候補プールの中では
        # 「最も暗い色」を primary、残りの中で最も彩度が高い色を accent とする。
        pool_n = min(4, len(colored))
        pool = [c for c, _ in colored[:pool_n]]
        primary = min(pool, key=lum)
        rest = [c for c in pool if c != primary]
        acc_cands = [c for c in rest if dist(c, primary) > 60]
        accent = max(acc_cands, key=sat) if acc_cands else (max(rest, key=sat) if rest else lighten(primary, 0.25))
        # 要確認シグナル: (a) 出現量の上位2件が僅差、または
        # (b) primaryとaccentの輝度差が小さく「暗い方=primary」の判定自体の確度が低い
        if len(colored) > 1 and colored[1][1] > 0:
            ambiguous = (colored[0][1] / colored[1][1]) < AMBIGUOUS_RATIO
        if abs(lum(primary) - lum(accent)) < 0.05:
            ambiguous = True
    else:
        primary = darken(text, 0.0) if lum(text) < 0.4 else (40, 40, 60)
        accent = default_accent
    return bg, text, primary, accent, bool(colored), ambiguous


# ---------------------------------------------------------------------------
# 形式別の抽出
# ---------------------------------------------------------------------------

_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _clr_from_elem(elem):
    """<a:srgbClr val=..> / <a:sysClr lastClr=..> から hex を得る。"""
    srgb = elem.find(f"{_A}srgbClr")
    if srgb is not None and srgb.get("val"):
        return hex2rgb(srgb.get("val"))
    sysclr = elem.find(f"{_A}sysClr")
    if sysclr is not None and sysclr.get("lastClr"):
        return hex2rgb(sysclr.get("lastClr"))
    return None


def _pptx_scheme_and_fonts(path):
    """theme1.xml から clrScheme（テーマ色スロット→rgb）と fontScheme を読む。"""
    with zipfile.ZipFile(path) as z:
        themes = [n for n in z.namelist() if re.match(r"ppt/theme/theme\d+\.xml$", n)]
        xml = z.read(sorted(themes)[0]) if themes else None
    scheme, fonts = {}, {}
    if xml:
        root = ET.fromstring(xml)
        clr = root.find(f"{_A}themeElements/{_A}clrScheme")
        for child in (list(clr) if clr is not None else []):
            rgb = _clr_from_elem(child)
            if rgb:
                scheme[child.tag.replace(_A, "")] = rgb
        fs = root.find(f"{_A}themeElements/{_A}fontScheme")
        if fs is not None:
            for slot, key in (("majorFont", "heading"), ("minorFont", "body")):
                latin = fs.find(f"{_A}{slot}/{_A}latin")
                if latin is not None and latin.get("typeface"):
                    fonts[key] = latin.get("typeface")
    return scheme, fonts


def extract_pptx(path):
    """スライド内で実際に使われている塗り色・文字色・フォントを面積/文字量で集計する。
    ビルダー生成PPTXやブランドPPTXは色を直接シェイプに適用することが多く、埋め込み
    clrScheme（既定Office色のことがある）より実使用色のほうが見た目に忠実なため。"""
    try:
        from pptx import Presentation
        from pptx.enum.dml import MSO_FILL, MSO_COLOR_TYPE, MSO_THEME_COLOR
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        sys.exit("error: PPTX 抽出には python-pptx が必要です。/slide-deck:setup（または pip install python-pptx）")

    scheme, scheme_fonts = _pptx_scheme_and_fonts(path)
    THEME_SLOT = {
        MSO_THEME_COLOR.DARK_1: "dk1", MSO_THEME_COLOR.TEXT_1: "dk1",
        MSO_THEME_COLOR.LIGHT_1: "lt1", MSO_THEME_COLOR.BACKGROUND_1: "lt1",
        MSO_THEME_COLOR.DARK_2: "dk2", MSO_THEME_COLOR.TEXT_2: "dk2",
        MSO_THEME_COLOR.LIGHT_2: "lt2", MSO_THEME_COLOR.BACKGROUND_2: "lt2",
        MSO_THEME_COLOR.ACCENT_1: "accent1", MSO_THEME_COLOR.ACCENT_2: "accent2",
        MSO_THEME_COLOR.ACCENT_3: "accent3", MSO_THEME_COLOR.ACCENT_4: "accent4",
        MSO_THEME_COLOR.ACCENT_5: "accent5", MSO_THEME_COLOR.ACCENT_6: "accent6",
        MSO_THEME_COLOR.HYPERLINK: "hlink", MSO_THEME_COLOR.FOLLOWED_HYPERLINK: "folHlink",
    }

    def resolve(color):
        try:
            if color.type == MSO_COLOR_TYPE.RGB:
                return hex2rgb(str(color.rgb))
            if color.type == MSO_COLOR_TYPE.SCHEME:
                return scheme.get(THEME_SLOT.get(color.theme_color))
        except Exception:
            return None
        return None

    fill_ct, text_ct, font_ct = Counter(), Counter(), Counter()

    def walk(shapes):
        for sh in shapes:
            try:
                if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                    walk(sh.shapes)
                    continue
            except Exception:
                pass
            # 塗り（面積重み）
            try:
                if sh.fill.type == MSO_FILL.SOLID:
                    rgb = resolve(sh.fill.fore_color)
                    if rgb:
                        area = max(1, (sh.width or 0) // 9525) * max(1, (sh.height or 0) // 9525)
                        fill_ct[rgb] += area
            except Exception:
                pass
            # 文字色・フォント（文字量重み）
            try:
                if sh.has_text_frame:
                    for para in sh.text_frame.paragraphs:
                        for run in para.runs:
                            n = max(1, len(run.text))
                            if run.font.name:
                                font_ct[run.font.name] += n
                            rgb = resolve(run.font.color)
                            if rgb:
                                text_ct[rgb] += n
            except Exception:
                pass

    prs = Presentation(path)
    for slide in prs.slides:
        walk(slide.shapes)

    # パレット = 塗り色（面積） + 文字色（文字量を面積スケールに寄せる）
    palette = Counter()
    for rgb, w in fill_ct.items():
        palette[rgb] += w
    tscale = (sum(fill_ct.values()) / max(1, sum(text_ct.values()))) * 0.5 if text_ct else 1
    for rgb, w in text_ct.items():
        palette[rgb] += w * tscale

    colors = [(c, w) for c, w in palette.items()]
    ambiguous = False
    if colors:
        bg, text, primary, accent, colored, ambiguous = classify_palette(colors, DEFAULT_ACCENT)
    else:
        # 塗りが取れない場合は埋め込みテーマ色にフォールバック
        bg = scheme.get("lt1", (255, 255, 255)); text = scheme.get("dk1", (26, 34, 51))
        primary = scheme.get("accent1", (15, 61, 110)); accent = scheme.get("accent2", lighten(primary, 0.25))
        colored = True

    fonts = {}
    if font_ct:
        fonts["heading"] = fonts["body"] = font_ct.most_common(1)[0][0]
    fonts = {**scheme_fonts, **fonts}  # 実使用フォント優先、無ければ fontScheme

    report = {"source": "pptx", "detected_colored": colored, "ambiguous": ambiguous,
              "palette": [rgb2hex(c) for c, _ in sorted(colors, key=lambda x: -x[1])[:8]],
              "fonts_found": [n for n, _ in font_ct.most_common(5)] or list(scheme_fonts.values())}
    return bg, text, primary, accent, fonts, report


def _image_palette(pil_img, ncolors):
    from PIL import Image
    im = pil_img.convert("RGB")
    im.thumbnail((400, 400))
    q = im.quantize(colors=max(4, ncolors), method=Image.Quantize.FASTOCTREE)
    pal = q.getpalette()
    total = sum(c for c, _ in q.getcolors())
    out = []
    for count, idx in q.getcolors():
        rgb = tuple(pal[idx * 3: idx * 3 + 3])
        out.append((rgb, count / total))
    return out


def extract_image(path, ncolors):
    try:
        from PIL import Image
    except ImportError:
        fail("画像抽出には pillow が必要です。/slide-deck:setup --pillow（または pip install pillow）")
    colors = _image_palette(Image.open(path), ncolors)
    bg, text, primary, accent, colored, ambiguous = classify_palette(colors, DEFAULT_ACCENT)
    report = {"source": "image", "detected_colored": colored, "ambiguous": ambiguous,
              "palette": [rgb2hex(c) for c, _ in sorted(colors, key=lambda x: -x[1])[:8]]}
    return bg, text, primary, accent, {}, report


def extract_pdf(path, pages, ncolors):
    try:
        import fitz  # PyMuPDF
    except ImportError:
        fail("PDF 抽出には pymupdf が必要です。/slide-deck:setup --pdf（または pip install pymupdf）")
    from PIL import Image
    doc = fitz.open(path)
    counter = Counter()
    fonts_ct = Counter()
    for i in range(min(pages, doc.page_count)):
        page = doc[i]
        pix = page.get_pixmap(matrix=fitz.Matrix(0.35, 0.35))
        img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
        for rgb, w in _image_palette(img, ncolors):
            counter[rgb] += w
        for f in page.get_fonts(full=True):
            nm = (f[3] or "").split("+")[-1]
            if nm:
                fonts_ct[nm] += 1
    colors = [(c, w) for c, w in counter.items()]
    bg, text, primary, accent, colored, ambiguous = classify_palette(colors, DEFAULT_ACCENT)
    fonts = {}
    if fonts_ct:
        top = fonts_ct.most_common(1)[0][0]
        fonts = {"heading": top, "body": top}
    report = {"source": "pdf", "detected_colored": colored, "ambiguous": ambiguous,
              "palette": [rgb2hex(c) for c, _ in sorted(colors, key=lambda x: -x[1])[:8]],
              "fonts_found": [n for n, _ in fonts_ct.most_common(5)]}
    return bg, text, primary, accent, fonts, report


_HEX6_RE = re.compile(r"#[0-9a-fA-F]{6}\b")
_HEX3_RE = re.compile(r"#[0-9a-fA-F]{3}\b")
_RGB_RE = re.compile(r"rgba?\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)")
_FONT_SIZE_RE = re.compile(r"font-size\s*:\s*([\d.]+)\s*px", re.I)
_HEADING_TAG_RE = re.compile(r"^h[1-3]$", re.I)
_HEADING_CLASS_RE = re.compile(r"(title|heading|punch|eyebrow)", re.I)


def _hex_colors_in(text):
    out = [hex2rgb(h) for h in _HEX6_RE.findall(text)]
    out += [hex2rgb(h) for h in _HEX3_RE.findall(text)]
    out += [tuple(int(x) for x in m) for m in _RGB_RE.findall(text)]
    return out


def extract_html(path, ncolors):
    """HTML内の色を、見出しらしい要素・font-sizeが大きい宣言ほど重く数えて集計する。
    見出しタグ/大きい文字ほど「ブランドの主色」を反映しやすいという前提（出現回数だけで
    数えると、頻出するが小さい装飾文字(eyebrow等)にprimaryが引きずられてしまうため）。"""
    text_src = Path(path).read_text(encoding="utf-8", errors="ignore")
    counter = Counter()

    # (1) <style>...</style> 内の各ルール、および要素の style="..." 属性を
    #     「宣言ブロック」として集め、font-size/見出しらしさで重み付けする。
    blocks = []  # (selector_or_tag, declaration_text)
    for style_block in re.findall(r"<style[^>]*>(.*?)</style>", text_src, re.I | re.S):
        for selector, body in re.findall(r"([^{}]+)\{([^{}]+)\}", style_block):
            blocks.append((selector.strip(), body))
    for tag, attrs in re.findall(r"<([a-zA-Z][a-zA-Z0-9]*)\b([^>]*)>", text_src):
        m = re.search(r'style\s*=\s*"([^"]*)"', attrs)
        if m:
            cls = re.search(r'class\s*=\s*"([^"]*)"', attrs)
            selector = tag + (" " + cls.group(1) if cls else "")
            blocks.append((selector, m.group(1)))

    weighted_hits = 0
    for selector, body in blocks:
        hits = _hex_colors_in(body)
        if not hits:
            continue
        fs_m = _FONT_SIZE_RE.search(body)
        weight = max(1.0, float(fs_m.group(1)) / 16.0) if fs_m else 1.0
        if _HEADING_TAG_RE.search(selector) or _HEADING_CLASS_RE.search(selector):
            weight *= 3.0
        for rgb in hits:
            counter[rgb] += weight
            weighted_hits += 1

    # (2) フォールバック: どこにも属さない #hex / rgb() も低い重みで拾い、
    #     (1) で拾いきれない外部CSS等由来の宣言も判定材料に加える。
    for rgb in _hex_colors_in(text_src):
        counter[rgb] += 0.2
    if not counter:
        fail("HTML から色を抽出できませんでした（インラインの #hex / rgb() が必要）")

    colors = [(c, n) for c, n in counter.items()]
    bg, text, primary, accent, colored, ambiguous = classify_palette(colors, DEFAULT_ACCENT)
    # font-family: 引用符付きも拾い、先頭ファミリだけを取り出す
    fams = re.findall(r"font-family\s*:\s*([^;}\n]+)", text_src, re.I)
    first_fams = [f.split(",")[0].strip().strip("'\"") for f in fams]
    first_fams = [f for f in first_fams if f and "monospace" not in f.lower()]
    fonts = {}
    counts = Counter(first_fams)
    if counts:
        top = counts.most_common(1)[0][0]
        fonts = {"heading": top, "body": top}
    report = {"source": "html", "detected_colored": colored, "ambiguous": ambiguous,
              "palette": [rgb2hex(c) for c, _ in sorted(colors, key=lambda x: -x[1])[:8]],
              "fonts_found": [n for n, _ in counts.most_common(5)]}
    return bg, text, primary, accent, fonts, report


# ---------------------------------------------------------------------------

def load_default():
    return json.loads((BUNDLED_THEMES / "default.json").read_text(encoding="utf-8"))


DEFAULT_ACCENT = (0, 168, 204)  # default.json の accent（フォールバック用）。main で上書き


def resolve_out_dir(arg_dir):
    if arg_dir:
        return Path(arg_dir)
    env = os.environ.get("SLIDE_DECK_THEMES")
    if env:
        first = next((p for p in env.split(os.pathsep) if p.strip()), None)
        if first:
            return Path(first)
    return Path.cwd() / "themes"


DISPATCH = {
    ".pptx": "pptx", ".pptm": "pptx",
    ".png": "image", ".jpg": "image", ".jpeg": "image", ".webp": "image", ".bmp": "image",
    ".pdf": "pdf",
    ".html": "html", ".htm": "html",
}


def main():
    setup_console()
    global DEFAULT_ACCENT
    ap = argparse.ArgumentParser(description="サンプルスライドからテーマJSONを生成")
    ap.add_argument("file", help="サンプル(.pptx/.png/.jpg/.pdf/.html)")
    ap.add_argument("--name", required=True, help="作成するテーマ名（kebab-case）")
    ap.add_argument("--dir", help="出力先ディレクトリ（既定: SLIDE_DECK_THEMES / ./themes/）")
    ap.add_argument("--pages", type=int, default=2, help="PDFで走査するページ数（既定2）")
    ap.add_argument("--colors", type=int, default=12,
                     help="量子化する色数（既定12。小さすぎると本文文字のような面積の小さい"
                          "無彩色クラスタが消え、text判定を誤りやすい）")
    ap.add_argument("--force", action="store_true", help="既存を上書き")
    args = ap.parse_args()

    if not NAME_RE.match(args.name) or args.name == "default":
        fail(f"テーマ名は kebab-case で 'default' 以外にしてください: {args.name}")
    src = Path(args.file)
    if not src.exists():
        fail(f"ファイルがありません: {src}")
    kind = DISPATCH.get(src.suffix.lower())
    if not kind:
        fail(f"未対応の形式です: {src.suffix}（対応: {', '.join(sorted(DISPATCH))}）")

    default_theme = load_default()
    DEFAULT_ACCENT = hex2rgb(default_theme["colors"]["accent"])

    if kind == "pptx":
        bg, text, primary, accent, fonts, report = extract_pptx(src)
    elif kind == "image":
        bg, text, primary, accent, fonts, report = extract_image(src, args.colors)
    elif kind == "pdf":
        bg, text, primary, accent, fonts, report = extract_pdf(src, args.pages, args.colors)
    else:
        bg, text, primary, accent, fonts, report = extract_html(src, args.colors)

    # フォントは検出できた分だけ採用、残りは default から補完
    merged_fonts = dict(default_theme["fonts"])
    merged_fonts.update({k: v for k, v in fonts.items() if v})
    theme = derive_theme(args.name, bg, text, primary, accent, merged_fonts)

    out_dir = resolve_out_dir(args.dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / f"{args.name}.json"
    if out_path.exists() and not args.force:
        fail(f"すでに存在します: {out_path}（上書きは --force）")
    try:
        out_path.write_text(json.dumps(theme, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    except PermissionError:
        fail(f"{out_path} に書き込めません（他アプリで開いている場合は閉じてから再実行してください）")

    # レポート
    print(f"サンプル: {src}  ({report['source']})")
    if "scheme" in report:
        print(f"  検出テーマ色: {report['scheme']}")
    if "palette" in report:
        print(f"  主要色: {', '.join(report['palette'])}")
    if report.get("fonts_found"):
        print(f"  検出フォント: {', '.join(report['fonts_found'])}")
    if not report.get("detected_colored", True):
        warn("有彩色を検出できず、primary/accent を推定/既定で補いました。要調整。")
    if report.get("ambiguous"):
        warn("primary/accent の候補が僅差でした。誤って入れ替わっている可能性があるため要確認。")
    print("  -> マッピング:")
    for k in ("background", "text", "primary", "accent"):
        print(f"      {k:10} = {theme['colors'][k]}")
    print(f"  フォント: heading={theme['fonts']['heading']} / body={theme['fonts']['body']} / code={theme['fonts']['code']}")
    print(f"\n作成しました: {out_path}")
    if out_dir == BUNDLED_THEMES:
        note("同梱テーマ置き場に作成しました。更新で消したくない場合は SLIDE_DECK_THEMES を設定するか --dir で明示してください。")
    print("書き込み先: " + str(out_dir))
    print("ビルダーのテーマ探索順: 環境変数 SLIDE_DECK_THEMES / <deck_dir>/themes/ / 同梱 templates/themes/")
    print("自動推定なので colors/fonts を目視で微調整してください（トークンの意味は references/themes.md）。")
    print(f"使い方: deck.json の meta.theme に \"{args.name}\" を指定して build_deck.py で再ビルド。作成後は check_theme.py でコントラストも確認してください。")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
